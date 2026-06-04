"""
PowerPoint Generator for: Agentic AI Mission Control
Kids Computer Science Class - June 6, 2026

Uses the standard theme from tools/create_theme.py.

Run:
    py create_presentation.py

Dependencies:
    pip install python-pptx
"""

import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "..", "tools"))
from create_theme import *  # noqa: F401,F403

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


COLORS["agent_purple"] = RGBColor(103, 58, 183)
COLORS["mission_teal"] = RGBColor(0, 137, 123)
COLORS["safety_red"] = RGBColor(211, 47, 47)
COLORS["test_green"] = RGBColor(56, 142, 60)
COLORS["permission_amber"] = RGBColor(255, 179, 0)
COLORS["soft_blue"] = RGBColor(227, 242, 253)
COLORS["soft_teal"] = RGBColor(224, 242, 241)
COLORS["soft_red"] = RGBColor(255, 235, 238)


def add_card(slide, x, y, w, h, title, body, color):
    add_rounded_box(slide, Inches(x), Inches(y), Inches(w), Inches(h), color)
    add_styled_textbox(
        slide, Inches(x + 0.15), Inches(y + 0.12), Inches(w - 0.3), Inches(0.35),
        title, font_size=15, font_color=COLORS["white"], bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_styled_textbox(
        slide, Inches(x + 0.18), Inches(y + 0.55), Inches(w - 0.36), Inches(h - 0.65),
        body, font_size=12, font_color=COLORS["white"],
        alignment=PP_ALIGN.CENTER,
    )


def add_white_card(slide, x, y, w, h, title, body, color):
    add_rounded_box(
        slide, Inches(x), Inches(y), Inches(w), Inches(h),
        COLORS["light_gray"], border_color=color,
    )
    add_styled_textbox(
        slide, Inches(x + 0.18), Inches(y + 0.12), Inches(w - 0.36), Inches(0.32),
        title, font_size=14, font_color=color, bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_styled_textbox(
        slide, Inches(x + 0.2), Inches(y + 0.5), Inches(w - 0.4), Inches(h - 0.55),
        body, font_size=12, font_color=COLORS["dark_gray"],
        alignment=PP_ALIGN.CENTER,
    )


def slide_title(prs):
    create_title_slide(
        prs,
        "Agentic AI\nMission Control",
        "From Vibe Coding to Supervised Agents",
        "June 6, 2026",
        tagline="AI Phase  |  Kids Computer Science Class",
    )


def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Today's Journey")
    items = [
        ("Recap: building apps with AI", COLORS["medium_blue"]),
        ("What makes a workflow agentic?", COLORS["agent_purple"]),
        ("The agent job brief", COLORS["mission_teal"]),
        ("Tools, permissions, and approval", COLORS["safety_red"]),
        ("Teacher demo: mission control", COLORS["permission_amber"]),
        ("Classwork: improve and test an app", COLORS["test_green"]),
        ("Kahoot + homework", COLORS["orange"]),
    ]
    for i, (text, color) in enumerate(items):
        add_agenda_item(slide, i + 1, text, 1.55 + i * 0.72, color)


def slide_may30_recap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Quick Recap: May 30")
    cards = [
        ("Canvas + Lovable", "AI app-building tools", COLORS["medium_blue"]),
        ("HTML", "Structure: buttons,\ntext, sections", COLORS["mission_teal"]),
        ("CSS", "Style: colors,\nfonts, layout", COLORS["permission_amber"]),
        ("JavaScript", "Behavior: clicks,\nscore, movement", COLORS["agent_purple"]),
    ]
    for i, (title, body, color) in enumerate(cards):
        add_card(slide, 0.6 + i * 3.2, 1.7, 2.75, 1.55, title, body, color)

    add_styled_textbox(
        slide, Inches(1.0), Inches(4.0), Inches(11.3), Inches(0.7),
        "Last week: describe -> build -> test -> refine",
        font_size=24, font_color=COLORS["dark_blue"], bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_takeaway_bar(slide, "Today we add supervision: goals, limits, approvals, and tests.")


def slide_three_modes(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Three Ways to Use AI")
    modes = [
        ("Chatbot", "Answers questions\nor explains ideas", COLORS["medium_blue"]),
        ("App Builder", "Turns prompts into\napps or games", COLORS["mission_teal"]),
        ("Agentic Workflow", "Plans, uses tools,\nchecks, and refines", COLORS["agent_purple"]),
    ]
    for i, (title, body, color) in enumerate(modes):
        x = 0.9 + i * 4.1
        add_card(slide, x, 2.0, 3.35, 1.75, title, body, color)
        if i < 2:
            add_arrow(slide, Inches(x + 3.45), Inches(2.62), Inches(0.5), Inches(0.35), COLORS["orange"])
    add_takeaway_bar(slide, "Agentic does not mean unsupervised. The human is still in charge.")


def slide_agent_timeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "How We Got Here")

    add_styled_textbox(
        slide, Inches(0.75), Inches(1.82), Inches(11.85), Inches(0.44),
        "From chatbots to agents that can use code, browsers, and computers",
        font_size=20,
        font_color=COLORS["dark_blue"],
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    milestones = [
        ("Nov 2022", "ChatGPT", "Chat interface\nbecomes mainstream", COLORS["medium_blue"]),
        ("Mar 2023", "GPT-4 + Claude", "Stronger reasoning\nassistants", COLORS["agent_purple"]),
        ("Oct 2024", "Computer Use", "Claude can look,\nclick, and type", COLORS["mission_teal"]),
        ("Jan 2025", "Operator", "Browser agent for\nweb tasks", COLORS["permission_amber"]),
        ("Feb 2025", "Claude Code", "Coding agent in\nthe terminal", COLORS["test_green"]),
        ("May 2025", "Codex", "Cloud coding agent\ninside ChatGPT", COLORS["medium_blue"]),
        ("2026", "OpenClaw + Hermes", "Agent platforms\nconnect models to tools", COLORS["safety_red"]),
    ]

    x_positions = [0.5, 2.3, 4.1, 5.9, 7.7, 9.5, 11.3]
    line_y = 3.65

    # Timeline rail
    rail = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.75), Inches(line_y + 0.22),
        Inches(11.85), Inches(0.06),
    )
    set_shape_fill(rail, COLORS["light_gray"])
    rail.line.fill.background()

    for i, ((date, title, body, color), x) in enumerate(zip(milestones, x_positions)):
        y_card = 2.55 if i % 2 == 0 else 4.35
        add_circle(slide, Inches(x + 0.52), Inches(line_y), Inches(0.34), color)
        add_rounded_box(slide, Inches(x), Inches(y_card), Inches(1.45), Inches(1.18), color)
        add_styled_textbox(
            slide, Inches(x + 0.06), Inches(y_card + 0.08), Inches(1.33), Inches(0.25),
            date, font_size=9, font_color=COLORS["white"], bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_styled_textbox(
            slide, Inches(x + 0.06), Inches(y_card + 0.34), Inches(1.33), Inches(0.25),
            title, font_size=9, font_color=COLORS["white"], bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_styled_textbox(
            slide, Inches(x + 0.06), Inches(y_card + 0.64), Inches(1.33), Inches(0.42),
            body, font_size=7.4, font_color=COLORS["white"],
            alignment=PP_ALIGN.CENTER,
        )

    add_takeaway_bar(slide, "These are examples of the trend. Today we focus on safe supervision, not installing new tools.")


def slide_agentic_loop(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "The Supervised Agent Loop")
    steps = [
        ("PLAN", "What should I do?", COLORS["agent_purple"]),
        ("ACT", "Use one tool", COLORS["permission_amber"]),
        ("OBSERVE", "What happened?", COLORS["mission_teal"]),
        ("TEST", "Did it work?", COLORS["test_green"]),
        ("REFINE", "Fix or improve", COLORS["medium_blue"]),
    ]
    for i, (label, body, color) in enumerate(steps):
        x = 0.5 + i * 2.55
        add_card(slide, x, 2.05, 2.2, 1.45, label, body, color)
        if i < len(steps) - 1:
            add_arrow(slide, Inches(x + 2.22), Inches(2.55), Inches(0.38), Inches(0.3), COLORS["orange"])
    add_styled_textbox(
        slide, Inches(1.1), Inches(4.35), Inches(11.1), Inches(0.65),
        "A good supervisor slows the AI down at the risky moments.",
        font_size=22, font_color=COLORS["dark_blue"], bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_takeaway_bar(slide, "Same loop as debugging, but with clearer rules.")


def slide_job_brief(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "The Agent Job Brief")
    fields = [
        ("Goal", "What should happen?"),
        ("Success Criteria", "How will we know?"),
        ("Allowed Tools", "What can it use?"),
        ("Not Allowed", "What is off limits?"),
        ("Approval Required", "When must it ask?"),
        ("Test Checklist", "How do we verify?"),
    ]
    colors = [
        COLORS["mission_teal"], COLORS["test_green"], COLORS["medium_blue"],
        COLORS["safety_red"], COLORS["permission_amber"], COLORS["agent_purple"],
    ]
    for i, ((title, body), color) in enumerate(zip(fields, colors)):
        row = i // 3
        col = i % 3
        add_white_card(slide, 0.75 + col * 4.15, 1.7 + row * 1.85, 3.45, 1.35, title, body, color)
    add_takeaway_bar(slide, "Do not just prompt. Write the mission first.")


def slide_success_criteria(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Success Criteria Make Testing Possible")
    add_white_card(slide, 0.8, 1.7, 5.5, 3.0, "Vague", "\"Make the game better\"\n\nHard to test.\nDifferent people may imagine different results.", COLORS["safety_red"])
    add_white_card(slide, 7.0, 1.7, 5.5, 3.0, "Clear", "\"The restart button resets the score to 0\"\n\nEasy to test.\nEveryone knows what done means.", COLORS["test_green"])
    add_takeaway_bar(slide, "If you cannot test it, the agent cannot prove it finished.")


def slide_tools(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Tools Give Agents Hands")
    tools = [
        ("Browser", "Preview apps\nor open websites", COLORS["medium_blue"]),
        ("Files", "Read or edit\nproject files", COLORS["mission_teal"]),
        ("Code", "Run scripts\nor tests", COLORS["agent_purple"]),
        ("Search", "Look up live\ninformation", COLORS["permission_amber"]),
        ("Messages", "Send updates\nor reminders", COLORS["safety_red"]),
    ]
    for i, (title, body, color) in enumerate(tools):
        add_card(slide, 0.55 + i * 2.55, 2.0, 2.2, 1.6, title, body, color)
    add_takeaway_bar(slide, "More tools = more power and more responsibility.")


def slide_permissions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Permissions Before Power")
    columns = [
        ("Usually Safe", "Read a page\nPreview an app\nExplain code\nMake a checklist", COLORS["test_green"]),
        ("Needs Care", "Search the web\nEdit a file\nRun code\nUse memory", COLORS["permission_amber"]),
        ("Ask First", "Delete files\nSend messages\nSpend money\nShare private info", COLORS["safety_red"]),
    ]
    for i, (title, body, color) in enumerate(columns):
        add_white_card(slide, 0.8 + i * 4.15, 1.75, 3.55, 3.55, title, body, color)
    add_takeaway_bar(slide, "A safe agent knows when to stop and ask.")


def slide_future_tools_preview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Future Agent Tools Preview")
    add_styled_textbox(
        slide, Inches(0.8), Inches(1.65), Inches(11.8), Inches(0.6),
        "Some agents can connect an LLM to real tools.",
        font_size=24, font_color=COLORS["dark_blue"], bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    flow = [
        ("LLM", "reasons", COLORS["agent_purple"]),
        ("Tools", "files, browser,\ncommands, messages", COLORS["mission_teal"]),
        ("Skills", "instructions for\nrepeatable tasks", COLORS["medium_blue"]),
        ("Human", "sets limits\nand approves", COLORS["safety_red"]),
    ]
    for i, (title, body, color) in enumerate(flow):
        x = 0.8 + i * 3.1
        add_card(slide, x, 2.75, 2.55, 1.5, title, body, color)
        if i < 3:
            add_arrow(slide, Inches(x + 2.6), Inches(3.25), Inches(0.42), Inches(0.28), COLORS["orange"])
    add_takeaway_bar(slide, "Today: learn the workflow. Later: use tool-based agents safely.")


def slide_approval_gate(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "The Approval Gate")
    add_styled_textbox(
        slide, Inches(1.0), Inches(1.55), Inches(11.3), Inches(0.6),
        "Before risky actions, the agent must stop.",
        font_size=25, font_color=COLORS["dark_blue"], bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_card(slide, 1.2, 2.55, 3.0, 1.45, "ASK", "Explain the action\nand why it matters", COLORS["permission_amber"])
    add_arrow(slide, Inches(4.35), Inches(3.08), Inches(0.6), Inches(0.35), COLORS["orange"])
    add_card(slide, 5.15, 2.55, 3.0, 1.45, "APPROVE", "Human says yes,\nno, or change plan", COLORS["mission_teal"])
    add_arrow(slide, Inches(8.3), Inches(3.08), Inches(0.6), Inches(0.35), COLORS["orange"])
    add_card(slide, 9.1, 2.55, 3.0, 1.45, "ACT", "Only then use\nthe risky tool", COLORS["test_green"])
    add_takeaway_bar(slide, "Approval is not a slowdown. It is a safety feature.")


def slide_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Teacher Demo: Mission Control")
    steps = [
        ("1", "Build a small quiz app"),
        ("2", "Write the job brief"),
        ("3", "Ask for a plan first"),
        ("4", "Improve one thing"),
        ("5", "Test and refine"),
    ]
    for i, (num, text) in enumerate(steps):
        y = 1.7 + i * 0.78
        add_circle(slide, Inches(1.1), Inches(y), Inches(0.5), COLORS["agent_purple"], text=num, text_size=16)
        add_styled_textbox(slide, Inches(1.9), Inches(y + 0.07), Inches(10.7), Inches(0.42), text, font_size=19, font_color=COLORS["dark_gray"])
    add_takeaway_bar(slide, "Watch for: plan -> act -> observe -> test -> refine.")


def slide_supervisor_prompts(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Good Supervisor Prompts")
    prompts = [
        ("Plan First", "Do not change the app yet.\nMake a short plan."),
        ("Small Change", "Make only one change.\nKeep everything else working."),
        ("Explain", "Tell me what you changed\nand why."),
        ("Test", "Give me a checklist\nso I can verify it."),
    ]
    colors = [COLORS["agent_purple"], COLORS["mission_teal"], COLORS["medium_blue"], COLORS["test_green"]]
    for i, ((title, body), color) in enumerate(zip(prompts, colors)):
        add_white_card(slide, 0.75 + i * 3.15, 2.0, 2.65, 2.05, title, body, color)
    add_takeaway_bar(slide, "The best prompts make the AI easier to supervise.")


def slide_classwork(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Classwork Mission")
    items = [
        ("Start", "Use May 30 app\nor build fresh", COLORS["medium_blue"]),
        ("Brief", "Write goal,\nlimits, tests", COLORS["agent_purple"]),
        ("Improve", "Make 2 small\nimprovements", COLORS["mission_teal"]),
        ("Debug", "Fix 1 bug or\nconfusing behavior", COLORS["permission_amber"]),
        ("Test", "Use your\nchecklist", COLORS["test_green"]),
    ]
    for i, (title, body, color) in enumerate(items):
        add_card(slide, 0.55 + i * 2.55, 2.0, 2.2, 1.65, title, body, color)
    add_takeaway_bar(slide, "Total: 100 points + 10 bonus. Submit on Teams and copy Ishwari Raut ma'am.")


def slide_testing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Test Like a Mission Controller")
    checks = [
        "Open Preview",
        "Click every button",
        "Check score/result",
        "Restart or reset",
        "Read the text",
        "Try one weird input",
    ]
    for i, text in enumerate(checks):
        row = i // 2
        col = i % 2
        add_white_card(slide, 1.0 + col * 5.7, 1.65 + row * 1.35, 4.9, 0.9, f"Test {i + 1}", text, COLORS["test_green"])
    add_takeaway_bar(slide, "Looking is not testing. Testing means trying the app.")


def slide_failures(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Common Agent Failure Modes")
    failures = [
        ("Too Big", "Changed everything\nat once", COLORS["safety_red"]),
        ("Skipped Test", "Looks done but\nbutton is broken", COLORS["permission_amber"]),
        ("Wrong Tool", "Used the tool that\ndid not fit", COLORS["medium_blue"]),
        ("Overconfident", "Claims fixed but\nstill wrong", COLORS["agent_purple"]),
    ]
    for i, (title, body, color) in enumerate(failures):
        add_card(slide, 0.85 + i * 3.1, 2.0, 2.55, 1.75, title, body, color)
    add_takeaway_bar(slide, "The fix is simple: smaller steps, clearer limits, better tests.")


def slide_key_takeaways(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Today's Key Takeaways")
    takeaways = [
        "Agentic AI needs goals, tools, limits, approvals, and tests.",
        "The human is the mission controller.",
        "Small changes are easier to test than giant changes.",
        "Tool-using agents are powerful because they can touch real systems.",
        "Safety rules come before autonomy.",
    ]
    for i, text in enumerate(takeaways):
        y = 1.65 + i * 0.82
        add_circle(slide, Inches(0.9), Inches(y), Inches(0.5), COLORS["mission_teal"], text=str(i + 1), text_size=15)
        add_styled_textbox(slide, Inches(1.65), Inches(y + 0.06), Inches(11.0), Inches(0.45), text, font_size=18, font_color=COLORS["dark_gray"])
    add_takeaway_bar(slide, "Next: Kahoot, then homework explanation.")


def slide_questions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLORS["dark_blue"])
    add_styled_textbox(
        slide, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.85),
        "Questions?",
        font_size=54,
        font_color=COLORS["white"],
        bold=True,
        alignment=PP_ALIGN.CENTER,
        font_name="Segoe UI Semibold",
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(3.0), Inches(11.33), Inches(0.65),
        "Before we start the hands-on mission:",
        font_size=24,
        font_color=COLORS["sky_blue"],
        alignment=PP_ALIGN.CENTER,
    )
    add_rounded_box(
        slide, Inches(2.0), Inches(4.15), Inches(9.33), Inches(0.95),
        COLORS["mission_teal"],
        text="What should an agent ask permission for?",
        text_size=24,
    )
    add_styled_textbox(
        slide, Inches(0.8), Inches(6.25), Inches(11.8), Inches(0.45),
        "Then: classwork, share/debug gallery, Kahoot, and homework.",
        font_size=18,
        font_color=COLORS["white"],
        alignment=PP_ALIGN.CENTER,
    )


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_agenda(prs)
    slide_may30_recap(prs)
    slide_three_modes(prs)
    slide_agent_timeline(prs)
    slide_agentic_loop(prs)
    slide_job_brief(prs)
    slide_success_criteria(prs)
    slide_tools(prs)
    slide_permissions(prs)
    slide_future_tools_preview(prs)
    slide_approval_gate(prs)
    slide_demo(prs)
    slide_supervisor_prompts(prs)
    slide_classwork(prs)
    slide_testing(prs)
    slide_failures(prs)
    slide_key_takeaways(prs)
    slide_questions(prs)

    return prs


def main():
    script_dir = os.path.dirname(os.path.abspath(__file__))
    parent_dir = os.path.dirname(script_dir)
    output_path = os.path.join(parent_dir, "2026-06-06_Agentic_AI_Mission_Control.pptx")
    prs = build_presentation()
    prs.save(output_path)
    print(f"[SUCCESS] Presentation created: {output_path}")


if __name__ == "__main__":
    main()
