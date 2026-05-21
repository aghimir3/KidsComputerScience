"""
PowerPoint Generator for: AI Image & Multimodal Models
Kids Computer Science Class - May 23, 2026

Uses the standard theme from tools/create_theme.py.
Week 7 of the AI phase — from text to images and beyond.

Run:
    py create_presentation.py

Dependencies:
    pip install python-pptx
"""

import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', 'tools'))
from create_theme import *

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


# ---- Theme-specific colors ----
COLORS["image_purple"] = RGBColor(126, 87, 194)
COLORS["diffusion_blue"] = RGBColor(33, 150, 243)
COLORS["multimodal_teal"] = RGBColor(0, 150, 136)
COLORS["ethics_red"] = RGBColor(229, 57, 53)
COLORS["creative_pink"] = RGBColor(236, 64, 122)
COLORS["prompt_gold"] = RGBColor(255, 179, 0)


# =============================================================================
# SLIDE FUNCTIONS
# =============================================================================

def slide_title(prs):
    create_title_slide(
        prs,
        "AI Image &\nMultimodal Models",
        "From Text to Pictures — How AI Sees and Creates",
        "May 23, 2026",
        tagline="Week 7 of the AI Phase  |  Kids Computer Science Class"
    )


def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Today's Journey")

    items = [
        ("Quick Recap: How LLMs Generate Text", COLORS["light_blue"]),
        ("What Are AI Image Models?", COLORS["image_purple"]),
        ("How They Work: From Noise to Art", COLORS["diffusion_blue"]),
        ("Multimodal AI: Text + Images + More", COLORS["multimodal_teal"]),
        ("Prompting for Images vs. Text", COLORS["prompt_gold"]),
        ("Live Demo: Generate Images Together!", COLORS["green"]),
        ("Ethics: Deepfakes & Responsible Use", COLORS["ethics_red"]),
        ("Kahoot!", COLORS["orange"]),
    ]
    for i, (text, color) in enumerate(items):
        add_agenda_item(slide, i + 1, text, 1.6 + i * 0.65, color)


def slide_recap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Quick Recap: How LLMs Generate Text")

    # Three step flow
    steps = [
        ("TOKENS", COLORS["light_blue"], "Text split into\nsmall pieces"),
        ("PREDICT", COLORS["medium_blue"], "Next token based\non context"),
        ("GENERATE", COLORS["dark_blue"], "One token at a time\nbuilds the answer"),
    ]
    for i, (label, color, desc) in enumerate(steps):
        x = 1.0 + i * 4.0
        add_rounded_box(slide, Inches(x), Inches(2.0), Inches(3.2), Inches(1.4),
                       color, text=label, text_size=22)
        add_styled_textbox(slide, Inches(x), Inches(3.5), Inches(3.2), Inches(0.8),
                          desc, font_size=14, font_color=COLORS["dark_gray"],
                          alignment=PP_ALIGN.CENTER)
        if i < 2:
            add_arrow(slide, Inches(x + 3.3), Inches(2.5), Inches(0.6), Inches(0.4),
                     COLORS["orange"])

    add_takeaway_bar(slide, "LLMs work with TEXT. But what about IMAGES?")


def slide_what_are_image_models(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "What Are AI Image Models?")

    add_styled_textbox(
        slide, Inches(0.5), Inches(1.6), Inches(12.5), Inches(0.6),
        "AI systems that can CREATE or UNDERSTAND images",
        font_size=22, font_color=COLORS["dark_gray"], bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # Three boxes
    models = [
        ("DALL-E 3", "by OpenAI", "Built into ChatGPT\nHigh quality, safe", COLORS["image_purple"]),
        ("Midjourney", "Independent", "Beautiful art style\nUsed by designers", COLORS["diffusion_blue"]),
        ("Stable Diffusion", "Open Source", "Free, runs on your PC\nHighly customizable", COLORS["multimodal_teal"]),
    ]
    for i, (name, maker, desc, color) in enumerate(models):
        x = 0.8 + i * 4.2
        add_rounded_box(slide, Inches(x), Inches(2.5), Inches(3.6), Inches(1.0),
                       color, text=name, text_size=20)
        add_styled_textbox(slide, Inches(x), Inches(3.55), Inches(3.6), Inches(0.35),
                          maker, font_size=12, font_color=COLORS["dark_gray"],
                          alignment=PP_ALIGN.CENTER)
        add_styled_textbox(slide, Inches(x + 0.2), Inches(3.9), Inches(3.2), Inches(0.8),
                          desc, font_size=13, font_color=COLORS["dark_gray"],
                          alignment=PP_ALIGN.CENTER)

    add_takeaway_bar(slide,
        "You type a text description (prompt) and the AI creates an image from scratch!")


def slide_how_it_works_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "How Do Image Models Work?",
                 subtitle="The Big Idea: Diffusion (Removing Noise)")

    # Analogy box
    add_rounded_box(slide, Inches(0.5), Inches(1.6), Inches(12.33), Inches(1.2),
                   COLORS["sky_blue"])
    add_styled_textbox(slide, Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.5),
                      "Imagine a sculptor starting with a rough block of marble...",
                      font_size=18, font_color=COLORS["dark_blue"], bold=True)
    add_styled_textbox(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(0.5),
                      "They chip away noise until a beautiful statue appears. AI image models do the same with pixels!",
                      font_size=16, font_color=COLORS["dark_gray"])

    # Steps
    steps = [
        ("1. Start with\nRandom Noise", COLORS["dark_gray"]),
        ("2. AI Removes Noise\n(Many Steps)", COLORS["diffusion_blue"]),
        ("3. Text Prompt\nGuides the Shape", COLORS["prompt_gold"]),
        ("4. Final Image\nAppears!", COLORS["green"]),
    ]
    for i, (text, color) in enumerate(steps):
        x = 0.6 + i * 3.2
        add_rounded_box(slide, Inches(x), Inches(3.3), Inches(2.9), Inches(1.3),
                       color, text=text, text_size=14, text_color=COLORS["white"])
        if i < 3:
            add_arrow(slide, Inches(x + 3.0), Inches(3.8), Inches(0.5), Inches(0.3),
                     COLORS["orange"])

    add_takeaway_bar(slide,
        "Diffusion = Start with noise, subtract it step by step, guided by your text prompt")


def slide_how_it_works_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Text-to-Image: Step by Step")

    # Flow diagram
    flow = [
        ("Your Prompt", "A dragon reading\na book in a library", COLORS["prompt_gold"]),
        ("Text Encoder", "Converts words into\nnumbers AI understands", COLORS["medium_blue"]),
        ("Diffusion Model", "Removes noise step\nby step (50-100 steps)", COLORS["diffusion_blue"]),
        ("Final Image", "A brand new picture\nthat never existed!", COLORS["green"]),
    ]
    for i, (label, desc, color) in enumerate(flow):
        y = 1.7 + i * 1.3
        add_rounded_box(slide, Inches(1.0), Inches(y), Inches(3.0), Inches(0.9),
                       color, text=label, text_size=16)
        add_styled_textbox(slide, Inches(4.3), Inches(y + 0.1), Inches(8.0), Inches(0.8),
                          desc, font_size=15, font_color=COLORS["dark_gray"])
        if i < 3:
            add_arrow(slide, Inches(2.3), Inches(y + 0.95), Inches(0.4), Inches(0.25),
                     COLORS["orange"], direction="down")

    add_styled_textbox(
        slide, Inches(1.0), Inches(6.5), Inches(11.0), Inches(0.5),
        "Key: The image is NOT copied from anywhere. It's generated from learned patterns!",
        font_size=14, font_color=COLORS["dark_blue"], bold=True
    )


def slide_multimodal(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Multimodal AI: Beyond Just Text")

    add_styled_textbox(
        slide, Inches(0.5), Inches(1.6), Inches(12.5), Inches(0.5),
        "Multimodal = AI that understands AND generates MULTIPLE types of data",
        font_size=18, font_color=COLORS["dark_gray"], bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # Modalities
    modes = [
        ("Text", "Read & write\nwords", COLORS["medium_blue"]),
        ("Images", "See & create\npictures", COLORS["image_purple"]),
        ("Audio", "Hear & speak\n(voice)", COLORS["multimodal_teal"]),
        ("Video", "Watch & make\nclips", COLORS["creative_pink"]),
    ]
    for i, (label, desc, color) in enumerate(modes):
        x = 0.8 + i * 3.2
        add_circle(slide, Inches(x + 0.6), Inches(2.4), Inches(1.5), color,
                  text=label, text_size=16)
        add_styled_textbox(slide, Inches(x), Inches(4.1), Inches(2.8), Inches(0.6),
                          desc, font_size=13, font_color=COLORS["dark_gray"],
                          alignment=PP_ALIGN.CENTER)

    # Examples
    add_rounded_box(slide, Inches(0.5), Inches(5.0), Inches(12.33), Inches(1.2),
                   COLORS["sky_blue"])
    add_styled_textbox(slide, Inches(0.8), Inches(5.1), Inches(11.7), Inches(0.4),
                      "Examples You Can Try Today:", font_size=16,
                      font_color=COLORS["dark_blue"], bold=True)
    add_styled_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.6),
                      "GPT-4o: Upload a photo and ask questions about it  |  "
                      "Gemini: Describe an image and it explains what's in it  |  "
                      "Claude: Analyze charts, screenshots, and documents",
                      font_size=13, font_color=COLORS["dark_gray"])


def slide_prompting_images(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Prompting for Images vs. Text")

    # Comparison table
    add_rounded_box(slide, Inches(0.5), Inches(1.7), Inches(6.0), Inches(0.6),
                   COLORS["medium_blue"], text="Text Prompts (what you know)",
                   text_size=15)
    add_rounded_box(slide, Inches(6.8), Inches(1.7), Inches(6.0), Inches(0.6),
                   COLORS["image_purple"], text="Image Prompts (new today!)",
                   text_size=15)

    text_tips = [
        "Role, Task, Context, Format",
        "Be specific about what you want",
        "Chain-of-thought for reasoning",
        "More words = more detail",
    ]
    image_tips = [
        "Subject, Style, Lighting, Mood",
        "Describe what the image LOOKS like",
        "Mention art style (cartoon, realistic, pixel art)",
        "More descriptive = more detailed image",
    ]

    for i, (t, im) in enumerate(zip(text_tips, image_tips)):
        y = 2.5 + i * 0.55
        add_styled_textbox(slide, Inches(0.7), Inches(y), Inches(5.8), Inches(0.5),
                          f"• {t}", font_size=14, font_color=COLORS["dark_gray"])
        add_styled_textbox(slide, Inches(7.0), Inches(y), Inches(5.8), Inches(0.5),
                          f"• {im}", font_size=14, font_color=COLORS["dark_gray"])

    # Example
    add_rounded_box(slide, Inches(0.5), Inches(4.9), Inches(12.33), Inches(1.8),
                   COLORS["light_gray"], border_color=COLORS["medium_blue"])
    add_styled_textbox(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(0.4),
                      "Bad image prompt:", font_size=14,
                      font_color=COLORS["red"], bold=True)
    add_styled_textbox(slide, Inches(0.8), Inches(5.35), Inches(5.5), Inches(0.4),
                      '"Make me a dog picture"', font_size=13,
                      font_color=COLORS["dark_gray"])
    add_styled_textbox(slide, Inches(6.8), Inches(5.0), Inches(5.8), Inches(0.4),
                      "Great image prompt:", font_size=14,
                      font_color=COLORS["green"], bold=True)
    add_styled_textbox(slide, Inches(6.8), Inches(5.35), Inches(5.8), Inches(0.8),
                      '"A golden retriever puppy sitting in a field of sunflowers,\n'
                      'watercolor painting style, warm sunset lighting, happy mood"',
                      font_size=13, font_color=COLORS["dark_gray"])


def slide_demo_intro(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Live Demo Time!")

    add_styled_textbox(
        slide, Inches(0.5), Inches(1.8), Inches(12.33), Inches(0.6),
        "Let's generate images together and see what happens!",
        font_size=22, font_color=COLORS["dark_gray"], bold=True,
        alignment=PP_ALIGN.CENTER
    )

    demos = [
        ("Demo 1", "Same prompt, different tools\n(ChatGPT vs. Gemini)", COLORS["diffusion_blue"]),
        ("Demo 2", "Vague prompt vs. detailed prompt\n(see the quality difference!)", COLORS["prompt_gold"]),
        ("Demo 3", "Upload a photo and ask AI about it\n(multimodal in action)", COLORS["multimodal_teal"]),
        ("Demo 4", "Try to make AI generate something\nimpossible or weird", COLORS["creative_pink"]),
    ]
    for i, (label, desc, color) in enumerate(demos):
        y = 2.6 + i * 1.1
        add_rounded_box(slide, Inches(1.0), Inches(y), Inches(2.2), Inches(0.85),
                       color, text=label, text_size=16)
        add_styled_textbox(slide, Inches(3.5), Inches(y + 0.1), Inches(9.0), Inches(0.75),
                          desc, font_size=15, font_color=COLORS["dark_gray"])


def slide_ethics_deepfakes(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "The Dark Side: Deepfakes & Misuse",
                 subtitle="With great power comes great responsibility")

    concerns = [
        ("Deepfakes", "Fake photos/videos of real people — used for bullying or scams",
         COLORS["ethics_red"]),
        ("Misinformation", "AI-generated images can spread false news and fool people",
         COLORS["creative_pink"]),
        ("Artist Concerns", "AI trained on artists' work without permission — is that fair?",
         COLORS["image_purple"]),
        ("Identity Theft", "Fake profile pictures, impersonation, fraud",
         COLORS["dark_gray"]),
    ]
    for i, (title, desc, color) in enumerate(concerns):
        y = 1.7 + i * 1.15
        add_rounded_box(slide, Inches(0.5), Inches(y), Inches(2.5), Inches(0.85),
                       color, text=title, text_size=14)
        add_styled_textbox(slide, Inches(3.3), Inches(y + 0.15), Inches(9.5), Inches(0.7),
                          desc, font_size=14, font_color=COLORS["dark_gray"])

    add_takeaway_bar(slide,
        "RULE: Never create fake images of real people. Always think before you share.")


def slide_how_to_spot(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "How to Spot AI-Generated Images")

    tips = [
        ("Hands & Fingers", "Often too many or too few fingers, weird poses"),
        ("Text in Images", "AI-generated text is usually garbled or misspelled"),
        ("Background Details", "Look for melting objects, impossible architecture"),
        ("Too Perfect", "Skin too smooth, lighting too even, no imperfections"),
        ("Context Clues", "Check the source — who shared it and why?"),
    ]
    for i, (tip, desc) in enumerate(tips):
        y = 1.7 + i * 1.0
        color = [COLORS["diffusion_blue"], COLORS["image_purple"],
                COLORS["multimodal_teal"], COLORS["prompt_gold"],
                COLORS["ethics_red"]][i]
        add_circle(slide, Inches(0.7), Inches(y), Inches(0.7), color,
                  text=str(i+1), text_size=18)
        add_styled_textbox(slide, Inches(1.7), Inches(y + 0.05), Inches(2.5), Inches(0.5),
                          tip, font_size=15, font_color=COLORS["dark_blue"], bold=True)
        add_styled_textbox(slide, Inches(4.3), Inches(y + 0.05), Inches(8.5), Inches(0.5),
                          desc, font_size=14, font_color=COLORS["dark_gray"])

    add_takeaway_bar(slide, "When in doubt, reverse image search or check the source!")


def slide_key_vocab(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Key Vocabulary")

    vocab = [
        ("Text-to-Image", "AI model that creates pictures from text descriptions"),
        ("Diffusion", "The process of removing noise step by step to create an image"),
        ("Multimodal", "AI that works with multiple types of data (text + image + audio)"),
        ("Deepfake", "AI-generated fake image or video of a real person"),
        ("Image Prompt", "The text description you give to an image generator"),
        ("Style Transfer", "Applying one artistic style to another image"),
    ]
    for i, (term, definition) in enumerate(vocab):
        y = 1.6 + i * 0.85
        color = [COLORS["image_purple"], COLORS["diffusion_blue"],
                COLORS["multimodal_teal"], COLORS["ethics_red"],
                COLORS["prompt_gold"], COLORS["creative_pink"]][i]
        add_rounded_box(slide, Inches(0.5), Inches(y), Inches(2.8), Inches(0.65),
                       color, text=term, text_size=13)
        add_styled_textbox(slide, Inches(3.6), Inches(y + 0.1), Inches(9.2), Inches(0.55),
                          definition, font_size=14, font_color=COLORS["dark_gray"])


def slide_takeaways(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Today's Key Takeaways")

    takeaways = [
        "AI can generate images from text using diffusion (noise removal)",
        "Multimodal AI understands text, images, audio, and video together",
        "Better prompts = better images (describe subject, style, lighting, mood)",
        "AI images are NOT copied — they're generated from learned patterns",
        "Deepfakes are dangerous — never make fake images of real people",
        "Always check if an image is AI-generated before trusting or sharing it",
    ]
    for i, text in enumerate(takeaways):
        y = 1.6 + i * 0.85
        colors = [COLORS["medium_blue"], COLORS["multimodal_teal"],
                 COLORS["prompt_gold"], COLORS["image_purple"],
                 COLORS["ethics_red"], COLORS["green"]]
        add_circle(slide, Inches(0.8), Inches(y + 0.05), Inches(0.55),
                  colors[i], text=str(i+1), text_size=16)
        add_styled_textbox(slide, Inches(1.7), Inches(y + 0.1), Inches(11.0), Inches(0.6),
                          text, font_size=17, font_color=COLORS["dark_gray"])


def slide_questions(prs):
    create_questions_slide(prs, "Any questions before the Kahoot?")


# =============================================================================
# MAIN
# =============================================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_agenda(prs)
    slide_recap(prs)
    slide_what_are_image_models(prs)
    slide_how_it_works_1(prs)
    slide_how_it_works_2(prs)
    slide_multimodal(prs)
    slide_prompting_images(prs)
    slide_demo_intro(prs)
    slide_ethics_deepfakes(prs)
    slide_how_to_spot(prs)
    slide_key_vocab(prs)
    slide_takeaways(prs)
    slide_questions(prs)

    script_dir = os.path.dirname(os.path.abspath(__file__))
    parent_dir = os.path.dirname(script_dir)
    output = os.path.join(parent_dir, "2026-05-23_AI_Image_and_Multimodal_Models.pptx")
    prs.save(output)
    print(f"[SUCCESS] Presentation saved: {output}")


if __name__ == "__main__":
    main()
