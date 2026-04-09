"""
EverestIT Class Theme — Standard Edition
Reusable helpers for creating clean, colorful PowerPoint slides.

Style: Dark blue title bars, orange accents, colorful rounded boxes,
white backgrounds. Matches the 2/07 and 2/28 lesson style.

Usage:
    import sys, os
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', 'tools'))
    from create_theme import *

Dependencies:
    pip install python-pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ============================================================================
# COLOR PALETTE
# ============================================================================

COLORS = {
    # Primary blues
    "dark_blue": RGBColor(30, 58, 95),
    "medium_blue": RGBColor(52, 103, 166),
    "light_blue": RGBColor(100, 149, 237),
    "sky_blue": RGBColor(173, 216, 230),

    # Accents
    "orange": RGBColor(255, 153, 51),
    "green": RGBColor(76, 175, 80),
    "emerald": RGBColor(46, 125, 50),
    "red": RGBColor(244, 67, 54),
    "purple": RGBColor(156, 39, 176),
    "teal": RGBColor(0, 150, 136),
    "violet": RGBColor(103, 58, 183),
    "amber": RGBColor(255, 160, 0),
    "coral": RGBColor(255, 107, 107),
    "pink": RGBColor(233, 30, 99),

    # Neutrals
    "white": RGBColor(255, 255, 255),
    "light_gray": RGBColor(245, 245, 245),
    "dark_gray": RGBColor(66, 66, 66),
    "black": RGBColor(33, 33, 33),
}


# ============================================================================
# CORE SHAPE HELPERS
# ============================================================================

def set_shape_fill(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_shape_line(shape, color, width_pt=1):
    line = shape.line
    line.color.rgb = color
    line.width = Pt(width_pt)


def add_styled_textbox(slide, left, top, width, height, text,
                       font_size=18, font_color=None, bold=False,
                       alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Add a text box with consistent styling."""
    if font_color is None:
        font_color = COLORS["dark_gray"]

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment

    return textbox


def add_rounded_box(slide, left, top, width, height, fill_color, text=None,
                    text_size=16, text_color=None, border_color=None):
    """Add a rounded rectangle with optional text and border."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    set_shape_fill(shape, fill_color)

    if border_color:
        set_shape_line(shape, border_color, 2)
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(text_size)
        p.font.color.rgb = text_color or COLORS["white"]
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(8)

    shape.text_frame.anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_circle(slide, left, top, size, fill_color, text=None,
               text_size=14, text_color=None):
    """Add a circle shape with optional centered text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left, top, size, size
    )
    set_shape_fill(shape, fill_color)
    shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(text_size)
        p.font.color.rgb = text_color or COLORS["white"]
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
    shape.text_frame.anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_arrow(slide, left, top, width, height, fill_color, direction="right"):
    """Add a directional arrow shape."""
    shape_map = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
    }
    shape = slide.shapes.add_shape(
        shape_map.get(direction, MSO_SHAPE.RIGHT_ARROW),
        left, top, width, height
    )
    set_shape_fill(shape, fill_color)
    shape.line.fill.background()
    return shape


def add_table_row(slide, y, cols, col_widths, col_starts, bg_color,
                  text_color=None, font_size=14, bold=False, row_height=0.55):
    """Add a row of boxes that looks like a table row."""
    if text_color is None:
        text_color = COLORS["white"]
    for text, w, x in zip(cols, col_widths, col_starts):
        add_rounded_box(
            slide, Inches(x), Inches(y), Inches(w), Inches(row_height),
            bg_color, text, text_size=font_size, text_color=text_color
        )


# ============================================================================
# SLIDE COMPONENT HELPERS
# ============================================================================

def add_title_bar(slide, text, subtitle=None):
    """Dark blue title bar with orange accent line."""
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(1.3)
    )
    set_shape_fill(title_bar, COLORS["dark_blue"])
    title_bar.line.fill.background()

    add_styled_textbox(
        slide, Inches(0.5), Inches(0.35), Inches(12.33), Inches(0.7),
        text, font_size=36, font_color=COLORS["white"], bold=True,
        alignment=PP_ALIGN.LEFT, font_name="Segoe UI Semibold"
    )

    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.3),
        Inches(13.33), Inches(0.08)
    )
    set_shape_fill(accent_line, COLORS["orange"])
    accent_line.line.fill.background()

    if subtitle:
        add_styled_textbox(
            slide, Inches(0.5), Inches(0.95), Inches(12.33), Inches(0.4),
            subtitle, font_size=16, font_color=COLORS["sky_blue"],
            alignment=PP_ALIGN.LEFT
        )


def add_agenda_item(slide, num, text, y, color):
    """Colored circle with number + text label for agenda slides."""
    add_circle(slide, Inches(1), Inches(y), Inches(0.55), color,
               text=str(num), text_size=20)
    add_styled_textbox(
        slide, Inches(2), Inches(y + 0.1), Inches(10.5), Inches(0.5),
        text, font_size=20, font_color=COLORS["dark_gray"]
    )


def add_takeaway_bar(slide, text, color=None):
    """Colored rounded box at the bottom of a slide for key message."""
    if color is None:
        color = COLORS["dark_blue"]
    add_rounded_box(
        slide, Inches(1.0), Inches(6.3), Inches(11.33), Inches(0.7),
        color, text=text, text_size=18, text_color=COLORS["white"]
    )


def add_full_bg(slide, color=None):
    """Full-slide solid color background."""
    if color is None:
        color = COLORS["dark_blue"]
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    set_shape_fill(bg, color)
    bg.line.fill.background()
    return bg


# ============================================================================
# COMMON SLIDE TEMPLATES
# ============================================================================

def create_title_slide(prs, title, subtitle, date, tagline=None):
    """Standard title slide: dark blue bg, light blue center box."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_full_bg(slide, COLORS["dark_blue"])

    # Count newlines to adjust spacing for multi-line titles
    title_lines = title.count('\n') + 1
    box_h = 2.4 + (title_lines - 1) * 0.7
    subtitle_y = 1.7 + box_h - 1.1

    add_rounded_box(
        slide, Inches(1.2), Inches(1.7), Inches(10.9), Inches(box_h),
        COLORS["light_blue"]
    )

    add_styled_textbox(
        slide, Inches(0.6), Inches(2.0), Inches(12.13), Inches(1.2),
        title, font_size=44, font_color=COLORS["dark_blue"], bold=True,
        alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold"
    )

    add_styled_textbox(
        slide, Inches(0.6), Inches(subtitle_y), Inches(12.13), Inches(0.6),
        subtitle, font_size=22, font_color=COLORS["dark_blue"],
        alignment=PP_ALIGN.CENTER
    )

    add_styled_textbox(
        slide, Inches(0.6), Inches(6.6), Inches(12.13), Inches(0.4),
        date, font_size=18, font_color=COLORS["sky_blue"],
        alignment=PP_ALIGN.CENTER
    )

    if tagline:
        # Orange accent line
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(3.0), Inches(5.2), Inches(7.33), Inches(0.06)
        )
        set_shape_fill(accent, COLORS["orange"])
        accent.line.fill.background()

        add_styled_textbox(
            slide, Inches(2.0), Inches(5.6), Inches(9.33), Inches(0.5),
            tagline, font_size=16, font_color=COLORS["sky_blue"],
            alignment=PP_ALIGN.CENTER
        )

    return slide


def create_questions_slide(prs, prompt_text="Ask away before we start the Kahoot!"):
    """Standard questions slide: dark blue bg, big ? mark."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_full_bg(slide, COLORS["dark_blue"])

    add_styled_textbox(
        slide, Inches(0), Inches(1.6), Inches(13.33), Inches(2),
        "?", font_size=120, font_color=COLORS["white"],
        alignment=PP_ALIGN.CENTER
    )

    add_styled_textbox(
        slide, Inches(0), Inches(3.5), Inches(13.33), Inches(1),
        "Questions?", font_size=60, font_color=COLORS["white"], bold=True,
        alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold"
    )

    add_styled_textbox(
        slide, Inches(0), Inches(5.1), Inches(13.33), Inches(0.6),
        prompt_text, font_size=24, font_color=COLORS["sky_blue"],
        alignment=PP_ALIGN.CENTER
    )

    return slide


# ============================================================================
# DEMO PRESENTATION
# ============================================================================

def main():
    script_dir = os.path.dirname(os.path.abspath(__file__))
    project_dir = os.path.dirname(script_dir)
    theme_dir = os.path.join(project_dir, "theme")
    os.makedirs(theme_dir, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    create_title_slide(prs,
        "EverestIT Class Theme",
        "Standard Edition — Clean & Colorful",
        "Theme Demo",
        tagline="Kids Computer Science Class  |  EverestIT"
    )

    # Slide 2: Agenda demo
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Today's Journey")
    items = [
        ("Title Bars & Accent Lines", COLORS["light_blue"]),
        ("Rounded Boxes & Colors", COLORS["green"]),
        ("Agenda Items", COLORS["orange"]),
        ("Content Cards (2x2 Grid)", COLORS["purple"]),
        ("Takeaway Bars", COLORS["teal"]),
        ("Questions Slide", COLORS["medium_blue"]),
    ]
    for i, (text, color) in enumerate(items):
        add_agenda_item(slide, i + 1, text, 1.6 + i * 0.68, color)

    # Slide 3: Color palette demo
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Color Palette", "All available accent colors")

    palette = [
        ("dark_blue", COLORS["dark_blue"]),
        ("medium_blue", COLORS["medium_blue"]),
        ("light_blue", COLORS["light_blue"]),
        ("orange", COLORS["orange"]),
        ("green", COLORS["green"]),
        ("red", COLORS["red"]),
        ("purple", COLORS["purple"]),
        ("teal", COLORS["teal"]),
        ("violet", COLORS["violet"]),
        ("pink", COLORS["pink"]),
    ]

    for i, (name, color) in enumerate(palette):
        x = Inches(0.5 + (i % 5) * 2.5)
        y = Inches(1.7 + (i // 5) * 2.2)
        add_rounded_box(slide, x, y, Inches(2.2), Inches(1.8), color,
                        text=name, text_size=14)

    add_takeaway_bar(slide, "Mix and match these colors for unique slides!",
                     COLORS["orange"])

    # Slide 4: 2x2 grid demo
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Content Cards (2x2 Grid)")

    cards = [
        ("Topic One", "Description of the first concept",
         COLORS["medium_blue"]),
        ("Topic Two", "Description of the second concept",
         COLORS["purple"]),
        ("Topic Three", "Description of the third concept",
         COLORS["teal"]),
        ("Topic Four", "Description of the fourth concept",
         COLORS["orange"]),
    ]

    for i, (title, desc, color) in enumerate(cards):
        x = Inches(0.7 + (i % 2) * 6.2)
        y = Inches(1.7 + (i // 2) * 2.0)
        add_rounded_box(slide, x, y, Inches(5.8), Inches(1.7), color)
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.2),
            Inches(5.4), Inches(0.5),
            title, font_size=22, font_color=COLORS["white"], bold=True,
            alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.85),
            Inches(5.4), Inches(0.5),
            desc, font_size=16, font_color=COLORS["light_gray"],
            alignment=PP_ALIGN.CENTER
        )

    add_takeaway_bar(slide, "Use 2x2 grids for comparing or grouping concepts")

    # Slide 5: Questions
    create_questions_slide(prs, "Let's discuss what you've learned!")

    demo_path = os.path.join(theme_dir, "EverestIT_Theme_Demo.pptx")
    prs.save(demo_path)
    print(f"[SUCCESS] Theme demo saved to: {demo_path}")
    print(f"  Slides: {len(prs.slides)}")
    print(f"\nImport with: from create_theme import *")


if __name__ == "__main__":
    main()
