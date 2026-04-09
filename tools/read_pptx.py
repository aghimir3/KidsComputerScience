"""
PowerPoint Reader Utility
Extracts text content from .pptx files for easy reference.

Usage:
    python read_pptx.py <filename.pptx>
    python read_pptx.py                     # Lists all .pptx files and lets you choose
    python read_pptx.py --all               # Reads all .pptx files in the folder

Output is saved to a .txt file with the same name.
"""

import os
import sys
from pptx import Presentation


def extract_text_from_pptx(filepath):
    """Extract all text content from a PowerPoint file."""
    prs = Presentation(filepath)
    output = []

    filename = os.path.basename(filepath)
    output.append("=" * 60)
    output.append(f"FILE: {filename}")
    output.append(f"SLIDES: {len(prs.slides)}")
    output.append("=" * 60)
    output.append("")

    for i, slide in enumerate(prs.slides, 1):
        output.append(f"{'='*20} SLIDE {i} {'='*20}")

        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                # Clean the text (remove emojis for cleaner output)
                text = shape.text.strip()
                # Keep the text as-is but handle encoding issues
                try:
                    clean_text = text.encode('utf-8').decode('utf-8')
                except:
                    clean_text = text.encode('ascii', 'ignore').decode('ascii')

                if clean_text:
                    output.append(clean_text)

        output.append("")

    return "\n".join(output)


def list_pptx_files(folder):
    """List all .pptx files in a folder."""
    files = []
    for f in os.listdir(folder):
        if f.endswith('.pptx') and not f.startswith('~$'):
            files.append(f)
    return sorted(files)


def save_output(content, original_filepath):
    """Save extracted content to a text file."""
    base = os.path.splitext(original_filepath)[0]
    output_path = f"{base}_content.txt"

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(content)

    return output_path


def main():
    folder = os.path.dirname(os.path.abspath(__file__))

    if len(sys.argv) > 1:
        arg = sys.argv[1]

        if arg == '--all':
            # Read all pptx files
            files = list_pptx_files(folder)
            if not files:
                print("No .pptx files found in this folder.")
                return

            print(f"Found {len(files)} PowerPoint files.\n")

            for filename in files:
                filepath = os.path.join(folder, filename)
                print(f"Reading: {filename}...")

                try:
                    content = extract_text_from_pptx(filepath)
                    output_path = save_output(content, filepath)
                    print(f"  Saved to: {os.path.basename(output_path)}")
                except Exception as e:
                    print(f"  Error: {e}")

            print("\nDone!")

        elif arg == '--help' or arg == '-h':
            print(__doc__)

        else:
            # Read specific file
            if not arg.endswith('.pptx'):
                arg += '.pptx'

            filepath = os.path.join(folder, arg) if not os.path.isabs(arg) else arg

            if not os.path.exists(filepath):
                print(f"File not found: {filepath}")
                return

            print(f"Reading: {arg}...")
            content = extract_text_from_pptx(filepath)

            # Print to console
            print("\n" + content)

            # Save to file
            output_path = save_output(content, filepath)
            print(f"\nSaved to: {output_path}")

    else:
        # Interactive mode - list files and let user choose
        files = list_pptx_files(folder)

        if not files:
            print("No .pptx files found in this folder.")
            return

        print("Available PowerPoint files:\n")
        for i, f in enumerate(files, 1):
            print(f"  {i}. {f}")

        print(f"\n  A. Read ALL files")
        print(f"  Q. Quit\n")

        choice = input("Enter number or letter: ").strip()

        if choice.upper() == 'Q':
            return
        elif choice.upper() == 'A':
            for filename in files:
                filepath = os.path.join(folder, filename)
                print(f"\nReading: {filename}...")

                try:
                    content = extract_text_from_pptx(filepath)
                    output_path = save_output(content, filepath)
                    print(f"  Saved to: {os.path.basename(output_path)}")
                except Exception as e:
                    print(f"  Error: {e}")
        else:
            try:
                idx = int(choice) - 1
                if 0 <= idx < len(files):
                    filename = files[idx]
                    filepath = os.path.join(folder, filename)

                    print(f"\nReading: {filename}...\n")
                    content = extract_text_from_pptx(filepath)

                    # Print to console
                    print(content)

                    # Save to file
                    output_path = save_output(content, filepath)
                    print(f"\nSaved to: {output_path}")
                else:
                    print("Invalid selection.")
            except ValueError:
                print("Invalid input.")


if __name__ == "__main__":
    main()
