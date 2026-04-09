"""
PowerPoint Layout Validator
Checks for overlapping elements, out-of-bounds shapes, text overflow,
and common layout issues in .pptx files.

Usage:
    python test_pptx_layout.py <pptx_file>   # Test specific file
    python test_pptx_layout.py --all         # Test all .pptx files in folder
    python test_pptx_layout.py               # Interactive mode

Dependencies:
    pip install python-pptx
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Emu


# ============================================================================
# CONSTANTS
# ============================================================================

# Standard 16:9 widescreen dimensions
SLIDE_WIDTH_IN = 13.33
SLIDE_HEIGHT_IN = 7.5

# Margins (inches) — elements beyond these are flagged
MARGIN = 0.0  # Allow shapes right to the edge (many designs use full-bleed)
OVERLAP_THRESHOLD_IN = 0.05  # Ignore overlaps smaller than this (rounding)

# Text overflow estimation
AVG_CHAR_WIDTH_PT = {
    10: 5.5, 11: 6.0, 12: 6.5, 13: 7.0, 14: 7.5,
    15: 8.0, 16: 8.5, 17: 9.0, 18: 9.5, 20: 10.5,
    22: 11.5, 24: 12.5, 26: 13.5, 28: 14.5, 30: 15.5,
    32: 16.5, 36: 18.5, 40: 21.0, 44: 23.0, 48: 25.0,
    60: 31.0, 120: 62.0,
}
AVG_LINE_HEIGHT_PT = {
    10: 14, 11: 15, 12: 16, 13: 18, 14: 19,
    15: 20, 16: 22, 17: 23, 18: 24, 20: 27,
    22: 30, 24: 32, 26: 35, 28: 38, 30: 40,
    32: 43, 36: 48, 40: 54, 44: 59, 48: 64,
    60: 80, 120: 160,
}


def emu_to_inches(emu):
    """Convert EMU to inches."""
    return emu / 914400


def get_char_width(font_size_pt):
    """Estimate average character width for a font size."""
    if font_size_pt in AVG_CHAR_WIDTH_PT:
        return AVG_CHAR_WIDTH_PT[font_size_pt]
    # Interpolate
    sizes = sorted(AVG_CHAR_WIDTH_PT.keys())
    for i in range(len(sizes) - 1):
        if sizes[i] <= font_size_pt <= sizes[i + 1]:
            ratio = (font_size_pt - sizes[i]) / (sizes[i + 1] - sizes[i])
            return AVG_CHAR_WIDTH_PT[sizes[i]] + ratio * (
                AVG_CHAR_WIDTH_PT[sizes[i + 1]] - AVG_CHAR_WIDTH_PT[sizes[i]])
    return font_size_pt * 0.52  # fallback


def get_line_height(font_size_pt):
    """Estimate line height for a font size."""
    if font_size_pt in AVG_LINE_HEIGHT_PT:
        return AVG_LINE_HEIGHT_PT[font_size_pt]
    sizes = sorted(AVG_LINE_HEIGHT_PT.keys())
    for i in range(len(sizes) - 1):
        if sizes[i] <= font_size_pt <= sizes[i + 1]:
            ratio = (font_size_pt - sizes[i]) / (sizes[i + 1] - sizes[i])
            return AVG_LINE_HEIGHT_PT[sizes[i]] + ratio * (
                AVG_LINE_HEIGHT_PT[sizes[i + 1]] - AVG_LINE_HEIGHT_PT[sizes[i]])
    return font_size_pt * 1.35


# ============================================================================
# SHAPE INFO EXTRACTION
# ============================================================================

class ShapeInfo:
    """Extracted position and text info for a shape."""

    def __init__(self, shape, index):
        self.index = index
        self.name = shape.name or f"Shape {index}"
        self.shape_type = str(shape.shape_type)

        self.left = emu_to_inches(shape.left) if shape.left is not None else 0
        self.top = emu_to_inches(shape.top) if shape.top is not None else 0
        self.width = emu_to_inches(shape.width) if shape.width is not None else 0
        self.height = emu_to_inches(shape.height) if shape.height is not None else 0

        self.right = self.left + self.width
        self.bottom = self.top + self.height

        # Extract text info
        self.text = ""
        self.font_size = None
        self.has_text = False

        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.text:
                    self.text += para.text + "\n"
                    for run in para.runs:
                        if run.font.size is not None:
                            self.font_size = run.font.size.pt
            self.text = self.text.strip()
            self.has_text = bool(self.text)

    def overlaps(self, other, threshold=OVERLAP_THRESHOLD_IN):
        """Check if this shape overlaps with another (beyond threshold)."""
        h_overlap = min(self.right, other.right) - max(self.left, other.left)
        v_overlap = min(self.bottom, other.bottom) - max(self.top, other.top)
        if h_overlap > threshold and v_overlap > threshold:
            return True
        return False

    def overlap_area(self, other):
        """Calculate overlap area in square inches."""
        h_overlap = max(0, min(self.right, other.right) - max(self.left, other.left))
        v_overlap = max(0, min(self.bottom, other.bottom) - max(self.top, other.top))
        return h_overlap * v_overlap

    def __repr__(self):
        return (f"{self.name} @ ({self.left:.2f}, {self.top:.2f}) "
                f"size ({self.width:.2f} x {self.height:.2f})")


# ============================================================================
# VALIDATORS
# ============================================================================

class LayoutValidator:
    """Validates layout of a single PowerPoint file."""

    def __init__(self, pptx_path):
        self.pptx_path = pptx_path
        self.prs = Presentation(pptx_path)
        self.errors = []
        self.warnings = []

        # Get actual slide dimensions
        self.slide_w = emu_to_inches(self.prs.slide_width)
        self.slide_h = emu_to_inches(self.prs.slide_height)

    def validate(self):
        """Run all validations."""
        print(f"\n{'=' * 60}")
        print(f"VALIDATING: {os.path.basename(self.pptx_path)}")
        print(f"{'=' * 60}")
        print(f"  Slides: {len(self.prs.slides)}")
        print(f"  Dimensions: {self.slide_w:.2f}\" x {self.slide_h:.2f}\"")

        for slide_num, slide in enumerate(self.prs.slides, 1):
            self._validate_slide(slide_num, slide)

        self._print_results()
        return len(self.errors) == 0

    def _validate_slide(self, slide_num, slide):
        """Validate a single slide."""
        shapes = []
        for i, shape in enumerate(slide.shapes):
            try:
                info = ShapeInfo(shape, i)
                shapes.append(info)
            except Exception:
                pass  # Skip shapes we can't parse

        # Check bounds
        self._check_bounds(slide_num, shapes)

        # Check text-on-text overlaps
        self._check_text_overlaps(slide_num, shapes)

        # Check text overflow
        self._check_text_overflow(slide_num, shapes)

        # Check zero-size shapes
        self._check_zero_size(slide_num, shapes)

    def _check_bounds(self, slide_num, shapes):
        """Check if shapes are within slide boundaries."""
        for s in shapes:
            # Allow small overflow (decorative elements often go slightly off-edge)
            overflow_threshold = 0.5  # inches

            if s.right > self.slide_w + overflow_threshold:
                self.errors.append(
                    f"Slide {slide_num}: {s.name} extends {s.right - self.slide_w:.2f}\" "
                    f"past right edge (right={s.right:.2f}\", slide width={self.slide_w:.2f}\")")

            if s.bottom > self.slide_h + overflow_threshold:
                self.errors.append(
                    f"Slide {slide_num}: {s.name} extends {s.bottom - self.slide_h:.2f}\" "
                    f"past bottom edge (bottom={s.bottom:.2f}\", slide height={self.slide_h:.2f}\")")

            if s.left < -overflow_threshold:
                self.warnings.append(
                    f"Slide {slide_num}: {s.name} starts {abs(s.left):.2f}\" "
                    f"past left edge (left={s.left:.2f}\")")

            if s.top < -overflow_threshold:
                self.warnings.append(
                    f"Slide {slide_num}: {s.name} starts {abs(s.top):.2f}\" "
                    f"past top edge (top={s.top:.2f}\")")

    def _check_text_overlaps(self, slide_num, shapes):
        """Check if text-containing shapes overlap other text shapes."""
        text_shapes = [s for s in shapes if s.has_text]

        for i in range(len(text_shapes)):
            for j in range(i + 1, len(text_shapes)):
                a = text_shapes[i]
                b = text_shapes[j]

                if a.overlaps(b):
                    area = a.overlap_area(b)
                    # Only flag significant overlaps
                    smaller_area = min(a.width * a.height, b.width * b.height)
                    if smaller_area > 0 and area / smaller_area > 0.3:
                        # Major overlap — more than 30% of the smaller shape
                        a_text = a.text[:40] + "..." if len(a.text) > 40 else a.text
                        b_text = b.text[:40] + "..." if len(b.text) > 40 else b.text
                        self.errors.append(
                            f"Slide {slide_num}: Major text overlap "
                            f"({area / smaller_area:.0%} of smaller shape)\n"
                            f"      Shape A: \"{a_text}\" @ y={a.top:.2f}\"-{a.bottom:.2f}\"\n"
                            f"      Shape B: \"{b_text}\" @ y={b.top:.2f}\"-{b.bottom:.2f}\"")
                    elif area > 0.1:
                        # Minor but notable overlap
                        a_text = a.text[:30] + "..." if len(a.text) > 30 else a.text
                        b_text = b.text[:30] + "..." if len(b.text) > 30 else b.text
                        self.warnings.append(
                            f"Slide {slide_num}: Minor text overlap ({area:.2f} sq in)\n"
                            f"      \"{a_text}\" & \"{b_text}\"")

    def _check_text_overflow(self, slide_num, shapes):
        """Estimate if text might overflow its container."""
        for s in shapes:
            if not s.has_text or not s.font_size:
                continue
            if s.width < 0.1 or s.height < 0.1:
                continue

            font_pt = s.font_size
            char_w_pt = get_char_width(font_pt)
            line_h_pt = get_line_height(font_pt)

            # Available width in points (1 inch = 72 points)
            avail_w_pt = s.width * 72
            avail_h_pt = s.height * 72

            # Estimate lines needed
            lines = s.text.split('\n')
            total_lines = 0
            for line in lines:
                line_width_pt = len(line) * char_w_pt
                wrapped_lines = max(1, int(line_width_pt / avail_w_pt) + 1)
                total_lines += wrapped_lines

            needed_h_pt = total_lines * line_h_pt

            if needed_h_pt > avail_h_pt * 1.3:  # 30% overflow tolerance
                overflow_pct = (needed_h_pt / avail_h_pt - 1) * 100
                text_preview = s.text[:50] + "..." if len(s.text) > 50 else s.text
                self.warnings.append(
                    f"Slide {slide_num}: Text may overflow by ~{overflow_pct:.0f}% "
                    f"in {s.name}\n"
                    f"      Text: \"{text_preview}\"\n"
                    f"      Box: {s.width:.2f}\" x {s.height:.2f}\", "
                    f"font: {font_pt:.0f}pt, ~{total_lines} lines needed")

    def _check_zero_size(self, slide_num, shapes):
        """Check for shapes with zero width or height."""
        for s in shapes:
            if s.width == 0 and s.height == 0:
                self.warnings.append(
                    f"Slide {slide_num}: {s.name} has zero size "
                    f"at ({s.left:.2f}\", {s.top:.2f}\")")

    def _print_results(self):
        """Print validation results."""
        print()

        if self.errors:
            print(f"  ERRORS ({len(self.errors)}):")
            for err in self.errors:
                for line in err.split('\n'):
                    print(f"    [ERROR] {line}")
            print()

        if self.warnings:
            print(f"  WARNINGS ({len(self.warnings)}):")
            for warn in self.warnings:
                for line in warn.split('\n'):
                    print(f"    [WARN]  {line}")
            print()

        if not self.errors and not self.warnings:
            print("  [PASS] No layout issues found!")
        elif not self.errors:
            print(f"  [PASS] No errors. {len(self.warnings)} warning(s).")
        else:
            print(f"  [FAIL] {len(self.errors)} error(s), {len(self.warnings)} warning(s).")


# ============================================================================
# CLI
# ============================================================================

def find_pptx_files(directory="."):
    """Find all .pptx files in directory (not temp files)."""
    files = []
    for f in os.listdir(directory):
        if f.endswith(".pptx") and not f.startswith("~$"):
            files.append(os.path.join(directory, f))
    return sorted(files)


def main():
    if len(sys.argv) > 1:
        arg = sys.argv[1]

        if arg == "--all":
            files = find_pptx_files()
            if not files:
                print("No .pptx files found in current directory.")
                return
            print(f"Found {len(files)} PowerPoint file(s).")
            all_pass = True
            for f in files:
                validator = LayoutValidator(f)
                if not validator.validate():
                    all_pass = False
            print(f"\n{'=' * 60}")
            if all_pass:
                print("ALL FILES PASSED")
            else:
                print("SOME FILES HAD ERRORS")
            return

        if os.path.isfile(arg):
            validator = LayoutValidator(arg)
            validator.validate()
            return

        print(f"File not found: {arg}")
        return

    # Interactive mode
    files = find_pptx_files()
    if not files:
        print("No .pptx files found in current directory.")
        return

    print("PowerPoint Layout Validator")
    print("-" * 40)
    for i, f in enumerate(files, 1):
        print(f"  {i}. {os.path.basename(f)}")
    print(f"  {len(files) + 1}. Test ALL files")
    print()

    try:
        choice = input("Choose a file (number): ").strip()
        choice = int(choice)
    except (ValueError, EOFError):
        print("Invalid choice.")
        return

    if choice == len(files) + 1:
        for f in files:
            validator = LayoutValidator(f)
            validator.validate()
    elif 1 <= choice <= len(files):
        validator = LayoutValidator(files[choice - 1])
        validator.validate()
    else:
        print("Invalid choice.")


if __name__ == "__main__":
    main()
