"""
PowerPoint Generator for: Cybersecurity & Staying Safe Online
Kids Computer Science Class - March 21, 2026

Run:
  pip install python-pptx
  python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ============================================================================
# COLOR PALETTE
# ============================================================================
COLORS = {
    # Primary
    "dark_blue": RGBColor(30, 58, 95),
    "medium_blue": RGBColor(52, 103, 166),
    "light_blue": RGBColor(100, 149, 237),
    "sky_blue": RGBColor(173, 216, 230),

    # Accents
    "orange": RGBColor(255, 153, 51),
    "green": RGBColor(76, 175, 80),
    "red": RGBColor(244, 67, 54),
    "purple": RGBColor(156, 39, 176),
    "teal": RGBColor(0, 150, 136),

    # Theme-specific
    "shield_blue": RGBColor(25, 118, 210),
    "danger_red": RGBColor(211, 47, 47),
    "safe_green": RGBColor(56, 142, 60),
    "warning_yellow": RGBColor(255, 193, 7),
    "dark_red": RGBColor(183, 28, 28),

    # Neutrals
    "white": RGBColor(255, 255, 255),
    "light_gray": RGBColor(245, 245, 245),
    "dark_gray": RGBColor(66, 66, 66),
    "black": RGBColor(33, 33, 33),
}

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def set_shape_fill(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_shape_line(shape, color, width_pt=1):
    line = shape.line
    line.color.rgb = color
    line.width = Pt(width_pt)

def add_styled_textbox(slide, left, top, width, height, text,
                       font_size=18, font_color=None, bold=False,
                       alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    if font_color is None:
        font_color = COLORS["dark_gray"]

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment

    return textbox

def add_title_shape(slide, text, subtitle=None):
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(1.3)
    )
    set_shape_fill(title_bar, COLORS["dark_blue"])
    title_bar.line.fill.background()

    add_styled_textbox(
        slide, Inches(0.5), Inches(0.35), Inches(12.33), Inches(0.7),
        text, font_size=36, font_color=COLORS["white"], bold=True,
        alignment=PP_ALIGN.LEFT, font_name="Segoe UI Semibold"
    )

    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.3),
        Inches(13.33), Inches(0.08)
    )
    set_shape_fill(accent_line, COLORS["orange"])
    accent_line.line.fill.background()

    if subtitle:
        add_styled_textbox(
            slide, Inches(0.5), Inches(0.95), Inches(12.33), Inches(0.4),
            subtitle, font_size=16, font_color=COLORS["sky_blue"],
            alignment=PP_ALIGN.LEFT
        )

def add_rounded_box(slide, left, top, width, height, fill_color, text=None,
                    text_size=16, text_color=None, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    set_shape_fill(shape, fill_color)

    if border_color:
        set_shape_line(shape, border_color, 2)
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(text_size)
        p.font.color.rgb = text_color or COLORS["white"]
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(8)

    shape.text_frame.anchor = MSO_ANCHOR.MIDDLE
    return shape

def add_arrow(slide, left, top, width, height, fill_color, direction="right"):
    if direction == "right":
        shape_type = MSO_SHAPE.RIGHT_ARROW
    elif direction == "left":
        shape_type = MSO_SHAPE.LEFT_ARROW
    elif direction == "down":
        shape_type = MSO_SHAPE.DOWN_ARROW
    else:
        shape_type = MSO_SHAPE.UP_ARROW

    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    set_shape_fill(shape, fill_color)
    shape.line.fill.background()
    return shape


# ============================================================================
# SLIDES
# ============================================================================

def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    set_shape_fill(background, COLORS["dark_blue"])
    background.line.fill.background()

    # Shield shape
    add_rounded_box(
        slide, Inches(5.67), Inches(0.8), Inches(2.0), Inches(2.0),
        COLORS["shield_blue"], "🛡️", text_size=60, text_color=COLORS["white"]
    )

    add_rounded_box(
        slide, Inches(1.2), Inches(3.2), Inches(10.9), Inches(2.2),
        COLORS["shield_blue"]
    )

    add_styled_textbox(
        slide, Inches(0.6), Inches(3.4), Inches(12.13), Inches(1.0),
        "Cybersecurity & Staying Safe Online",
        font_size=42, font_color=COLORS["white"], bold=True,
        alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold"
    )

    add_styled_textbox(
        slide, Inches(0.6), Inches(4.4), Inches(12.13), Inches(0.6),
        "Protecting Yourself in the Digital World",
        font_size=22, font_color=COLORS["sky_blue"],
        alignment=PP_ALIGN.CENTER
    )

    add_styled_textbox(
        slide, Inches(0.6), Inches(6.6), Inches(12.13), Inches(0.4),
        "March 21, 2026",
        font_size=18, font_color=COLORS["sky_blue"],
        alignment=PP_ALIGN.CENTER
    )


def create_agenda_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Today's Journey")

    topics = [
        ("1", "What is Cybersecurity?", COLORS["shield_blue"]),
        ("2", "Passwords: Your First Line of Defense", COLORS["safe_green"]),
        ("3", "Phishing: Don't Take the Bait!", COLORS["danger_red"]),
        ("4", "Social Engineering: Tricks People Use", COLORS["orange"]),
        ("5", "Two-Factor Authentication (2FA)", COLORS["teal"]),
        ("6", "Safe Browsing & Digital Footprint", COLORS["purple"]),
        ("7", "Activity: Can You Spot the Scam?", COLORS["green"]),
    ]

    for i, (num, text, color) in enumerate(topics):
        y = Inches(1.7 + i * 0.72)

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1), y,
            Inches(0.55), Inches(0.55)
        )
        set_shape_fill(circle, color)
        circle.line.fill.background()

        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS["white"]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        add_styled_textbox(
            slide, Inches(2), y + Inches(0.1), Inches(10.5), Inches(0.5),
            text, font_size=20, font_color=COLORS["dark_gray"]
        )


def create_what_is_cybersecurity_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "What is Cybersecurity?", "Protecting computers, networks, and data from threats")

    add_rounded_box(
        slide, Inches(0.8), Inches(1.7), Inches(11.73), Inches(1.0),
        COLORS["shield_blue"],
        "Cybersecurity = Keeping your digital life safe from bad actors",
        text_size=22, text_color=COLORS["white"]
    )

    threats = [
        ("Hackers", "People who try to break into systems", COLORS["danger_red"]),
        ("Malware", "Harmful software (viruses, ransomware)", COLORS["dark_red"]),
        ("Scammers", "People who trick you into giving up info", COLORS["orange"]),
        ("Data Breaches", "When companies lose your personal data", COLORS["purple"]),
    ]

    for i, (title, desc, color) in enumerate(threats):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(3.1 + (i // 2) * 1.7)

        add_rounded_box(slide, x, y, Inches(5.8), Inches(1.4), color)
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.15), Inches(5.4), Inches(0.5),
            title, font_size=20, font_color=COLORS["white"],
            bold=True, alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.7), Inches(5.4), Inches(0.5),
            desc, font_size=15, font_color=COLORS["white"],
            alignment=PP_ALIGN.CENTER
        )

    add_rounded_box(
        slide, Inches(0.8), Inches(6.6), Inches(11.73), Inches(0.65),
        COLORS["light_gray"],
        "Cybersecurity is everyone's job -- not just IT professionals!",
        text_size=16, text_color=COLORS["dark_blue"], border_color=COLORS["shield_blue"]
    )


def create_why_it_matters_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Why Does This Matter to YOU?")

    scenarios = [
        ("Your Social Media", "Someone could post as you, message your friends, or lock you out", COLORS["shield_blue"]),
        ("Your Gaming Accounts", "Hackers steal accounts with rare items and sell them", COLORS["safe_green"]),
        ("Your Personal Info", "Name, school, address -- can be used for identity theft", COLORS["orange"]),
        ("Your Family", "Scammers target teens to get to parents' credit cards", COLORS["purple"]),
    ]

    for i, (title, desc, color) in enumerate(scenarios):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(1.7 + (i // 2) * 2.2)

        add_rounded_box(slide, x, y, Inches(5.8), Inches(1.8), color)
        add_styled_textbox(
            slide, x + Inches(0.15), y + Inches(0.25), Inches(5.5), Inches(0.5),
            title, font_size=20, font_color=COLORS["white"],
            bold=True, alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x + Inches(0.15), y + Inches(0.9), Inches(5.5), Inches(0.6),
            desc, font_size=15, font_color=COLORS["white"],
            alignment=PP_ALIGN.CENTER
        )

    add_styled_textbox(
        slide, Inches(1), Inches(6.5), Inches(11.33), Inches(0.5),
        "This is not hypothetical -- it happens to teens every day!",
        font_size=20, font_color=COLORS["danger_red"], bold=True,
        alignment=PP_ALIGN.CENTER
    )


def create_passwords_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Passwords: Your First Line of Defense")

    # Bad password examples
    add_rounded_box(
        slide, Inches(0.5), Inches(1.7), Inches(5.9), Inches(4.0),
        COLORS["light_gray"], border_color=COLORS["danger_red"]
    )
    add_rounded_box(
        slide, Inches(0.5), Inches(1.7), Inches(5.9), Inches(0.65),
        COLORS["danger_red"], "BAD Passwords", text_size=22, text_color=COLORS["white"]
    )
    bad_passwords = [
        "password123",
        "123456",
        "your birthday (01012010)",
        "your pet's name (fluffy)",
        "qwerty",
        "Same password for everything!",
    ]
    for i, pw in enumerate(bad_passwords):
        add_styled_textbox(
            slide, Inches(0.8), Inches(2.6 + i * 0.45), Inches(5.3), Inches(0.4),
            f"  {pw}", font_size=15, font_color=COLORS["dark_red"]
        )

    # Good password examples
    add_rounded_box(
        slide, Inches(6.93), Inches(1.7), Inches(5.9), Inches(4.0),
        COLORS["light_gray"], border_color=COLORS["safe_green"]
    )
    add_rounded_box(
        slide, Inches(6.93), Inches(1.7), Inches(5.9), Inches(0.65),
        COLORS["safe_green"], "GOOD Passwords", text_size=22, text_color=COLORS["white"]
    )
    good_passwords = [
        "Long: 12+ characters",
        "Mix: uppercase, lowercase, numbers",
        "Add symbols: !@#$%",
        "Use a passphrase: PurpleTiger$Eats42Tacos!",
        "Different for each account",
        "Use a password manager!",
    ]
    for i, pw in enumerate(good_passwords):
        add_styled_textbox(
            slide, Inches(7.23), Inches(2.6 + i * 0.45), Inches(5.3), Inches(0.4),
            f"  {pw}", font_size=15, font_color=COLORS["safe_green"]
        )

    add_rounded_box(
        slide, Inches(0.5), Inches(6.0), Inches(12.33), Inches(1.0),
        COLORS["shield_blue"],
        "A strong password is like a strong lock on your front door!",
        text_size=20, text_color=COLORS["white"]
    )


def create_password_strength_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "How Fast Can Passwords Be Cracked?")

    # Table-like rows
    col_starts = [0.5, 5.0, 9.0]
    col_widths = [4.3, 3.8, 3.8]

    headers = ["Password Type", "Example", "Time to Crack"]
    for i, (text, w, x) in enumerate(zip(headers, col_widths, col_starts)):
        add_rounded_box(
            slide, Inches(x), Inches(1.7), Inches(w), Inches(0.55),
            COLORS["dark_blue"], text, text_size=14, text_color=COLORS["white"]
        )

    rows = [
        ("6 numbers only", "123456", "Instantly", COLORS["danger_red"]),
        ("6 lowercase letters", "fluffy", "10 seconds", COLORS["danger_red"]),
        ("8 mixed characters", "Fluffy12", "8 hours", COLORS["orange"]),
        ("12 mixed + symbols", "Fl!ffy_12cats", "200 years", COLORS["safe_green"]),
        ("Passphrase (20+ chars)", "MyDog$Ate42Pizzas!", "Millions of years", COLORS["safe_green"]),
    ]

    row_colors = [COLORS["light_gray"], COLORS["sky_blue"]]
    for i, (ptype, example, time, indicator) in enumerate(rows):
        y = 2.4 + i * 0.65
        bg = row_colors[i % 2]
        for j, (text, w, x) in enumerate(zip([ptype, example, time], col_widths, col_starts)):
            color = indicator if j == 2 else bg
            text_c = COLORS["white"] if j == 2 else COLORS["dark_gray"]
            add_rounded_box(
                slide, Inches(x), Inches(y), Inches(w), Inches(0.55),
                color, text, text_size=13, text_color=text_c
            )

    add_rounded_box(
        slide, Inches(0.8), Inches(5.9), Inches(11.73), Inches(1.2),
        COLORS["shield_blue"]
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(6.0), Inches(11.33), Inches(0.5),
        "The longer and more complex, the safer you are!",
        font_size=22, font_color=COLORS["white"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(6.5), Inches(11.33), Inches(0.5),
        "Pro tip: A passphrase you can remember beats a short random password",
        font_size=16, font_color=COLORS["sky_blue"],
        alignment=PP_ALIGN.CENTER
    )


def create_phishing_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Phishing: Don't Take the Bait!", "When attackers pretend to be someone you trust")

    add_rounded_box(
        slide, Inches(0.8), Inches(1.7), Inches(11.73), Inches(1.0),
        COLORS["danger_red"],
        "Phishing = Fake emails, texts, or websites designed to steal your information",
        text_size=20, text_color=COLORS["white"]
    )

    examples = [
        ("Fake Email", '"Your account has been compromised! Click here immediately to verify."', COLORS["danger_red"]),
        ("Fake Text", '"You won a $500 gift card! Click this link to claim it now!"', COLORS["dark_red"]),
        ("Fake Website", 'g00gle.com instead of google.com -- looks almost the same!', COLORS["orange"]),
    ]

    for i, (title, desc, color) in enumerate(examples):
        y = Inches(3.1 + i * 1.2)

        add_rounded_box(slide, Inches(0.8), y, Inches(11.73), Inches(1.0), color)
        add_styled_textbox(
            slide, Inches(1.0), y + Inches(0.1), Inches(2.5), Inches(0.4),
            title, font_size=18, font_color=COLORS["white"],
            bold=True
        )
        add_styled_textbox(
            slide, Inches(3.5), y + Inches(0.1), Inches(8.8), Inches(0.7),
            desc, font_size=14, font_color=COLORS["white"]
        )

    add_rounded_box(
        slide, Inches(0.8), Inches(6.6), Inches(11.73), Inches(0.65),
        COLORS["safe_green"],
        "Rule #1: If it seems too urgent or too good to be true, it probably is!",
        text_size=16, text_color=COLORS["white"]
    )


def create_spot_phishing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "How to Spot a Phishing Attack")

    red_flags = [
        ("Urgency / Fear", '"Act NOW or your account will be deleted!"', COLORS["danger_red"]),
        ("Spelling Mistakes", '"Congradulations! You have ben selected..."', COLORS["orange"]),
        ("Suspicious Links", "Hover to check: amaz0n-secure-login.sketchy-site.com", COLORS["dark_red"]),
        ("Asks for Passwords", "No real company will ever ask for your password by email", COLORS["purple"]),
        ("Generic Greeting", '"Dear Customer" instead of your actual name', COLORS["shield_blue"]),
        ("Too Good to Be True", '"You won an iPhone!" -- Did you enter a contest?', COLORS["teal"]),
    ]

    for i, (flag, example, color) in enumerate(red_flags):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(1.7 + (i // 2) * 1.7)

        add_rounded_box(slide, x, y, Inches(5.8), Inches(1.4), color)
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.15), Inches(5.4), Inches(0.4),
            f"🚩 {flag}", font_size=17, font_color=COLORS["white"],
            bold=True
        )
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.65), Inches(5.4), Inches(0.5),
            example, font_size=13, font_color=COLORS["white"]
        )

    add_rounded_box(
        slide, Inches(0.8), Inches(6.8), Inches(11.73), Inches(0.5),
        COLORS["safe_green"],
        "When in doubt, go directly to the website by typing the URL yourself!",
        text_size=16, text_color=COLORS["white"]
    )


def create_social_engineering_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Social Engineering: Tricks People Use", "Hacking the human, not the computer")

    add_rounded_box(
        slide, Inches(0.8), Inches(1.7), Inches(11.73), Inches(0.9),
        COLORS["orange"],
        "Social Engineering = Manipulating people into giving up confidential information",
        text_size=18, text_color=COLORS["white"]
    )

    tricks = [
        ("Pretexting", "Pretending to be someone\nthey're not (IT support,\na teacher, a friend)", COLORS["danger_red"]),
        ("Baiting", 'Leaving a USB drive labeled\n"Employee Salaries" where\nsomeone will find it', COLORS["dark_red"]),
        ("Tailgating", "Following someone through\na locked door without\nscanning their own badge", COLORS["orange"]),
        ("Quid Pro Quo", '"I\'ll fix your computer\nif you give me your\nlogin credentials"', COLORS["purple"]),
    ]

    for i, (title, desc, color) in enumerate(tricks):
        x = Inches(0.8 + i * 3.1)
        y = Inches(3.0)

        add_rounded_box(slide, x, y, Inches(2.8), Inches(3.0), color)
        add_styled_textbox(
            slide, x + Inches(0.1), y + Inches(0.2), Inches(2.6), Inches(0.5),
            title, font_size=18, font_color=COLORS["white"],
            bold=True, alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x + Inches(0.1), y + Inches(0.9), Inches(2.6), Inches(1.8),
            desc, font_size=13, font_color=COLORS["white"],
            alignment=PP_ALIGN.CENTER
        )

    add_rounded_box(
        slide, Inches(0.8), Inches(6.4), Inches(11.73), Inches(0.8),
        COLORS["shield_blue"],
        "The best defense: THINK before you share any information!",
        text_size=20, text_color=COLORS["white"]
    )


def create_2fa_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Two-Factor Authentication (2FA)", "Adding a second lock to your door")

    # Without 2FA
    add_rounded_box(
        slide, Inches(0.5), Inches(1.7), Inches(5.9), Inches(2.5),
        COLORS["light_gray"], border_color=COLORS["danger_red"]
    )
    add_rounded_box(
        slide, Inches(0.5), Inches(1.7), Inches(5.9), Inches(0.65),
        COLORS["danger_red"], "Without 2FA", text_size=22, text_color=COLORS["white"]
    )
    add_styled_textbox(
        slide, Inches(0.8), Inches(2.6), Inches(5.3), Inches(0.4),
        "Password only", font_size=18, font_color=COLORS["dark_gray"], bold=True
    )
    add_styled_textbox(
        slide, Inches(0.8), Inches(3.1), Inches(5.3), Inches(0.4),
        "If someone gets your password...", font_size=15, font_color=COLORS["dark_gray"]
    )
    add_styled_textbox(
        slide, Inches(0.8), Inches(3.5), Inches(5.3), Inches(0.4),
        "They have FULL access to your account!", font_size=16, font_color=COLORS["danger_red"], bold=True
    )

    # With 2FA
    add_rounded_box(
        slide, Inches(6.93), Inches(1.7), Inches(5.9), Inches(2.5),
        COLORS["light_gray"], border_color=COLORS["safe_green"]
    )
    add_rounded_box(
        slide, Inches(6.93), Inches(1.7), Inches(5.9), Inches(0.65),
        COLORS["safe_green"], "With 2FA", text_size=22, text_color=COLORS["white"]
    )
    add_styled_textbox(
        slide, Inches(7.23), Inches(2.6), Inches(5.3), Inches(0.4),
        "Password + verification code", font_size=18, font_color=COLORS["dark_gray"], bold=True
    )
    add_styled_textbox(
        slide, Inches(7.23), Inches(3.1), Inches(5.3), Inches(0.4),
        "If someone gets your password...", font_size=15, font_color=COLORS["dark_gray"]
    )
    add_styled_textbox(
        slide, Inches(7.23), Inches(3.5), Inches(5.3), Inches(0.4),
        "They STILL can't get in without your phone!", font_size=16, font_color=COLORS["safe_green"], bold=True
    )

    # Types of 2FA
    add_styled_textbox(
        slide, Inches(0.8), Inches(4.5), Inches(11.73), Inches(0.5),
        "Common Types of 2FA:", font_size=20, font_color=COLORS["dark_blue"], bold=True
    )

    types_2fa = [
        ("Text Message (SMS)", "A code sent to your phone", COLORS["shield_blue"]),
        ("Authenticator App", "Google/Microsoft Authenticator", COLORS["teal"]),
        ("Biometrics", "Fingerprint or Face ID", COLORS["purple"]),
    ]

    for i, (title, desc, color) in enumerate(types_2fa):
        x = Inches(0.5 + i * 4.2)
        add_rounded_box(slide, x, Inches(5.1), Inches(3.8), Inches(1.2), color)
        add_styled_textbox(
            slide, x + Inches(0.1), Inches(5.2), Inches(3.6), Inches(0.5),
            title, font_size=16, font_color=COLORS["white"],
            bold=True, alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x + Inches(0.1), Inches(5.7), Inches(3.6), Inches(0.4),
            desc, font_size=13, font_color=COLORS["white"],
            alignment=PP_ALIGN.CENTER
        )

    add_rounded_box(
        slide, Inches(0.8), Inches(6.6), Inches(11.73), Inches(0.65),
        COLORS["safe_green"],
        "Turn on 2FA for your email, social media, and gaming accounts TODAY!",
        text_size=16, text_color=COLORS["white"]
    )


def create_safe_browsing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Safe Browsing & Your Digital Footprint")

    # Safe browsing tips
    add_rounded_box(
        slide, Inches(0.5), Inches(1.7), Inches(6.0), Inches(4.5),
        COLORS["light_gray"], border_color=COLORS["shield_blue"]
    )
    add_rounded_box(
        slide, Inches(0.5), Inches(1.7), Inches(6.0), Inches(0.65),
        COLORS["shield_blue"], "Safe Browsing Tips", text_size=20, text_color=COLORS["white"]
    )

    tips = [
        "Look for HTTPS (the S = Secure)",
        "Check the padlock icon in the address bar",
        "Don't download from unknown sites",
        "Keep your browser and OS updated",
        "Use private/incognito for sensitive searches",
        "Don't click pop-ups that say 'You have a virus!'",
    ]
    for i, tip in enumerate(tips):
        add_styled_textbox(
            slide, Inches(0.8), Inches(2.6 + i * 0.45), Inches(5.4), Inches(0.4),
            f"  {tip}", font_size=14, font_color=COLORS["dark_gray"]
        )

    # Digital footprint
    add_rounded_box(
        slide, Inches(6.83), Inches(1.7), Inches(6.0), Inches(4.5),
        COLORS["light_gray"], border_color=COLORS["purple"]
    )
    add_rounded_box(
        slide, Inches(6.83), Inches(1.7), Inches(6.0), Inches(0.65),
        COLORS["purple"], "Your Digital Footprint", text_size=20, text_color=COLORS["white"]
    )

    footprint = [
        "Everything you post online stays forever",
        "Colleges & employers check social media",
        "Photos contain hidden location data",
        "Think before you post: Would you show this",
        "  to your grandma? Your future boss?",
        "Your online reputation IS your reputation",
    ]
    for i, tip in enumerate(footprint):
        add_styled_textbox(
            slide, Inches(7.13), Inches(2.6 + i * 0.45), Inches(5.4), Inches(0.4),
            f"  {tip}", font_size=14, font_color=COLORS["dark_gray"]
        )

    add_rounded_box(
        slide, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.75),
        COLORS["shield_blue"],
        "The internet never forgets. Be intentional about what you share!",
        text_size=18, text_color=COLORS["white"]
    )


def create_real_world_breaches_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Real-World Data Breaches", "These actually happened!")

    breaches = [
        ("Yahoo (2013)", "3 BILLION accounts exposed", "Passwords, emails, security questions", COLORS["danger_red"]),
        ("Fortnite (2019)", "Millions of player accounts", "Hackers could buy V-Bucks with your card", COLORS["orange"]),
        ("T-Mobile (2023)", "37 million customers", "Names, addresses, phone numbers leaked", COLORS["dark_red"]),
        ("23andMe (2023)", "6.9 million users", "DNA data and family info exposed", COLORS["purple"]),
    ]

    for i, (company, scale, detail, color) in enumerate(breaches):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(1.7 + (i // 2) * 2.2)

        add_rounded_box(slide, x, y, Inches(5.8), Inches(1.8), color)
        add_styled_textbox(
            slide, x + Inches(0.15), y + Inches(0.15), Inches(5.5), Inches(0.5),
            company, font_size=20, font_color=COLORS["white"],
            bold=True, alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x + Inches(0.15), y + Inches(0.65), Inches(5.5), Inches(0.4),
            scale, font_size=16, font_color=COLORS["white"],
            bold=True, alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x + Inches(0.15), y + Inches(1.1), Inches(5.5), Inches(0.4),
            detail, font_size=14, font_color=COLORS["white"],
            alignment=PP_ALIGN.CENTER
        )

    add_rounded_box(
        slide, Inches(0.8), Inches(6.3), Inches(11.73), Inches(0.9),
        COLORS["shield_blue"],
        "This is why unique passwords and 2FA matter -- breaches WILL happen!",
        text_size=18, text_color=COLORS["white"]
    )


def create_ai_and_cybersecurity_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "AI & Cybersecurity", "AI is changing both the attacks AND the defenses")

    # Left side: AI threats
    add_rounded_box(
        slide, Inches(0.5), Inches(1.7), Inches(5.9), Inches(4.5),
        COLORS["light_gray"], border_color=COLORS["danger_red"]
    )
    add_rounded_box(
        slide, Inches(0.5), Inches(1.7), Inches(5.9), Inches(0.65),
        COLORS["danger_red"], "AI-Powered Attacks", text_size=20, text_color=COLORS["white"]
    )
    threats = [
        "Deepfakes: AI-generated fake videos",
        "  of real people (celebrities, politicians)",
        "AI Phishing: ChatGPT can write perfect",
        "  scam emails with no spelling mistakes",
        "Voice Cloning: AI can copy someone's",
        "  voice from a short audio clip",
        "AI Password Cracking: AI makes",
        "  guessing passwords even faster",
    ]
    for i, pt in enumerate(threats):
        add_styled_textbox(
            slide, Inches(0.8), Inches(2.6 + i * 0.43), Inches(5.3), Inches(0.4),
            f"  {pt}", font_size=14, font_color=COLORS["dark_red"]
        )

    # Right side: AI defenses
    add_rounded_box(
        slide, Inches(6.93), Inches(1.7), Inches(5.9), Inches(4.5),
        COLORS["light_gray"], border_color=COLORS["safe_green"]
    )
    add_rounded_box(
        slide, Inches(6.93), Inches(1.7), Inches(5.9), Inches(0.65),
        COLORS["safe_green"], "AI-Powered Defenses", text_size=20, text_color=COLORS["white"]
    )
    defenses = [
        "Spam Filters: AI detects phishing",
        "  emails before they reach your inbox",
        "Threat Detection: AI monitors networks",
        "  for suspicious activity 24/7",
        "Fraud Alerts: Banks use AI to flag",
        "  unusual transactions on your card",
        "CAPTCHA: AI tells bots apart from",
        "  humans to protect websites",
    ]
    for i, pt in enumerate(defenses):
        add_styled_textbox(
            slide, Inches(7.23), Inches(2.6 + i * 0.43), Inches(5.3), Inches(0.4),
            f"  {pt}", font_size=14, font_color=COLORS["safe_green"]
        )

    add_rounded_box(
        slide, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.75),
        COLORS["shield_blue"],
        "AI is a tool -- it can be used for good OR bad. Knowing cybersecurity helps you stay ahead!",
        text_size=18, text_color=COLORS["white"]
    )


def create_ai_deepfakes_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Deepfakes & AI Scams", "When you can't trust your own eyes and ears")

    add_rounded_box(
        slide, Inches(0.8), Inches(1.7), Inches(11.73), Inches(0.9),
        COLORS["danger_red"],
        "Deepfake = AI-generated video or audio that looks and sounds like a real person",
        text_size=18, text_color=COLORS["white"]
    )

    examples = [
        ("Fake Celebrity Videos", "AI creates videos of famous people saying things they never said -- used to spread misinformation", COLORS["dark_red"]),
        ("Voice Clone Scams", "Scammers clone a family member's voice from social media and call saying 'I need money!'", COLORS["orange"]),
        ("AI-Written Phishing", "AI writes phishing emails with perfect grammar and personalized details -- harder to spot!", COLORS["purple"]),
    ]

    for i, (title, desc, color) in enumerate(examples):
        y = Inches(3.0 + i * 1.2)
        add_rounded_box(slide, Inches(0.8), y, Inches(11.73), Inches(1.0), color)
        add_styled_textbox(
            slide, Inches(1.0), y + Inches(0.1), Inches(2.8), Inches(0.4),
            title, font_size=18, font_color=COLORS["white"], bold=True
        )
        add_styled_textbox(
            slide, Inches(4.0), y + Inches(0.1), Inches(8.3), Inches(0.7),
            desc, font_size=14, font_color=COLORS["white"]
        )

    add_rounded_box(
        slide, Inches(0.8), Inches(6.6), Inches(11.73), Inches(0.65),
        COLORS["safe_green"],
        "Defense: Verify through a different channel! Call the person directly. Check multiple sources.",
        text_size=16, text_color=COLORS["white"]
    )


def create_protect_yourself_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "How to Protect Yourself: Summary")

    tips = [
        ("1", "Use strong, unique passwords for every account", COLORS["safe_green"]),
        ("2", "Turn on Two-Factor Authentication (2FA) everywhere", COLORS["shield_blue"]),
        ("3", "Never click suspicious links -- type URLs yourself", COLORS["danger_red"]),
        ("4", "Think before you share personal info online", COLORS["purple"]),
        ("5", "Keep your software and apps updated", COLORS["teal"]),
        ("6", "When in doubt, ask a trusted adult!", COLORS["orange"]),
    ]

    for i, (num, text, color) in enumerate(tips):
        y = Inches(1.6 + i * 0.85)

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8), y + Inches(0.05),
            Inches(0.55), Inches(0.55)
        )
        set_shape_fill(circle, color)
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS["white"]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        add_styled_textbox(
            slide, Inches(1.7), y + Inches(0.1), Inches(10.5), Inches(0.5),
            text, font_size=20, font_color=COLORS["dark_gray"]
        )

    add_styled_textbox(
        slide, Inches(1), Inches(6.8), Inches(11.33), Inches(0.5),
        "You are now a Cybersecurity-Aware Digital Citizen!",
        font_size=22, font_color=COLORS["shield_blue"], bold=True,
        alignment=PP_ALIGN.CENTER
    )


def create_activity_preview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Activity: Can You Spot the Scam?")

    add_rounded_box(
        slide, Inches(1), Inches(1.9), Inches(11.33), Inches(1.2),
        COLORS["orange"],
        "Time to put your skills to the test!",
        text_size=24, text_color=COLORS["white"]
    )

    steps = [
        ("1", "Look at sample emails, texts, and websites", COLORS["shield_blue"]),
        ("2", "Decide: Is it REAL or a SCAM?", COLORS["danger_red"]),
        ("3", "Identify the red flags that gave it away", COLORS["safe_green"]),
        ("4", "Create your own strong password using a passphrase", COLORS["purple"]),
    ]

    for i, (num, desc, color) in enumerate(steps):
        y = Inches(3.5 + i * 0.8)

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.2), y + Inches(0.05),
            Inches(0.5), Inches(0.5)
        )
        set_shape_fill(circle, color)
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS["white"]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        add_styled_textbox(
            slide, Inches(2.0), y + Inches(0.1), Inches(10.0), Inches(0.5),
            desc, font_size=18, font_color=COLORS["dark_gray"]
        )


def create_recap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Today's Key Takeaways")

    takeaways = [
        ("1", "Cybersecurity protects your data, accounts, and identity"),
        ("2", "Strong passwords: long, unique, use a passphrase"),
        ("3", "Phishing tricks you with urgency and fake trust"),
        ("4", "Social engineering hacks people, not computers"),
        ("5", "2FA adds a second layer of protection"),
        ("6", "Your digital footprint follows you forever -- post wisely!"),
    ]

    for i, (num, text) in enumerate(takeaways):
        y = Inches(1.6 + i * 0.82)

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8), y + Inches(0.05),
            Inches(0.5), Inches(0.5)
        )
        set_shape_fill(circle, COLORS["shield_blue"])
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS["white"]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        add_styled_textbox(
            slide, Inches(1.6), y + Inches(0.1), Inches(10.5), Inches(0.5),
            text, font_size=18, font_color=COLORS["dark_gray"]
        )

    add_styled_textbox(
        slide, Inches(1), Inches(6.8), Inches(11.33), Inches(0.5),
        "Great job today! Time for Kahoot!",
        font_size=24, font_color=COLORS["dark_blue"], bold=True,
        alignment=PP_ALIGN.CENTER
    )


def create_questions_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    set_shape_fill(background, COLORS["dark_blue"])
    background.line.fill.background()

    add_styled_textbox(
        slide, Inches(0), Inches(1.6), Inches(13.33), Inches(2),
        "?", font_size=120, font_color=COLORS["white"],
        alignment=PP_ALIGN.CENTER
    )

    add_styled_textbox(
        slide, Inches(0), Inches(3.5), Inches(13.33), Inches(1),
        "Questions?", font_size=60, font_color=COLORS["white"], bold=True,
        alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold"
    )

    add_styled_textbox(
        slide, Inches(0), Inches(5.1), Inches(13.33), Inches(0.6),
        "Ask away before we start the Kahoot!",
        font_size=24, font_color=COLORS["sky_blue"],
        alignment=PP_ALIGN.CENTER
    )


# ============================================================================
# MAIN
# ============================================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs)                  # 1
    create_agenda_slide(prs)                 # 2
    create_what_is_cybersecurity_slide(prs)  # 3
    create_why_it_matters_slide(prs)         # 4
    create_passwords_intro_slide(prs)        # 5
    create_password_strength_slide(prs)      # 6
    create_phishing_intro_slide(prs)         # 7
    create_spot_phishing_slide(prs)          # 8
    create_social_engineering_slide(prs)     # 9
    create_2fa_slide(prs)                    # 10
    create_safe_browsing_slide(prs)          # 11
    create_real_world_breaches_slide(prs)    # 12
    create_ai_and_cybersecurity_slide(prs)   # 13
    create_ai_deepfakes_slide(prs)           # 14
    create_protect_yourself_slide(prs)       # 15
    create_activity_preview_slide(prs)       # 16
    create_recap_slide(prs)                  # 17
    create_questions_slide(prs)              # 18

    script_dir = os.path.dirname(os.path.abspath(__file__))
    parent_dir = os.path.dirname(script_dir)
    output_path = os.path.join(parent_dir, "2026-03-21_Cybersecurity_Staying_Safe_Online_BEAUTIFUL.pptx")

    prs.save(output_path)
    print(f"[SUCCESS] Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
