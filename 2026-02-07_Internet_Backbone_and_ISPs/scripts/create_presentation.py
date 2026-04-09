"""
PowerPoint Generator for: How the Internet Connects Everything
Kids Computer Science Class - February 7, 2026

Run:
  pip install python-pptx
  python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ============================================================================
# COLOR PALETTE
# ============================================================================
COLORS = {
    # Primary
    "dark_blue": RGBColor(30, 58, 95),
    "medium_blue": RGBColor(52, 103, 166),
    "light_blue": RGBColor(100, 149, 237),
    "sky_blue": RGBColor(173, 216, 230),

    # Accents
    "orange": RGBColor(255, 153, 51),
    "green": RGBColor(76, 175, 80),
    "red": RGBColor(244, 67, 54),
    "purple": RGBColor(156, 39, 176),
    "teal": RGBColor(0, 150, 136),

    # Neutrals
    "white": RGBColor(255, 255, 255),
    "light_gray": RGBColor(245, 245, 245),
    "dark_gray": RGBColor(66, 66, 66),
    "black": RGBColor(33, 33, 33),

    # Theme
    "cable_blue": RGBColor(0, 120, 212),
    "ocean_blue": RGBColor(0, 102, 153),
    "cdn_pink": RGBColor(233, 30, 99),
}

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def set_shape_fill(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_shape_line(shape, color, width_pt=1):
    line = shape.line
    line.color.rgb = color
    line.width = Pt(width_pt)

def add_styled_textbox(slide, left, top, width, height, text,
                       font_size=18, font_color=None, bold=False,
                       alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    if font_color is None:
        font_color = COLORS["dark_gray"]

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment

    return textbox

def add_title_shape(slide, text, subtitle=None):
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(1.3)
    )
    set_shape_fill(title_bar, COLORS["dark_blue"])
    title_bar.line.fill.background()

    add_styled_textbox(
        slide, Inches(0.5), Inches(0.35), Inches(12.33), Inches(0.7),
        text, font_size=36, font_color=COLORS["white"], bold=True,
        alignment=PP_ALIGN.LEFT, font_name="Segoe UI Semibold"
    )

    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.3),
        Inches(13.33), Inches(0.08)
    )
    set_shape_fill(accent_line, COLORS["orange"])
    accent_line.line.fill.background()

    if subtitle:
        add_styled_textbox(
            slide, Inches(0.5), Inches(0.95), Inches(12.33), Inches(0.4),
            subtitle, font_size=16, font_color=COLORS["sky_blue"],
            alignment=PP_ALIGN.LEFT
        )

def add_rounded_box(slide, left, top, width, height, fill_color, text=None,
                    text_size=16, text_color=None, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    set_shape_fill(shape, fill_color)

    if border_color:
        set_shape_line(shape, border_color, 2)
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(text_size)
        p.font.color.rgb = text_color or COLORS["white"]
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(8)

    shape.text_frame.anchor = MSO_ANCHOR.MIDDLE
    return shape

def add_arrow(slide, left, top, width, height, fill_color, direction="right"):
    if direction == "right":
        shape_type = MSO_SHAPE.RIGHT_ARROW
    elif direction == "left":
        shape_type = MSO_SHAPE.LEFT_ARROW
    elif direction == "down":
        shape_type = MSO_SHAPE.DOWN_ARROW
    else:
        shape_type = MSO_SHAPE.UP_ARROW

    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    set_shape_fill(shape, fill_color)
    shape.line.fill.background()
    return shape

# ============================================================================
# SLIDES
# ============================================================================

def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    set_shape_fill(background, COLORS["dark_blue"])
    background.line.fill.background()

    add_rounded_box(
        slide, Inches(1.2), Inches(1.7), Inches(10.9), Inches(3.1),
        COLORS["light_blue"]
    )

    add_styled_textbox(
        slide, Inches(0.6), Inches(2.2), Inches(12.13), Inches(1.2),
        "How the Internet Connects Everything",
        font_size=44, font_color=COLORS["dark_blue"], bold=True,
        alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold"
    )

    add_styled_textbox(
        slide, Inches(0.6), Inches(3.5), Inches(12.13), Inches(0.6),
        "ISPs, Backbones, and the Cables Under the Ocean",
        font_size=22, font_color=COLORS["dark_blue"],
        alignment=PP_ALIGN.CENTER
    )

    add_styled_textbox(
        slide, Inches(0.6), Inches(6.6), Inches(12.13), Inches(0.4),
        "February 7, 2026",
        font_size=18, font_color=COLORS["sky_blue"],
        alignment=PP_ALIGN.CENTER
    )


def create_agenda_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Today's Journey")

    topics = [
        ("1", "Quick Review: Cloud & Data Centers", COLORS["light_blue"]),
        ("2", "What is an ISP?", COLORS["green"]),
        ("3", "Internet Backbones & Fiber", COLORS["cable_blue"]),
        ("4", "Undersea Cables & What If They Break?", COLORS["ocean_blue"]),
        ("5", "Internet Exchange Points (IXPs)", COLORS["orange"]),
        ("6", "CDNs: Content Close to You", COLORS["cdn_pink"]),
        ("7", "Asking ChatGPT: A Real-World Trace", COLORS["purple"]),
        ("8", "Hands-on: Internet Map Explorer", COLORS["teal"]),
    ]

    for i, (num, text, color) in enumerate(topics):
        y = Inches(1.6 + i * 0.68)

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1), y,
            Inches(0.55), Inches(0.55)
        )
        set_shape_fill(circle, color)
        circle.line.fill.background()

        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS["white"]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        add_styled_textbox(
            slide, Inches(2), y + Inches(0.1), Inches(10.5), Inches(0.5),
            text, font_size=20, font_color=COLORS["dark_gray"]
        )


def create_review_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Quick Review", "From last week")

    add_rounded_box(
        slide, Inches(0.8), Inches(1.8), Inches(11.73), Inches(1.1),
        COLORS["light_blue"],
        "The cloud is real buildings filled with servers (data centers)",
        text_size=24, text_color=COLORS["dark_blue"]
    )

    add_styled_textbox(
        slide, Inches(1), Inches(3.5), Inches(11.33), Inches(0.6),
        "Today we answer a new question:",
        font_size=20, font_color=COLORS["dark_gray"], bold=True,
        alignment=PP_ALIGN.CENTER
    )
    add_rounded_box(
        slide, Inches(1.5), Inches(4.2), Inches(10.33), Inches(1.3),
        COLORS["orange"],
        "How do those data centers connect to the whole world?",
        text_size=24, text_color=COLORS["white"]
    )

    add_styled_textbox(
        slide, Inches(1), Inches(6.2), Inches(11.33), Inches(0.6),
        "Answer: ISPs, backbones, and huge cables!",
        font_size=20, font_color=COLORS["medium_blue"], bold=True,
        alignment=PP_ALIGN.CENTER
    )


def create_isp_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "What is an ISP?")

    add_rounded_box(
        slide, Inches(0.8), Inches(1.7), Inches(11.73), Inches(1.0),
        COLORS["green"],
        "ISP = Internet Service Provider",
        text_size=28, text_color=COLORS["white"]
    )

    points = [
        "The company that connects your home or school to the internet",
        "Examples: Comcast, Spectrum, AT&T, Verizon, Cox",
        "They own local cables, towers, and neighborhood equipment",
        "Your ISP is your first hop to the wider internet",
    ]

    for i, item in enumerate(points):
        add_styled_textbox(
            slide, Inches(1), Inches(3 + i * 0.6), Inches(11.33), Inches(0.5),
            f"- {item}", font_size=18, font_color=COLORS["dark_gray"]
        )

    add_rounded_box(
        slide, Inches(1.2), Inches(6.4), Inches(10.93), Inches(0.75),
        COLORS["light_blue"],
        "Think of an ISP like your neighborhood road system",
        text_size=18, text_color=COLORS["dark_blue"]
    )


def create_home_to_isp_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "From Your Home to the Internet")

    # Home device
    add_rounded_box(
        slide, Inches(0.8), Inches(2.2), Inches(2.6), Inches(1.4),
        COLORS["light_blue"], "Your Device",
        text_size=16, text_color=COLORS["dark_blue"]
    )

    add_arrow(slide, Inches(3.6), Inches(2.7), Inches(1.0), Inches(0.5), COLORS["orange"], "right")

    add_rounded_box(
        slide, Inches(4.8), Inches(2.2), Inches(2.6), Inches(1.4),
        COLORS["green"], "Home Router",
        text_size=16, text_color=COLORS["white"]
    )

    add_arrow(slide, Inches(7.5), Inches(2.7), Inches(1.0), Inches(0.5), COLORS["orange"], "right")

    add_rounded_box(
        slide, Inches(8.7), Inches(2.2), Inches(3.6), Inches(1.4),
        COLORS["medium_blue"], "ISP Network",
        text_size=16, text_color=COLORS["white"]
    )

    add_styled_textbox(
        slide, Inches(1), Inches(4.2), Inches(11.33), Inches(0.6),
        "Your data leaves your house and enters your ISP's network.",
        font_size=18, font_color=COLORS["dark_gray"], alignment=PP_ALIGN.CENTER
    )

    add_styled_textbox(
        slide, Inches(1), Inches(5.1), Inches(11.33), Inches(0.6),
        "From there, it travels to bigger and bigger networks.",
        font_size=18, font_color=COLORS["dark_gray"], alignment=PP_ALIGN.CENTER
    )


def create_backbone_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Internet Backbones")

    add_rounded_box(
        slide, Inches(0.8), Inches(1.7), Inches(11.73), Inches(1.0),
        COLORS["cable_blue"],
        "Backbones = Super-fast fiber highways that carry internet traffic",
        text_size=20, text_color=COLORS["white"]
    )

    # Visual lanes
    for i in range(3):
        lane = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(3 + i * 0.8),
            Inches(11.33), Inches(0.45)
        )
        set_shape_fill(lane, COLORS["light_gray"])
        lane.line.fill.background()

        dash = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.5), Inches(3.18 + i * 0.8),
            Inches(1.2), Inches(0.1)
        )
        set_shape_fill(dash, COLORS["orange"])
        dash.line.fill.background()

    add_styled_textbox(
        slide, Inches(1), Inches(5.8), Inches(11.33), Inches(0.6),
        "Fiber optic cables use light to send data super fast.",
        font_size=18, font_color=COLORS["dark_gray"], alignment=PP_ALIGN.CENTER
    )


def create_undersea_cables_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Undersea Cables")

    add_rounded_box(
        slide, Inches(0.8), Inches(1.7), Inches(11.73), Inches(1.0),
        COLORS["ocean_blue"],
        "Most of the world's internet travels through cables under the ocean",
        text_size=18, text_color=COLORS["white"]
    )

    ocean = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(3.0),
        Inches(11.33), Inches(2.6)
    )
    set_shape_fill(ocean, COLORS["sky_blue"])
    ocean.line.fill.background()

    for i in range(4):
        cable = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.2), Inches(3.35 + i * 0.55),
            Inches(10.8), Inches(0.08)
        )
        set_shape_fill(cable, COLORS["cable_blue"])
        cable.line.fill.background()

    add_styled_textbox(
        slide, Inches(1), Inches(5.9), Inches(11.33), Inches(0.6),
        "These cables connect continents so data can cross the globe.",
        font_size=18, font_color=COLORS["dark_gray"], alignment=PP_ALIGN.CENTER
    )


def create_ixp_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Internet Exchange Points (IXPs)")

    add_rounded_box(
        slide, Inches(0.8), Inches(1.7), Inches(11.73), Inches(1.0),
        COLORS["orange"],
        "IXP = A meeting place where networks connect and swap traffic",
        text_size=18, text_color=COLORS["white"]
    )

    for i, label in enumerate(["ISP A", "ISP B", "Cloud Provider", "University"]):
        x = Inches(1 + (i % 2) * 6)
        y = Inches(3 + (i // 2) * 1.6)
        add_rounded_box(
            slide, x, y, Inches(5.2), Inches(1.1),
            COLORS["light_gray"], label,
            text_size=16, text_color=COLORS["dark_blue"], border_color=COLORS["orange"]
        )

    add_styled_textbox(
        slide, Inches(1), Inches(6.2), Inches(11.33), Inches(0.6),
        "IXPs make the internet faster and cheaper by shortening routes.",
        font_size=18, font_color=COLORS["dark_gray"], alignment=PP_ALIGN.CENTER
    )


def create_cdn_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "CDNs: Content Close to You")

    add_rounded_box(
        slide, Inches(0.8), Inches(1.7), Inches(11.73), Inches(0.85),
        COLORS["cdn_pink"],
        "CDN = Content Delivery Network",
        text_size=26, text_color=COLORS["white"]
    )

    # --- WITHOUT CDN (left side) ---
    add_rounded_box(
        slide, Inches(0.8), Inches(2.85), Inches(5.6), Inches(2.6),
        COLORS["light_gray"], border_color=COLORS["red"]
    )
    add_styled_textbox(
        slide, Inches(0.8), Inches(2.95), Inches(5.6), Inches(0.4),
        "Without CDN", font_size=18, font_color=COLORS["red"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    # You box
    add_rounded_box(
        slide, Inches(1.1), Inches(3.5), Inches(1.6), Inches(0.7),
        COLORS["light_blue"], "You",
        text_size=14, text_color=COLORS["dark_blue"]
    )
    # Long arrow
    add_arrow(slide, Inches(2.85), Inches(3.65), Inches(1.4), Inches(0.4), COLORS["red"], "right")
    # Far server
    add_rounded_box(
        slide, Inches(4.4), Inches(3.5), Inches(1.7), Inches(0.7),
        COLORS["red"], "Server",
        text_size=14, text_color=COLORS["white"]
    )
    add_styled_textbox(
        slide, Inches(0.8), Inches(4.4), Inches(5.6), Inches(0.4),
        "3,000+ miles away", font_size=14, font_color=COLORS["red"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(0.8), Inches(4.85), Inches(5.6), Inches(0.4),
        "Slow loading, buffering, lag",
        font_size=13, font_color=COLORS["dark_gray"], alignment=PP_ALIGN.CENTER
    )

    # --- WITH CDN (right side) ---
    add_rounded_box(
        slide, Inches(6.93), Inches(2.85), Inches(5.6), Inches(2.6),
        COLORS["light_gray"], border_color=COLORS["green"]
    )
    add_styled_textbox(
        slide, Inches(6.93), Inches(2.95), Inches(5.6), Inches(0.4),
        "With CDN", font_size=18, font_color=COLORS["green"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    # You box
    add_rounded_box(
        slide, Inches(7.23), Inches(3.5), Inches(1.6), Inches(0.7),
        COLORS["light_blue"], "You",
        text_size=14, text_color=COLORS["dark_blue"]
    )
    # Short arrow
    add_arrow(slide, Inches(8.97), Inches(3.65), Inches(0.8), Inches(0.4), COLORS["green"], "right")
    # Nearby cache
    add_rounded_box(
        slide, Inches(9.9), Inches(3.5), Inches(2.3), Inches(0.7),
        COLORS["green"], "Nearby Cache",
        text_size=14, text_color=COLORS["white"]
    )
    add_styled_textbox(
        slide, Inches(6.93), Inches(4.4), Inches(5.6), Inches(0.4),
        "~50 miles away", font_size=14, font_color=COLORS["green"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(6.93), Inches(4.85), Inches(5.6), Inches(0.4),
        "Fast loading, no buffering!",
        font_size=13, font_color=COLORS["dark_gray"], alignment=PP_ALIGN.CENTER
    )

    # Examples row at bottom
    add_rounded_box(
        slide, Inches(0.8), Inches(5.8), Inches(11.73), Inches(0.7),
        COLORS["light_gray"]
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(5.9), Inches(11.33), Inches(0.5),
        "CDN examples: YouTube, Netflix, TikTok, Spotify, Fortnite updates, school portals",
        font_size=14, font_color=COLORS["dark_gray"], alignment=PP_ALIGN.CENTER, bold=True
    )

    add_styled_textbox(
        slide, Inches(1), Inches(6.7), Inches(11.33), Inches(0.5),
        "Result: Faster loading, less buffering, and smoother games.",
        font_size=18, font_color=COLORS["dark_gray"], alignment=PP_ALIGN.CENTER
    )


def create_video_path_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Example: Loading a YouTube Video")

    steps = [
        "You tap a video",
        "Your request goes to your ISP",
        "It hops across the backbone",
        "A nearby CDN serves the video",
        "The video streams to your device"
    ]

    for i, step in enumerate(steps):
        y = Inches(1.9 + i * 1.0)
        add_rounded_box(
            slide, Inches(1), y, Inches(11.33), Inches(0.8),
            COLORS["light_gray"],
            step, text_size=18, text_color=COLORS["dark_blue"], border_color=COLORS["medium_blue"]
        )


def create_latency_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Latency: Distance Matters")

    add_styled_textbox(
        slide, Inches(1), Inches(2.0), Inches(11.33), Inches(0.7),
        "Latency = the delay before data arrives.",
        font_size=22, font_color=COLORS["dark_blue"], bold=True,
        alignment=PP_ALIGN.CENTER
    )

    add_rounded_box(
        slide, Inches(1.5), Inches(3.1), Inches(10.33), Inches(1.2),
        COLORS["light_blue"],
        "Closer servers = faster response",
        text_size=20, text_color=COLORS["dark_blue"]
    )

    add_styled_textbox(
        slide, Inches(1), Inches(4.8), Inches(11.33), Inches(0.6),
        "That is why CDNs and local data centers are so important.",
        font_size=18, font_color=COLORS["dark_gray"], alignment=PP_ALIGN.CENTER
    )


def create_slow_internet_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Why Can the Internet Feel Slow?")

    reasons = [
        ("Busy Wi-Fi", "Too many devices at once"),
        ("Long Distance", "Data travels farther"),
        ("Backbone Congestion", "Too much traffic"),
        ("Server Load", "Popular sites get crowded"),
    ]

    for i, (title, desc) in enumerate(reasons):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(1.9 + (i // 2) * 2.0)

        add_rounded_box(
            slide, x, y, Inches(5.8), Inches(1.5),
            COLORS["green"]
        )
        add_styled_textbox(
            slide, x, y + Inches(0.2), Inches(5.8), Inches(0.5),
            title, font_size=16, font_color=COLORS["white"], bold=True, alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x, y + Inches(0.8), Inches(5.8), Inches(0.5),
            desc, font_size=12, font_color=COLORS["white"], alignment=PP_ALIGN.CENTER
        )

    add_styled_textbox(
        slide, Inches(1), Inches(6.4), Inches(11.33), Inches(0.6),
        "Most slowdowns are about distance, traffic, or Wi-Fi.",
        font_size=18, font_color=COLORS["dark_gray"], alignment=PP_ALIGN.CENTER
    )


def create_activity_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Hands-on: Internet Map Explorer")

    add_rounded_box(
        slide, Inches(1), Inches(1.9), Inches(11.33), Inches(1.2),
        COLORS["purple"],
        "We will explore real internet maps together",
        text_size=22, text_color=COLORS["white"]
    )

    items = [
        "Find where undersea cables land near North America",
        "Look for big internet hubs (IXPs)",
        "See how many cables connect different continents",
        "Notice: the internet is physical, not magic"
    ]

    for i, item in enumerate(items):
        add_styled_textbox(
            slide, Inches(1), Inches(3.6 + i * 0.6), Inches(11.33), Inches(0.5),
            f"- {item}", font_size=18, font_color=COLORS["dark_gray"]
        )

    add_styled_textbox(
        slide, Inches(1), Inches(6.6), Inches(11.33), Inches(0.5),
        "Get ready for the class activity!",
        font_size=20, font_color=COLORS["medium_blue"], bold=True,
        alignment=PP_ALIGN.CENTER
    )

def create_ai_physical_internet_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "What Happens When You Ask ChatGPT a Question?")

    # Step-by-step trace through the physical internet
    steps = [
        ("1", "You type a question on your phone or laptop", COLORS["light_blue"], COLORS["dark_blue"]),
        ("2", "Your router sends it to your ISP", COLORS["green"], COLORS["white"]),
        ("3", "Your ISP passes it along the backbone", COLORS["cable_blue"], COLORS["white"]),
        ("4", "It may cross undersea cables to reach a data center", COLORS["ocean_blue"], COLORS["white"]),
        ("5", "Thousands of GPUs process your question", COLORS["purple"], COLORS["white"]),
        ("6", "The answer travels all the way back to you!", COLORS["orange"], COLORS["white"]),
    ]

    for i, (num, desc, bg_color, text_color) in enumerate(steps):
        y = Inches(1.65 + i * 0.78)

        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.9), y + Inches(0.05),
            Inches(0.55), Inches(0.55)
        )
        set_shape_fill(circle, bg_color)
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.color.rgb = text_color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Step description
        add_rounded_box(
            slide, Inches(1.7), y, Inches(10.63), Inches(0.65),
            bg_color, desc,
            text_size=16, text_color=text_color
        )

        # Arrow between steps (except after last)
        if i < len(steps) - 1:
            add_arrow(slide, Inches(6.5), y + Inches(0.6), Inches(0.4), Inches(0.2), COLORS["dark_gray"], "down")

    add_styled_textbox(
        slide, Inches(1), Inches(6.6), Inches(11.33), Inches(0.5),
        "Total time: ~1-2 seconds! All thanks to fast cables and backbones.",
        font_size=18, font_color=COLORS["dark_blue"], bold=True,
        alignment=PP_ALIGN.CENTER
    )


def create_cable_breaks_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "What If a Cable Breaks?")

    add_rounded_box(
        slide, Inches(0.8), Inches(1.7), Inches(11.73), Inches(0.85),
        COLORS["ocean_blue"],
        "It happens more often than you think!",
        text_size=22, text_color=COLORS["white"]
    )

    # Causes - 2x2 grid
    causes = [
        ("Shark Bites", "Sharks sometimes bite\nunderwater cables!", COLORS["red"]),
        ("Ship Anchors", "Anchors can snag and\ncut cables on the seafloor", COLORS["orange"]),
        ("Earthquakes", "Undersea quakes can\nsnap cables apart", COLORS["purple"]),
        ("Wear & Tear", "Salt water and pressure\ndamage cables over time", COLORS["dark_gray"]),
    ]

    for i, (title, desc, color) in enumerate(causes):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(2.85 + (i // 2) * 1.55)

        add_rounded_box(slide, x, y, Inches(5.8), Inches(1.3), color)
        add_styled_textbox(
            slide, x + Inches(0.15), y + Inches(0.12), Inches(5.5), Inches(0.4),
            title, font_size=16, font_color=COLORS["white"],
            bold=True, alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x + Inches(0.15), y + Inches(0.55), Inches(5.5), Inches(0.6),
            desc, font_size=12, font_color=COLORS["white"],
            alignment=PP_ALIGN.CENTER
        )

    # Reassurance box
    add_rounded_box(
        slide, Inches(0.8), Inches(6.15), Inches(11.73), Inches(0.9),
        COLORS["green"]
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(6.25), Inches(11.33), Inches(0.7),
        "Good news: There are 500+ undersea cables! If one breaks, traffic reroutes through others.",
        font_size=16, font_color=COLORS["white"], bold=True,
        alignment=PP_ALIGN.CENTER
    )


def create_recap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Today's Key Takeaways")

    takeaways = [
        ("ISP", "Your first connection to the internet"),
        ("BACKBONE", "Fast fiber highways connecting cities"),
        ("UNDERSEA CABLES", "Connect continents under the ocean"),
        ("IXP", "Meeting point where networks swap traffic"),
        ("CDN", "Brings content closer for speed"),
    ]

    for i, (term, definition) in enumerate(takeaways):
        y = Inches(1.7 + i * 0.9)
        add_rounded_box(
            slide, Inches(0.8), y, Inches(3.0), Inches(0.75),
            COLORS["light_blue"],
            term, text_size=16, text_color=COLORS["dark_blue"]
        )
        add_styled_textbox(
            slide, Inches(4.1), y + Inches(0.15), Inches(8.2), Inches(0.6),
            definition, font_size=18, font_color=COLORS["dark_gray"]
        )

    add_styled_textbox(
        slide, Inches(1), Inches(6.8), Inches(11.33), Inches(0.5),
        "Great job today! Time for Kahoot!",
        font_size=24, font_color=COLORS["dark_blue"], bold=True,
        alignment=PP_ALIGN.CENTER
    )


def create_questions_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    set_shape_fill(background, COLORS["dark_blue"])
    background.line.fill.background()

    add_styled_textbox(
        slide, Inches(0), Inches(1.6), Inches(13.33), Inches(2),
        "?", font_size=120, font_color=COLORS["white"],
        alignment=PP_ALIGN.CENTER
    )

    add_styled_textbox(
        slide, Inches(0), Inches(3.5), Inches(13.33), Inches(1),
        "Questions?", font_size=60, font_color=COLORS["white"], bold=True,
        alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold"
    )

    add_styled_textbox(
        slide, Inches(0), Inches(5.1), Inches(13.33), Inches(0.6),
        "Ask away before we start the Kahoot!",
        font_size=24, font_color=COLORS["sky_blue"],
        alignment=PP_ALIGN.CENTER
    )

# ============================================================================
# MAIN
# ============================================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs)
    create_agenda_slide(prs)
    create_review_slide(prs)
    create_isp_slide(prs)
    create_home_to_isp_slide(prs)
    create_backbone_slide(prs)
    create_undersea_cables_slide(prs)
    create_cable_breaks_slide(prs)
    create_ixp_slide(prs)
    create_cdn_slide(prs)
    create_video_path_slide(prs)
    create_latency_slide(prs)
    create_slow_internet_slide(prs)
    create_ai_physical_internet_slide(prs)
    create_activity_slide(prs)
    create_recap_slide(prs)
    create_questions_slide(prs)

    script_dir = os.path.dirname(os.path.abspath(__file__))
    parent_dir = os.path.dirname(script_dir)
    output_path = os.path.join(parent_dir, "2026-02-07_Internet_Backbone_and_ISPs_BEAUTIFUL.pptx")

    prs.save(output_path)
    print(f"[SUCCESS] Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
