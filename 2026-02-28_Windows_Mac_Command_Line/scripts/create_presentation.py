"""
PowerPoint Generator for: Windows & Mac Command Line — Everything Connects
Kids Computer Science Class - February 28, 2026

Run:
  pip install python-pptx
  python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ============================================================================
# COLOR PALETTE
# ============================================================================
COLORS = {
    # Primary
    "dark_blue": RGBColor(30, 58, 95),
    "medium_blue": RGBColor(52, 103, 166),
    "light_blue": RGBColor(100, 149, 237),
    "sky_blue": RGBColor(173, 216, 230),

    # Accents
    "orange": RGBColor(255, 153, 51),
    "green": RGBColor(76, 175, 80),
    "red": RGBColor(244, 67, 54),
    "purple": RGBColor(156, 39, 176),
    "teal": RGBColor(0, 150, 136),

    # Theme-specific
    "terminal_green": RGBColor(0, 200, 83),
    "windows_blue": RGBColor(0, 120, 215),
    "mac_gray": RGBColor(142, 142, 147),

    # Neutrals
    "white": RGBColor(255, 255, 255),
    "light_gray": RGBColor(245, 245, 245),
    "dark_gray": RGBColor(66, 66, 66),
    "black": RGBColor(33, 33, 33),
}

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def set_shape_fill(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_shape_line(shape, color, width_pt=1):
    line = shape.line
    line.color.rgb = color
    line.width = Pt(width_pt)

def add_styled_textbox(slide, left, top, width, height, text,
                       font_size=18, font_color=None, bold=False,
                       alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    if font_color is None:
        font_color = COLORS["dark_gray"]

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment

    return textbox

def add_title_shape(slide, text, subtitle=None):
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(1.3)
    )
    set_shape_fill(title_bar, COLORS["dark_blue"])
    title_bar.line.fill.background()

    add_styled_textbox(
        slide, Inches(0.5), Inches(0.35), Inches(12.33), Inches(0.7),
        text, font_size=36, font_color=COLORS["white"], bold=True,
        alignment=PP_ALIGN.LEFT, font_name="Segoe UI Semibold"
    )

    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.3),
        Inches(13.33), Inches(0.08)
    )
    set_shape_fill(accent_line, COLORS["orange"])
    accent_line.line.fill.background()

    if subtitle:
        add_styled_textbox(
            slide, Inches(0.5), Inches(0.95), Inches(12.33), Inches(0.4),
            subtitle, font_size=16, font_color=COLORS["sky_blue"],
            alignment=PP_ALIGN.LEFT
        )

def add_rounded_box(slide, left, top, width, height, fill_color, text=None,
                    text_size=16, text_color=None, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    set_shape_fill(shape, fill_color)

    if border_color:
        set_shape_line(shape, border_color, 2)
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(text_size)
        p.font.color.rgb = text_color or COLORS["white"]
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(8)

    shape.text_frame.anchor = MSO_ANCHOR.MIDDLE
    return shape

def add_arrow(slide, left, top, width, height, fill_color, direction="right"):
    if direction == "right":
        shape_type = MSO_SHAPE.RIGHT_ARROW
    elif direction == "left":
        shape_type = MSO_SHAPE.LEFT_ARROW
    elif direction == "down":
        shape_type = MSO_SHAPE.DOWN_ARROW
    else:
        shape_type = MSO_SHAPE.UP_ARROW

    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    set_shape_fill(shape, fill_color)
    shape.line.fill.background()
    return shape


def add_table_row(slide, y, cols, col_widths, col_starts, bg_color,
                  text_color=None, font_size=14, bold=False, row_height=0.55):
    """Helper to add a row of boxes that looks like a table row."""
    if text_color is None:
        text_color = COLORS["white"]
    for i, (text, w, x) in enumerate(zip(cols, col_widths, col_starts)):
        add_rounded_box(
            slide, Inches(x), Inches(y), Inches(w), Inches(row_height),
            bg_color, text, text_size=font_size, text_color=text_color
        )


# ============================================================================
# SLIDES
# ============================================================================

def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    set_shape_fill(background, COLORS["dark_blue"])
    background.line.fill.background()

    add_rounded_box(
        slide, Inches(1.2), Inches(1.7), Inches(10.9), Inches(3.1),
        COLORS["light_blue"]
    )

    add_styled_textbox(
        slide, Inches(0.6), Inches(2.0), Inches(12.13), Inches(1.2),
        "Windows & Mac Command Line",
        font_size=44, font_color=COLORS["dark_blue"], bold=True,
        alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold"
    )

    add_styled_textbox(
        slide, Inches(0.6), Inches(3.3), Inches(12.13), Inches(0.6),
        "Everything Connects: Your Linux Skills Work Everywhere!",
        font_size=22, font_color=COLORS["dark_blue"],
        alignment=PP_ALIGN.CENTER
    )

    add_styled_textbox(
        slide, Inches(0.6), Inches(6.6), Inches(12.13), Inches(0.4),
        "February 28, 2026",
        font_size=18, font_color=COLORS["sky_blue"],
        alignment=PP_ALIGN.CENTER
    )


def create_agenda_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Today's Journey")

    topics = [
        ("1", "Quick Recap: Linux Commands You Know", COLORS["light_blue"]),
        ("2", "Opening the Terminal on Every OS", COLORS["terminal_green"]),
        ("3", "The Big Translation Tables", COLORS["orange"]),
        ("4", "File Paths & Key Differences", COLORS["windows_blue"]),
        ("5", "Windows Bonus Tools & Mac's Secret", COLORS["teal"]),
        ("6", "AI Meets the Command Line", COLORS["purple"]),
        ("7", "Hands-on: Practice on YOUR Computer", COLORS["green"]),
    ]

    for i, (num, text, color) in enumerate(topics):
        y = Inches(1.7 + i * 0.72)

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1), y,
            Inches(0.55), Inches(0.55)
        )
        set_shape_fill(circle, color)
        circle.line.fill.background()

        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS["white"]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        add_styled_textbox(
            slide, Inches(2), y + Inches(0.1), Inches(10.5), Inches(0.5),
            text, font_size=20, font_color=COLORS["dark_gray"]
        )


def create_recap_week1_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Quick Recap: Linux Week 1", "Navigation & File Management")

    commands = [
        ("ls", "List files and folders", COLORS["light_blue"]),
        ("cd", "Change directory", COLORS["green"]),
        ("mkdir", "Make a new directory", COLORS["orange"]),
        ("touch", "Create an empty file", COLORS["teal"]),
        ("rm", "Remove a file", COLORS["red"]),
        ("rmdir", "Remove an empty folder", COLORS["purple"]),
    ]

    for i, (cmd, desc, color) in enumerate(commands):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(1.7 + (i // 2) * 1.5)

        add_rounded_box(slide, x, y, Inches(1.8), Inches(1.2), color,
                        cmd, text_size=22, text_color=COLORS["white"])
        add_styled_textbox(
            slide, x + Inches(2.0), y + Inches(0.35), Inches(3.8), Inches(0.5),
            desc, font_size=16, font_color=COLORS["dark_gray"]
        )

    add_rounded_box(
        slide, Inches(1.2), Inches(6.3), Inches(10.93), Inches(0.75),
        COLORS["light_blue"],
        "This week: Use these same skills on Windows and Mac!",
        text_size=18, text_color=COLORS["dark_blue"]
    )


def create_recap_week2_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Quick Recap: Linux Week 2", "File Content & Operations")

    commands = [
        ("echo", "Print text or write to files", COLORS["light_blue"]),
        ("cat", "View file contents", COLORS["green"]),
        ("cp", "Copy files", COLORS["orange"]),
        ("mv", "Move or rename files", COLORS["teal"]),
        ("> / >>", "Overwrite / Append", COLORS["purple"]),
        ("clear", "Clear the screen", COLORS["mac_gray"]),
    ]

    for i, (cmd, desc, color) in enumerate(commands):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(1.7 + (i // 2) * 1.5)

        add_rounded_box(slide, x, y, Inches(1.8), Inches(1.2), color,
                        cmd, text_size=20, text_color=COLORS["white"])
        add_styled_textbox(
            slide, x + Inches(2.0), y + Inches(0.35), Inches(3.8), Inches(0.5),
            desc, font_size=16, font_color=COLORS["dark_gray"]
        )

    add_rounded_box(
        slide, Inches(1.2), Inches(6.3), Inches(10.93), Inches(0.75),
        COLORS["orange"],
        "Plus from January: ipconfig, ping, tracert, nslookup",
        text_size=18, text_color=COLORS["white"]
    )


def create_opening_terminal_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Opening the Terminal on Every OS")

    # Three columns
    os_data = [
        ("Linux", "Ctrl + Alt + T", "Opens Terminal app", COLORS["terminal_green"]),
        ("Windows", "Win + R, type 'cmd'", "Opens Command Prompt", COLORS["windows_blue"]),
        ("Mac", "Cmd + Space, type 'Terminal'", "Opens Terminal app", COLORS["mac_gray"]),
    ]

    for i, (name, shortcut, desc, color) in enumerate(os_data):
        x = Inches(0.5 + i * 4.2)

        add_rounded_box(slide, x, Inches(1.8), Inches(3.8), Inches(1.0),
                        color, name, text_size=28, text_color=COLORS["white"])
        add_styled_textbox(
            slide, x, Inches(3.1), Inches(3.8), Inches(0.5),
            shortcut, font_size=16, font_color=COLORS["dark_gray"],
            bold=True, alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x, Inches(3.7), Inches(3.8), Inches(0.5),
            desc, font_size=14, font_color=COLORS["dark_gray"],
            alignment=PP_ALIGN.CENTER
        )

    add_rounded_box(
        slide, Inches(1.2), Inches(5.0), Inches(10.93), Inches(1.6),
        COLORS["light_gray"], border_color=COLORS["medium_blue"]
    )
    add_styled_textbox(
        slide, Inches(1.5), Inches(5.1), Inches(10.33), Inches(0.5),
        "All three look similar:", font_size=18, font_color=COLORS["dark_blue"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(1.5), Inches(5.6), Inches(10.33), Inches(0.8),
        "A window where you type commands and the computer responds!",
        font_size=16, font_color=COLORS["dark_gray"],
        alignment=PP_ALIGN.CENTER
    )


def create_translation_table_1_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Translation Table: File Commands")

    col_starts = [0.5, 3.25, 6.45, 9.65]
    col_widths = [2.55, 3.0, 3.0, 3.0]

    # Header row
    headers = ["What it does", "Linux", "Windows CMD", "Mac Terminal"]
    add_table_row(slide, 1.55, headers, col_widths, col_starts,
                  COLORS["dark_blue"], COLORS["white"], 13, True, 0.5)

    rows = [
        ("List files", "ls", "dir", "ls"),
        ("Change directory", "cd folder", "cd folder", "cd folder"),
        ("Go up one level", "cd ..", "cd ..", "cd .."),
        ("Create folder", "mkdir name", "mkdir name", "mkdir name"),
        ("Create empty file", "touch file", "echo. > file", "touch file"),
        ("Delete file", "rm file", "del file", "rm file"),
        ("Remove folder", "rmdir folder", "rmdir folder", "rmdir folder"),
    ]

    row_colors = [COLORS["light_gray"], COLORS["sky_blue"]]
    for i, (action, linux, windows, mac) in enumerate(rows):
        y = 2.15 + i * 0.62
        bg = row_colors[i % 2]
        add_table_row(slide, y, [action, linux, windows, mac],
                      col_widths, col_starts, bg, COLORS["dark_gray"], 13, False, 0.52)

    add_rounded_box(
        slide, Inches(0.8), Inches(6.6), Inches(11.73), Inches(0.65),
        COLORS["terminal_green"],
        "Many commands are the SAME across all three!",
        text_size=18, text_color=COLORS["white"]
    )


def create_translation_table_2_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Translation Table: Content & Operations")

    col_starts = [0.5, 3.25, 6.45, 9.65]
    col_widths = [2.55, 3.0, 3.0, 3.0]

    headers = ["What it does", "Linux", "Windows CMD", "Mac Terminal"]
    add_table_row(slide, 1.55, headers, col_widths, col_starts,
                  COLORS["dark_blue"], COLORS["white"], 13, True, 0.5)

    rows = [
        ("Print text", 'echo "text"', "echo text", 'echo "text"'),
        ("View file contents", "cat file", "type file", "cat file"),
        ("Copy file", "cp src dst", "copy src dst", "cp src dst"),
        ("Move file", "mv src dst", "move src dst", "mv src dst"),
        ("Rename file", "mv old new", "ren old new", "mv old new"),
        ("Overwrite to file", ">", ">", ">"),
        ("Append to file", ">>", ">>", ">>"),
        ("Clear screen", "clear", "cls", "clear"),
        ("Current directory", "pwd", "cd (no args)", "pwd"),
    ]

    row_colors = [COLORS["light_gray"], COLORS["sky_blue"]]
    for i, (action, linux, windows, mac) in enumerate(rows):
        y = 2.15 + i * 0.52
        bg = row_colors[i % 2]
        add_table_row(slide, y, [action, linux, windows, mac],
                      col_widths, col_starts, bg, COLORS["dark_gray"], 12, False, 0.44)

    add_rounded_box(
        slide, Inches(0.8), Inches(6.95), Inches(11.73), Inches(0.45),
        COLORS["orange"],
        "Only a few are truly different: dir, del, type, copy, move, ren, cls",
        text_size=14, text_color=COLORS["white"]
    )


def create_networking_commands_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Networking Commands Across OSes")

    col_starts = [0.5, 3.25, 6.45, 9.65]
    col_widths = [2.55, 3.0, 3.0, 3.0]

    headers = ["What it does", "Linux", "Windows CMD", "Mac Terminal"]
    add_table_row(slide, 1.7, headers, col_widths, col_starts,
                  COLORS["dark_blue"], COLORS["white"], 14, True, 0.55)

    rows = [
        ("Network info", "ifconfig", "ipconfig", "ifconfig"),
        ("Test connectivity", "ping site", "ping site", "ping site"),
        ("Trace route", "traceroute site", "tracert site", "traceroute site"),
        ("DNS lookup", "nslookup site", "nslookup site", "nslookup site"),
    ]

    row_colors = [COLORS["light_gray"], COLORS["sky_blue"]]
    for i, (action, linux, windows, mac) in enumerate(rows):
        y = 2.4 + i * 0.7
        bg = row_colors[i % 2]
        add_table_row(slide, y, [action, linux, windows, mac],
                      col_widths, col_starts, bg, COLORS["dark_gray"], 14, False, 0.6)

    add_rounded_box(
        slide, Inches(0.8), Inches(5.5), Inches(11.73), Inches(1.3),
        COLORS["light_blue"]
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(5.65), Inches(11.33), Inches(0.5),
        "You already used ipconfig, ping, tracert, and nslookup in January!",
        font_size=18, font_color=COLORS["dark_blue"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(6.2), Inches(11.33), Inches(0.5),
        "Now you know these commands exist on ALL operating systems.",
        font_size=16, font_color=COLORS["dark_blue"],
        alignment=PP_ALIGN.CENTER
    )


def create_file_paths_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "File Paths: The Biggest Difference")

    # Windows side
    add_rounded_box(
        slide, Inches(0.5), Inches(1.7), Inches(5.9), Inches(3.5),
        COLORS["light_gray"], border_color=COLORS["windows_blue"]
    )
    add_rounded_box(
        slide, Inches(0.5), Inches(1.7), Inches(5.9), Inches(0.65),
        COLORS["windows_blue"], "Windows", text_size=24, text_color=COLORS["white"]
    )
    points_win = [
        'Uses backslash:  C:\\Users\\Student\\Docs',
        'Starts with a drive letter:  C:\\',
        'NOT case-sensitive:',
        'MyFile.txt = myfile.TXT',
    ]
    for i, pt in enumerate(points_win):
        add_styled_textbox(
            slide, Inches(0.8), Inches(2.6 + i * 0.5), Inches(5.3), Inches(0.45),
            pt, font_size=15, font_color=COLORS["dark_gray"],
            bold=(i == 0 or i == 2)
        )

    # Linux/Mac side
    add_rounded_box(
        slide, Inches(6.93), Inches(1.7), Inches(5.9), Inches(3.5),
        COLORS["light_gray"], border_color=COLORS["terminal_green"]
    )
    add_rounded_box(
        slide, Inches(6.93), Inches(1.7), Inches(5.9), Inches(0.65),
        COLORS["terminal_green"], "Linux / Mac", text_size=24, text_color=COLORS["white"]
    )
    points_unix = [
        'Uses forward slash:  /home/student/docs',
        'Starts from root:  /',
        'IS case-sensitive:',
        'MyFile.txt != myfile.TXT',
    ]
    for i, pt in enumerate(points_unix):
        add_styled_textbox(
            slide, Inches(7.23), Inches(2.6 + i * 0.5), Inches(5.3), Inches(0.45),
            pt, font_size=15, font_color=COLORS["dark_gray"],
            bold=(i == 0 or i == 2)
        )

    # Arrow between
    add_styled_textbox(
        slide, Inches(4.5), Inches(5.8), Inches(4.33), Inches(0.6),
        "Same idea, different notation!",
        font_size=20, font_color=COLORS["medium_blue"], bold=True,
        alignment=PP_ALIGN.CENTER
    )


def create_key_differences_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Key Differences Summary")

    diffs = [
        ("Backslash vs Forward Slash", "Windows uses \\   Linux/Mac use /", COLORS["windows_blue"]),
        ("Case Sensitivity", "Windows ignores case, Linux/Mac care", COLORS["terminal_green"]),
        ("Different Command Names", "dir vs ls, type vs cat, del vs rm", COLORS["orange"]),
        ("Drive Letters", "Windows has C:\\, D:\\   Linux/Mac have /", COLORS["purple"]),
    ]

    for i, (title, desc, color) in enumerate(diffs):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(1.7 + (i // 2) * 2.2)

        add_rounded_box(slide, x, y, Inches(5.8), Inches(1.8), color)
        add_styled_textbox(
            slide, x + Inches(0.15), y + Inches(0.25), Inches(5.5), Inches(0.5),
            title, font_size=20, font_color=COLORS["white"],
            bold=True, alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x + Inches(0.15), y + Inches(0.9), Inches(5.5), Inches(0.6),
            desc, font_size=15, font_color=COLORS["white"],
            alignment=PP_ALIGN.CENTER
        )

    add_styled_textbox(
        slide, Inches(1), Inches(6.5), Inches(11.33), Inches(0.5),
        "These are the only major differences!",
        font_size=20, font_color=COLORS["medium_blue"], bold=True,
        alignment=PP_ALIGN.CENTER
    )


def create_windows_bonus_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Windows Bonus Tools")

    # Task Manager box
    add_rounded_box(
        slide, Inches(0.8), Inches(1.7), Inches(5.6), Inches(2.8),
        COLORS["windows_blue"]
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(1.85), Inches(5.2), Inches(0.5),
        "Task Manager", font_size=22, font_color=COLORS["white"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(2.45), Inches(5.2), Inches(0.4),
        "Ctrl + Shift + Esc", font_size=18, font_color=COLORS["sky_blue"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    tm_points = [
        "See all running programs",
        "Check CPU and memory usage",
        "End programs that are stuck",
    ]
    for i, pt in enumerate(tm_points):
        add_styled_textbox(
            slide, Inches(1.2), Inches(3.1 + i * 0.4), Inches(4.8), Inches(0.35),
            f"- {pt}", font_size=14, font_color=COLORS["white"]
        )

    # systeminfo box
    add_rounded_box(
        slide, Inches(6.93), Inches(1.7), Inches(5.6), Inches(2.8),
        COLORS["teal"]
    )
    add_styled_textbox(
        slide, Inches(7.13), Inches(1.85), Inches(5.2), Inches(0.5),
        "systeminfo", font_size=22, font_color=COLORS["white"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(7.13), Inches(2.45), Inches(5.2), Inches(0.4),
        "Type in Command Prompt", font_size=18, font_color=COLORS["sky_blue"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    si_points = [
        "Your computer's name and OS",
        "How much RAM you have",
        "When Windows was installed",
    ]
    for i, pt in enumerate(si_points):
        add_styled_textbox(
            slide, Inches(7.35), Inches(3.1 + i * 0.4), Inches(4.8), Inches(0.35),
            f"- {pt}", font_size=14, font_color=COLORS["white"]
        )

    add_rounded_box(
        slide, Inches(0.8), Inches(5.0), Inches(11.73), Inches(0.75),
        COLORS["light_gray"],
        "These are Windows-only tools that are super useful for troubleshooting!",
        text_size=16, text_color=COLORS["dark_blue"]
    )

    # Mac equivalent note
    add_rounded_box(
        slide, Inches(0.8), Inches(6.0), Inches(11.73), Inches(1.0),
        COLORS["mac_gray"]
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(6.05), Inches(11.33), Inches(0.4),
        "Mac equivalents:", font_size=14, font_color=COLORS["white"],
        bold=True, alignment=PP_ALIGN.LEFT
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(6.45), Inches(11.33), Inches(0.4),
        "Activity Monitor (like Task Manager)   |   sw_vers (shows macOS version)   |   system_profiler (full specs)",
        font_size=13, font_color=COLORS["white"],
        alignment=PP_ALIGN.LEFT
    )


def create_mac_secret_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Mac's Secret: It IS Basically Linux!")

    add_rounded_box(
        slide, Inches(0.8), Inches(1.7), Inches(11.73), Inches(1.2),
        COLORS["terminal_green"],
        "macOS is built on Unix -- the same family as Linux!",
        text_size=24, text_color=COLORS["white"]
    )

    # Unix family tree
    add_rounded_box(
        slide, Inches(4.67), Inches(3.3), Inches(4.0), Inches(0.8),
        COLORS["dark_blue"], "Unix (1970s)", text_size=20, text_color=COLORS["white"]
    )

    add_arrow(slide, Inches(4.5), Inches(4.2), Inches(0.5), Inches(0.5), COLORS["dark_gray"], "down")
    add_arrow(slide, Inches(8.5), Inches(4.2), Inches(0.5), Inches(0.5), COLORS["dark_gray"], "down")

    add_rounded_box(
        slide, Inches(2.5), Inches(4.8), Inches(3.5), Inches(0.8),
        COLORS["terminal_green"], "Linux (1991)", text_size=18, text_color=COLORS["white"]
    )
    add_rounded_box(
        slide, Inches(7.33), Inches(4.8), Inches(3.5), Inches(0.8),
        COLORS["mac_gray"], "macOS", text_size=18, text_color=COLORS["white"]
    )

    add_styled_textbox(
        slide, Inches(1), Inches(6.0), Inches(11.33), Inches(0.5),
        "If you know Linux commands, you already know Mac commands!",
        font_size=20, font_color=COLORS["dark_blue"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(1), Inches(6.6), Inches(11.33), Inches(0.5),
        "ls, cd, mkdir, touch, rm, cat, cp, mv, echo, clear -- all work on Mac!",
        font_size=16, font_color=COLORS["dark_gray"],
        alignment=PP_ALIGN.CENTER
    )


def create_big_picture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "The Big Picture: Same Concepts, Different Syntax")

    # Three OS boxes connected
    add_rounded_box(
        slide, Inches(0.8), Inches(1.8), Inches(3.5), Inches(1.2),
        COLORS["terminal_green"], "Linux", text_size=26, text_color=COLORS["white"]
    )
    add_rounded_box(
        slide, Inches(4.92), Inches(1.8), Inches(3.5), Inches(1.2),
        COLORS["windows_blue"], "Windows", text_size=26, text_color=COLORS["white"]
    )
    add_rounded_box(
        slide, Inches(9.03), Inches(1.8), Inches(3.5), Inches(1.2),
        COLORS["mac_gray"], "Mac", text_size=26, text_color=COLORS["white"]
    )

    # Connecting arrows
    add_arrow(slide, Inches(4.4), Inches(2.2), Inches(0.4), Inches(0.3), COLORS["orange"], "right")
    add_arrow(slide, Inches(8.53), Inches(2.2), Inches(0.4), Inches(0.3), COLORS["orange"], "right")

    # Center message
    add_rounded_box(
        slide, Inches(1.5), Inches(3.5), Inches(10.33), Inches(1.0),
        COLORS["orange"],
        "Every OS has a command line. The concepts are the same!",
        text_size=22, text_color=COLORS["white"]
    )

    points = [
        "Listing files = same idea everywhere",
        "Creating and deleting folders = same idea everywhere",
        "Navigating directories = same idea everywhere",
        "Network troubleshooting = same idea everywhere",
    ]
    for i, pt in enumerate(points):
        add_styled_textbox(
            slide, Inches(1.2), Inches(5.0 + i * 0.5), Inches(10.93), Inches(0.4),
            f"- {pt}", font_size=18, font_color=COLORS["dark_gray"]
        )


def create_ai_connection_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "AI Meets the Command Line", "Your skills connect to the future")

    # Left box: AI trains on Linux servers
    add_rounded_box(
        slide, Inches(0.5), Inches(1.7), Inches(4.0), Inches(3.6),
        COLORS["dark_blue"]
    )
    add_styled_textbox(
        slide, Inches(0.7), Inches(1.85), Inches(3.6), Inches(0.5),
        "AI Trains on Linux Servers",
        font_size=17, font_color=COLORS["orange"], bold=True,
        alignment=PP_ALIGN.CENTER
    )
    ai_cloud_points = [
        "ChatGPT, Claude, Gemini all",
        "train on thousands of GPUs",
        "in cloud data centers",
        "",
        "Those servers run Linux",
        "and are controlled through",
        "the SAME terminal commands",
        "you learned this month!",
    ]
    for i, pt in enumerate(ai_cloud_points):
        add_styled_textbox(
            slide, Inches(0.8), Inches(2.5 + i * 0.33), Inches(3.4), Inches(0.3),
            pt, font_size=13, font_color=COLORS["white"],
            alignment=PP_ALIGN.CENTER
        )

    # Middle box: Developers use AI IN the terminal
    add_rounded_box(
        slide, Inches(4.67), Inches(1.7), Inches(4.0), Inches(3.6),
        COLORS["terminal_green"]
    )
    add_styled_textbox(
        slide, Inches(4.87), Inches(1.85), Inches(3.6), Inches(0.5),
        "Devs Use AI in the Terminal",
        font_size=17, font_color=COLORS["dark_blue"], bold=True,
        alignment=PP_ALIGN.CENTER
    )
    ai_terminal_points = [
        "Claude Code (Anthropic)",
        "  AI assistant that runs",
        "  entirely in the terminal!",
        "",
        "Codex CLI (OpenAI)",
        "  AI coding tool that runs",
        "  in your command line!",
        "",
    ]
    for i, pt in enumerate(ai_terminal_points):
        add_styled_textbox(
            slide, Inches(4.97), Inches(2.5 + i * 0.33), Inches(3.4), Inches(0.3),
            pt, font_size=13, font_color=COLORS["white"],
            alignment=PP_ALIGN.CENTER
        )

    # Right box: Your journey
    add_rounded_box(
        slide, Inches(8.83), Inches(1.7), Inches(4.0), Inches(3.6),
        COLORS["purple"]
    )
    add_styled_textbox(
        slide, Inches(9.03), Inches(1.85), Inches(3.6), Inches(0.5),
        "Your Journey So Far",
        font_size=17, font_color=COLORS["orange"], bold=True,
        alignment=PP_ALIGN.CENTER
    )
    journey_points = [
        "Jan: Networking & Cloud",
        "  How data travels,",
        "  where servers live",
        "",
        "Feb: Command Line",
        "  How to talk to ANY",
        "  computer",
        "",
    ]
    for i, pt in enumerate(journey_points):
        add_styled_textbox(
            slide, Inches(9.13), Inches(2.5 + i * 0.33), Inches(3.4), Inches(0.3),
            pt, font_size=13, font_color=COLORS["white"],
            alignment=PP_ALIGN.CENTER
        )

    # Bottom banner
    add_rounded_box(
        slide, Inches(0.5), Inches(5.6), Inches(12.33), Inches(0.8),
        COLORS["orange"],
        "The terminal you are learning is how real AI engineers work every day!",
        text_size=20, text_color=COLORS["white"]
    )

    # Preview April
    add_rounded_box(
        slide, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.65),
        COLORS["light_gray"], border_color=COLORS["purple"]
    )
    add_styled_textbox(
        slide, Inches(0.7), Inches(6.65), Inches(11.93), Inches(0.5),
        "Coming in April: YOU will start using AI tools!  Everything you learn now prepares you for that.",
        font_size=15, font_color=COLORS["purple"], bold=True,
        alignment=PP_ALIGN.CENTER
    )


def create_practice_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Let's Practice on YOUR Computer!")

    add_rounded_box(
        slide, Inches(1), Inches(1.9), Inches(11.33), Inches(1.2),
        COLORS["terminal_green"],
        "Time to try it on your OWN machine!",
        text_size=24, text_color=COLORS["white"]
    )

    steps = [
        ("1", "Open Command Prompt (Windows) or Terminal (Mac)", COLORS["windows_blue"]),
        ("2", "Follow along with the classwork worksheet", COLORS["orange"]),
        ("3", "Try your Linux commands -- which ones work? Which don't?", COLORS["purple"]),
        ("4", "Use the translation table when you get stuck!", COLORS["green"]),
    ]

    for i, (num, desc, color) in enumerate(steps):
        y = Inches(3.5 + i * 0.8)

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.2), y + Inches(0.05),
            Inches(0.5), Inches(0.5)
        )
        set_shape_fill(circle, color)
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS["white"]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        add_styled_textbox(
            slide, Inches(2.0), y + Inches(0.1), Inches(10.0), Inches(0.5),
            desc, font_size=18, font_color=COLORS["dark_gray"]
        )


def create_recap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Today's Key Takeaways")

    takeaways = [
        ("1", "Linux commands map directly to Windows CMD and Mac Terminal"),
        ("2", "Windows unique: dir, type, del, copy, move, ren, cls"),
        ("3", "Mac Terminal = same as Linux (both are Unix-based)"),
        ("4", "Paths: Windows uses \\ and C:\\ , Linux/Mac use /"),
        ("5", "Windows is NOT case-sensitive, Linux/Mac ARE"),
        ("6", "The concepts are universal -- only the spelling changes!"),
    ]

    for i, (num, text) in enumerate(takeaways):
        y = Inches(1.6 + i * 0.82)

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8), y + Inches(0.05),
            Inches(0.5), Inches(0.5)
        )
        set_shape_fill(circle, COLORS["light_blue"])
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS["dark_blue"]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        add_styled_textbox(
            slide, Inches(1.6), y + Inches(0.1), Inches(10.5), Inches(0.5),
            text, font_size=18, font_color=COLORS["dark_gray"]
        )

    add_styled_textbox(
        slide, Inches(1), Inches(6.8), Inches(11.33), Inches(0.5),
        "Great job today! Time for Kahoot!",
        font_size=24, font_color=COLORS["dark_blue"], bold=True,
        alignment=PP_ALIGN.CENTER
    )


def create_questions_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    set_shape_fill(background, COLORS["dark_blue"])
    background.line.fill.background()

    add_styled_textbox(
        slide, Inches(0), Inches(1.6), Inches(13.33), Inches(2),
        "?", font_size=120, font_color=COLORS["white"],
        alignment=PP_ALIGN.CENTER
    )

    add_styled_textbox(
        slide, Inches(0), Inches(3.5), Inches(13.33), Inches(1),
        "Questions?", font_size=60, font_color=COLORS["white"], bold=True,
        alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold"
    )

    add_styled_textbox(
        slide, Inches(0), Inches(5.1), Inches(13.33), Inches(0.6),
        "Ask away before we start the Kahoot!",
        font_size=24, font_color=COLORS["sky_blue"],
        alignment=PP_ALIGN.CENTER
    )


# ============================================================================
# MAIN
# ============================================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs)           # 1
    create_agenda_slide(prs)          # 2
    create_recap_week1_slide(prs)     # 3
    create_recap_week2_slide(prs)     # 4
    create_opening_terminal_slide(prs)  # 5
    create_translation_table_1_slide(prs)  # 6
    create_translation_table_2_slide(prs)  # 7
    create_networking_commands_slide(prs)   # 8
    create_file_paths_slide(prs)      # 9
    create_key_differences_slide(prs)  # 10
    create_windows_bonus_slide(prs)   # 11
    create_mac_secret_slide(prs)      # 12
    create_big_picture_slide(prs)     # 13
    create_ai_connection_slide(prs)   # 14
    create_practice_slide(prs)        # 15
    create_recap_slide(prs)           # 16
    create_questions_slide(prs)       # 17

    script_dir = os.path.dirname(os.path.abspath(__file__))
    parent_dir = os.path.dirname(script_dir)
    output_path = os.path.join(parent_dir, "2026-02-28_Windows_Mac_Command_Line_BEAUTIFUL.pptx")

    prs.save(output_path)
    print(f"[SUCCESS] Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
