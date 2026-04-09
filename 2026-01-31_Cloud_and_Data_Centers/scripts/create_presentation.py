"""
PowerPoint Generator for: The Cloud & Data Centers
Kids Computer Science Class - January 31, 2026

This script creates a beautifully designed presentation using python-pptx.
Run: pip install python-pptx
Then: python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ============================================================================
# COLOR PALETTE - Modern, Professional, Kid-Friendly
# ============================================================================
COLORS = {
    # Primary colors
    'dark_blue': RGBColor(30, 58, 95),
    'medium_blue': RGBColor(52, 103, 166),
    'light_blue': RGBColor(100, 149, 237),
    'sky_blue': RGBColor(173, 216, 230),

    # Accent colors
    'orange': RGBColor(255, 153, 51),
    'green': RGBColor(76, 175, 80),
    'red': RGBColor(244, 67, 54),
    'purple': RGBColor(156, 39, 176),
    'teal': RGBColor(0, 150, 136),

    # Neutrals
    'white': RGBColor(255, 255, 255),
    'light_gray': RGBColor(245, 245, 245),
    'dark_gray': RGBColor(66, 66, 66),
    'black': RGBColor(33, 33, 33),

    # Cloud theme colors
    'cloud_blue': RGBColor(135, 206, 250),
    'data_center_gray': RGBColor(96, 125, 139),
    'aws_orange': RGBColor(255, 153, 0),
    'azure_blue': RGBColor(0, 120, 212),
    'google_blue': RGBColor(66, 133, 244),
}

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def set_shape_fill(shape, color):
    """Set solid fill color for a shape."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_shape_line(shape, color, width_pt=1):
    """Set line/border color and width for a shape."""
    line = shape.line
    line.color.rgb = color
    line.width = Pt(width_pt)

def add_styled_textbox(slide, left, top, width, height, text,
                       font_size=18, font_color=None, bold=False,
                       alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Add a textbox with styled text."""
    if font_color is None:
        font_color = COLORS['dark_gray']

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment

    return textbox

def add_title_shape(slide, text, subtitle=None):
    """Add a styled title bar at the top of the slide."""
    # Title background bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(1.3)
    )
    set_shape_fill(title_bar, COLORS['dark_blue'])
    title_bar.line.fill.background()

    # Title text
    add_styled_textbox(
        slide, Inches(0.5), Inches(0.35), Inches(12.33), Inches(0.7),
        text, font_size=36, font_color=COLORS['white'], bold=True,
        alignment=PP_ALIGN.LEFT, font_name="Segoe UI Semibold"
    )

    # Accent line under title
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.3),
        Inches(13.33), Inches(0.08)
    )
    set_shape_fill(accent_line, COLORS['orange'])
    accent_line.line.fill.background()

    if subtitle:
        add_styled_textbox(
            slide, Inches(0.5), Inches(0.95), Inches(12.33), Inches(0.4),
            subtitle, font_size=16, font_color=COLORS['sky_blue'],
            alignment=PP_ALIGN.LEFT
        )

def add_bullet_points(slide, left, top, width, height, items, font_size=20):
    """Add a text frame with bullet points."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS['dark_gray']
        p.font.name = "Segoe UI"
        p.space_before = Pt(8)
        p.space_after = Pt(8)

    return textbox

def add_rounded_box(slide, left, top, width, height, fill_color, text=None,
                    text_size=16, text_color=None, border_color=None):
    """Add a rounded rectangle with optional text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    set_shape_fill(shape, fill_color)

    if border_color:
        set_shape_line(shape, border_color, 2)
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(text_size)
        p.font.color.rgb = text_color or COLORS['white']
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(10)

    shape.text_frame.anchor = MSO_ANCHOR.MIDDLE
    return shape

def add_arrow(slide, left, top, width, height, fill_color, direction='right'):
    """Add an arrow shape."""
    if direction == 'right':
        shape_type = MSO_SHAPE.RIGHT_ARROW
    elif direction == 'left':
        shape_type = MSO_SHAPE.LEFT_ARROW
    elif direction == 'down':
        shape_type = MSO_SHAPE.DOWN_ARROW
    else:
        shape_type = MSO_SHAPE.UP_ARROW

    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    set_shape_fill(shape, fill_color)
    shape.line.fill.background()
    return shape

def add_cloud_shape(slide, left, top, width, height, fill_color):
    """Add a cloud shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CLOUD,
        left, top, width, height
    )
    set_shape_fill(shape, fill_color)
    shape.line.fill.background()
    return shape

# ============================================================================
# SLIDE CREATION FUNCTIONS
# ============================================================================

def create_title_slide(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    set_shape_fill(background, COLORS['dark_blue'])
    background.line.fill.background()

    # Decorative clouds
    for x, y, size in [(10.5, 0.8, 2.2), (0.3, 5.5, 1.8), (-0.3, 6.2, 1.2)]:
        cloud = add_cloud_shape(slide, Inches(x), Inches(y), Inches(size), Inches(size * 0.6), COLORS['medium_blue'])

    # Main cloud behind title
    main_cloud = add_cloud_shape(slide, Inches(2), Inches(1.5), Inches(9), Inches(3.5), COLORS['cloud_blue'])

    # Main title
    add_styled_textbox(
        slide, Inches(0.5), Inches(2.4), Inches(12.33), Inches(1.2),
        "The Cloud & Data Centers",
        font_size=52, font_color=COLORS['dark_blue'], bold=True,
        alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold"
    )

    # Subtitle
    add_styled_textbox(
        slide, Inches(0.5), Inches(3.6), Inches(12.33), Inches(0.6),
        "Where the Internet Actually Lives",
        font_size=28, font_color=COLORS['dark_blue'],
        alignment=PP_ALIGN.CENTER
    )

    # Date
    add_styled_textbox(
        slide, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.5),
        "January 31, 2026",
        font_size=18, font_color=COLORS['sky_blue'],
        alignment=PP_ALIGN.CENTER
    )

def create_agenda_slide(prs):
    """Slide 2: Today's agenda."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Today's Journey")

    topics = [
        ("1", "Quick Review: Servers Refresher", COLORS['light_blue']),
        ("2", "What is 'The Cloud'?", COLORS['cloud_blue']),
        ("3", "Inside a Data Center", COLORS['data_center_gray']),
        ("4", "The Big Three: AWS, Google, Azure", COLORS['orange']),
        ("5", "AI + Cloud: ChatGPT, Grok & More!", COLORS['purple']),
        ("6", "Hands-on: Cloud Explorer", COLORS['green']),
    ]

    for i, (num, text, color) in enumerate(topics):
        y = Inches(1.7 + i * 0.85)

        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1), y,
            Inches(0.65), Inches(0.65)
        )
        set_shape_fill(circle, color)
        circle.line.fill.background()

        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS['white']
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Topic text
        add_styled_textbox(
            slide, Inches(2), y + Inches(0.12), Inches(9), Inches(0.5),
            text, font_size=22, font_color=COLORS['dark_gray']
        )

def create_review_slide(prs):
    """Slide 3: Quick review of servers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Quick Review: What's a Server?", "From last week...")

    # Server definition box
    add_rounded_box(
        slide, Inches(0.8), Inches(1.8), Inches(11.73), Inches(1.1),
        COLORS['medium_blue'],
        "A SERVER is a computer that stores data and serves it to clients",
        text_size=26, text_color=COLORS['white']
    )

    # Visual: client-server
    # Client
    client_box = add_rounded_box(
        slide, Inches(1), Inches(3.5), Inches(2.8), Inches(1.8),
        COLORS['light_blue']
    )
    add_styled_textbox(
        slide, Inches(1), Inches(3.7), Inches(2.8), Inches(0.6),
        "Your Device", font_size=18, font_color=COLORS['dark_blue'],
        alignment=PP_ALIGN.CENTER, bold=True
    )
    add_styled_textbox(
        slide, Inches(1), Inches(4.2), Inches(2.8), Inches(0.6),
        "(Client)", font_size=14, font_color=COLORS['dark_blue'],
        alignment=PP_ALIGN.CENTER
    )

    # Arrow
    add_arrow(slide, Inches(4.2), Inches(4.1), Inches(1.8), Inches(0.5), COLORS['orange'], 'right')

    # Server
    server_box = add_rounded_box(
        slide, Inches(6.5), Inches(3.5), Inches(2.8), Inches(1.8),
        COLORS['green']
    )
    add_styled_textbox(
        slide, Inches(6.5), Inches(3.7), Inches(2.8), Inches(0.6),
        "Server", font_size=18, font_color=COLORS['white'],
        alignment=PP_ALIGN.CENTER, bold=True
    )
    add_styled_textbox(
        slide, Inches(6.5), Inches(4.2), Inches(2.8), Inches(0.6),
        "(Stores websites)", font_size=14, font_color=COLORS['white'],
        alignment=PP_ALIGN.CENTER
    )

    # Question
    add_rounded_box(
        slide, Inches(9.8), Inches(3.5), Inches(3), Inches(1.8),
        COLORS['purple'],
        "But WHERE\nare these\nservers?",
        text_size=20, text_color=COLORS['white']
    )

    # Key question
    add_styled_textbox(
        slide, Inches(1), Inches(6), Inches(11.33), Inches(0.8),
        "Today's Big Question: Where do all these servers actually live?",
        font_size=22, font_color=COLORS['dark_blue'], bold=True,
        alignment=PP_ALIGN.CENTER
    )

def create_cloud_myth_slide(prs):
    """Slide 4: The cloud is NOT in the sky!"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Wait... What IS 'The Cloud'?")

    # Myth box - no text in box, use separate textboxes
    add_rounded_box(
        slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.2),
        COLORS['red']
    )
    add_styled_textbox(
        slide, Inches(0.8), Inches(1.95), Inches(5.5), Inches(0.5),
        "What People Think:",
        font_size=20, font_color=COLORS['white'],
        alignment=PP_ALIGN.CENTER, bold=True
    )
    add_styled_textbox(
        slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(1.2),
        "Some magical\nfloating computers\nin the sky",
        font_size=18, font_color=COLORS['white'],
        alignment=PP_ALIGN.CENTER
    )

    # Reality box - no text in box, use separate textboxes
    add_rounded_box(
        slide, Inches(7), Inches(1.8), Inches(5.5), Inches(2.2),
        COLORS['green']
    )
    add_styled_textbox(
        slide, Inches(7), Inches(1.95), Inches(5.5), Inches(0.5),
        "What It Actually Is:",
        font_size=20, font_color=COLORS['white'],
        alignment=PP_ALIGN.CENTER, bold=True
    )
    add_styled_textbox(
        slide, Inches(7), Inches(2.5), Inches(5.5), Inches(1.2),
        "Other people's\ncomputers in\nBIG buildings",
        font_size=18, font_color=COLORS['white'],
        alignment=PP_ALIGN.CENTER
    )

    # Big reveal
    add_rounded_box(
        slide, Inches(2), Inches(4.5), Inches(9.33), Inches(1.5),
        COLORS['orange'],
        '"The Cloud" = Servers in Data Centers Around the World',
        text_size=28, text_color=COLORS['white']
    )

    # Fun fact
    add_styled_textbox(
        slide, Inches(1), Inches(6.3), Inches(11.33), Inches(0.8),
        "There's a famous saying: \"There is no cloud... it's just someone else's computer!\"",
        font_size=18, font_color=COLORS['medium_blue'],
        alignment=PP_ALIGN.CENTER
    )

def create_data_center_intro_slide(prs):
    """Slide 5: What is a Data Center?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "What is a Data Center?")

    # Definition
    add_rounded_box(
        slide, Inches(0.8), Inches(1.8), Inches(11.73), Inches(1),
        COLORS['data_center_gray'],
        "A HUGE building filled with thousands of servers",
        text_size=28, text_color=COLORS['white']
    )

    # Features
    features = [
        ("Thousands of Servers", "Rows and rows of computers"),
        ("24/7 Power", "Never turns off - backup generators"),
        ("Massive Cooling", "Servers get HOT - needs AC"),
        ("High Security", "Guards, cameras, locked doors"),
        ("Super Fast Internet", "Connected to the backbone"),
    ]

    for i, (title, desc) in enumerate(features):
        x = Inches(0.8 + (i % 3) * 4.1)
        y = Inches(3.2 + (i // 3) * 1.8)

        box = add_rounded_box(
            slide, x, y, Inches(3.8), Inches(1.5),
            COLORS['dark_blue']
        )

        add_styled_textbox(
            slide, x, y + Inches(0.25), Inches(3.8), Inches(0.5),
            title, font_size=16, font_color=COLORS['orange'],
            alignment=PP_ALIGN.CENTER, bold=True
        )
        add_styled_textbox(
            slide, x, y + Inches(0.7), Inches(3.8), Inches(0.6),
            desc, font_size=13, font_color=COLORS['white'],
            alignment=PP_ALIGN.CENTER
        )

def create_data_center_scale_slide(prs):
    """Slide 6: Scale of data centers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "How Big Are Data Centers?", "Mind-blowing numbers!")

    # Fun facts
    facts = [
        ("Google", "Has 30+ data centers worldwide", "2+ million servers total!", COLORS['google_blue']),
        ("Amazon (AWS)", "Data centers in 30+ countries", "Millions of servers", COLORS['aws_orange']),
        ("Microsoft", "200+ data centers globally", "Powers Xbox, Teams, Office", COLORS['azure_blue']),
    ]

    for i, (company, fact1, fact2, color) in enumerate(facts):
        y = Inches(1.7 + i * 1.5)

        # Company badge
        badge = add_rounded_box(
            slide, Inches(0.8), y, Inches(2.5), Inches(1.25),
            color,
            company,
            text_size=22, text_color=COLORS['white']
        )

        # Facts
        add_styled_textbox(
            slide, Inches(3.6), y + Inches(0.15), Inches(9), Inches(0.5),
            fact1, font_size=20, font_color=COLORS['dark_gray'], bold=True
        )
        add_styled_textbox(
            slide, Inches(3.6), y + Inches(0.65), Inches(9), Inches(0.5),
            fact2, font_size=18, font_color=COLORS['dark_gray']
        )

    # Mind-blowing stat
    add_rounded_box(
        slide, Inches(1.5), Inches(6.4), Inches(10.33), Inches(0.85),
        COLORS['purple'],
        "A single data center can use as much electricity as a small CITY!",
        text_size=20, text_color=COLORS['white']
    )

def create_cloud_services_intro_slide(prs):
    """Slide 7: Types of cloud services."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "What Can the Cloud Do?", "Three main things...")

    services = [
        ("STORAGE", "Store your files", "Google Drive, iCloud,\nDropbox, OneDrive", COLORS['light_blue']),
        ("COMPUTE", "Run programs", "Netflix streaming,\nOnline games, AI", COLORS['green']),
        ("NETWORKING", "Connect things", "Video calls, messaging,\nwebsites", COLORS['orange']),
    ]

    for i, (title, subtitle, examples, color) in enumerate(services):
        x = Inches(0.8 + i * 4.1)

        # Main box
        box = add_rounded_box(
            slide, x, Inches(1.8), Inches(3.8), Inches(4.5),
            color
        )

        # Title
        add_styled_textbox(
            slide, x, Inches(2.1), Inches(3.8), Inches(0.6),
            title, font_size=26, font_color=COLORS['white'],
            alignment=PP_ALIGN.CENTER, bold=True
        )

        # Subtitle
        add_styled_textbox(
            slide, x, Inches(2.7), Inches(3.8), Inches(0.5),
            subtitle, font_size=16, font_color=COLORS['white'],
            alignment=PP_ALIGN.CENTER
        )

        # Divider
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.5), Inches(3.4),
            Inches(2.8), Inches(0.05)
        )
        set_shape_fill(divider, COLORS['white'])
        divider.line.fill.background()

        # Examples
        add_styled_textbox(
            slide, x, Inches(3.6), Inches(3.8), Inches(1.5),
            examples, font_size=15, font_color=COLORS['white'],
            alignment=PP_ALIGN.CENTER
        )

    # Key point
    add_styled_textbox(
        slide, Inches(1), Inches(6.6), Inches(11.33), Inches(0.6),
        "You use ALL of these every single day!",
        font_size=20, font_color=COLORS['medium_blue'], bold=True,
        alignment=PP_ALIGN.CENTER
    )

def create_cloud_you_use_slide(prs):
    """Slide 8: Cloud services you already use."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Cloud Services You Already Use!")

    services = [
        ("YouTube", "Videos stored on Google's servers", COLORS['red']),
        ("Netflix", "Movies streamed from data centers", COLORS['red']),
        ("Google Drive", "Your files in the cloud", COLORS['google_blue']),
        ("iCloud", "iPhone backups & photos", COLORS['light_blue']),
        ("Spotify", "Music streamed from servers", COLORS['green']),
        ("Fortnite/Roblox", "Game data on cloud servers", COLORS['purple']),
        ("TikTok", "Videos stored in data centers", COLORS['dark_gray']),
        ("Discord", "Messages saved in the cloud", COLORS['purple']),
    ]

    for i, (name, desc, color) in enumerate(services):
        col = i % 4
        row = i // 4
        x = Inches(0.5 + col * 3.2)
        y = Inches(1.8 + row * 2.6)

        # Service card
        card = add_rounded_box(
            slide, x, y, Inches(3), Inches(2.3),
            COLORS['light_gray'], border_color=color
        )

        # Icon bar
        icon_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x, y,
            Inches(3), Inches(0.6)
        )
        set_shape_fill(icon_bar, color)
        icon_bar.line.fill.background()

        add_styled_textbox(
            slide, x, y + Inches(0.1), Inches(3), Inches(0.5),
            name, font_size=16, font_color=COLORS['white'],
            alignment=PP_ALIGN.CENTER, bold=True
        )

        add_styled_textbox(
            slide, x + Inches(0.15), y + Inches(0.8), Inches(2.7), Inches(1.2),
            desc, font_size=13, font_color=COLORS['dark_gray'],
            alignment=PP_ALIGN.CENTER
        )

def create_no_cloud_slide(prs):
    """Slide 9: What happens without the cloud?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "What If There Was No Cloud?")

    scenarios = [
        ("No YouTube", "You'd have to download EVERY video first"),
        ("No Streaming", "Back to buying DVDs and CDs!"),
        ("No Cloud Saves", "Lose your phone? Lose EVERYTHING"),
        ("No Online Gaming", "No multiplayer, no updates"),
        ("No Google Docs", "Can't work on projects together"),
        ("No Video Calls", "No Zoom, Teams, or FaceTime"),
    ]

    for i, (title, desc) in enumerate(scenarios):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 6.3)
        y = Inches(1.7 + row * 1.5)

        # Box
        box = add_rounded_box(
            slide, x, y, Inches(5.8), Inches(1.25),
            COLORS['light_gray'], border_color=COLORS['red']
        )

        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.15), Inches(5.4), Inches(0.5),
            title, font_size=18, font_color=COLORS['red'],
            bold=True
        )
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.6), Inches(5.4), Inches(0.5),
            desc, font_size=14, font_color=COLORS['dark_gray']
        )

    # Summary
    add_rounded_box(
        slide, Inches(2), Inches(6.4), Inches(9.33), Inches(0.85),
        COLORS['green'],
        "The cloud makes our digital life possible!",
        text_size=22, text_color=COLORS['white']
    )

def create_big_three_slide(prs):
    """Slide 10: The Big Three cloud providers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "The Big Three Cloud Providers")

    providers = [
        ("Amazon Web Services\n(AWS)", "The biggest!\n~33% of cloud market", "Netflix, Twitch,\nNASA, McDonald's", COLORS['aws_orange']),
        ("Microsoft Azure", "Second biggest!\n~22% of cloud market", "Xbox Live, LinkedIn,\nAdobe, Samsung", COLORS['azure_blue']),
        ("Google Cloud", "Third biggest!\n~10% of cloud market", "Spotify, Twitter,\nSnapchat, PayPal", COLORS['google_blue']),
    ]

    for i, (name, share, users, color) in enumerate(providers):
        x = Inches(0.6 + i * 4.2)

        # Main card
        card = add_rounded_box(
            slide, x, Inches(1.8), Inches(4), Inches(5),
            color
        )

        # Name
        add_styled_textbox(
            slide, x, Inches(2.1), Inches(4), Inches(1),
            name, font_size=20, font_color=COLORS['white'],
            alignment=PP_ALIGN.CENTER, bold=True
        )

        # Market share
        add_styled_textbox(
            slide, x, Inches(3.2), Inches(4), Inches(0.9),
            share, font_size=14, font_color=COLORS['white'],
            alignment=PP_ALIGN.CENTER
        )

        # Divider
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.4), Inches(4.3),
            Inches(3.2), Inches(0.04)
        )
        set_shape_fill(divider, COLORS['white'])
        divider.line.fill.background()

        # Who uses it
        add_styled_textbox(
            slide, x, Inches(4.5), Inches(4), Inches(0.5),
            "Who uses it:", font_size=13, font_color=COLORS['white'],
            alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x, Inches(5), Inches(4), Inches(1.2),
            users, font_size=13, font_color=COLORS['white'],
            alignment=PP_ALIGN.CENTER
        )

def create_how_cloud_works_slide(prs):
    """Slide 11: How cloud computing works."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "How Does Cloud Computing Work?")

    steps = [
        ("1", "You open an app\n(YouTube, Drive, etc.)", COLORS['light_blue']),
        ("2", "Your request travels\nthrough the internet", COLORS['orange']),
        ("3", "Arrives at a\ndata center", COLORS['data_center_gray']),
        ("4", "Server processes\nyour request", COLORS['green']),
        ("5", "Response sent\nback to you!", COLORS['purple']),
    ]

    for i, (num, text, color) in enumerate(steps):
        x = Inches(0.4 + i * 2.5)

        # Step circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + Inches(0.75), Inches(2),
            Inches(0.9), Inches(0.9)
        )
        set_shape_fill(circle, color)
        circle.line.fill.background()

        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(28)
        p.font.color.rgb = COLORS['white']
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Step box
        box = add_rounded_box(
            slide, x, Inches(3.2), Inches(2.4), Inches(1.8),
            color
        )
        add_styled_textbox(
            slide, x, Inches(3.5), Inches(2.4), Inches(1.4),
            text, font_size=13, font_color=COLORS['white'],
            alignment=PP_ALIGN.CENTER
        )

        # Arrow between steps
        if i < 4:
            add_arrow(slide, x + Inches(2.45), Inches(2.3), Inches(0.45), Inches(0.35), COLORS['dark_gray'], 'right')

    # Speed note
    add_rounded_box(
        slide, Inches(2), Inches(5.5), Inches(9.33), Inches(1.2),
        COLORS['sky_blue'],
        "All of this happens in MILLISECONDS!\n(faster than you can blink)",
        text_size=20, text_color=COLORS['dark_blue']
    )

def create_local_vs_cloud_slide(prs):
    """Slide 12: Local storage vs cloud storage."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Local vs. Cloud Storage")

    # Local storage
    local_box = add_rounded_box(
        slide, Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.5),
        COLORS['light_gray'], border_color=COLORS['orange']
    )

    add_styled_textbox(
        slide, Inches(0.8), Inches(1.9), Inches(5.8), Inches(0.6),
        "LOCAL STORAGE", font_size=24, font_color=COLORS['orange'],
        alignment=PP_ALIGN.CENTER, bold=True
    )
    add_styled_textbox(
        slide, Inches(0.8), Inches(2.4), Inches(5.8), Inches(0.5),
        "(On YOUR device)", font_size=14, font_color=COLORS['dark_gray'],
        alignment=PP_ALIGN.CENTER
    )

    local_points = [
        " Your phone's storage",
        " Computer's hard drive",
        " USB flash drives",
        " Works offline",
        " Can be lost/damaged"
    ]
    for i, point in enumerate(local_points):
        add_styled_textbox(
            slide, Inches(1.2), Inches(3 + i * 0.5), Inches(5), Inches(0.5),
            point, font_size=16, font_color=COLORS['dark_gray']
        )

    # Cloud storage
    cloud_box = add_rounded_box(
        slide, Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.5),
        COLORS['light_gray'], border_color=COLORS['light_blue']
    )

    add_styled_textbox(
        slide, Inches(6.9), Inches(1.9), Inches(5.8), Inches(0.6),
        "CLOUD STORAGE", font_size=24, font_color=COLORS['light_blue'],
        alignment=PP_ALIGN.CENTER, bold=True
    )
    add_styled_textbox(
        slide, Inches(6.9), Inches(2.4), Inches(5.8), Inches(0.5),
        "(On remote servers)", font_size=14, font_color=COLORS['dark_gray'],
        alignment=PP_ALIGN.CENTER
    )

    cloud_points = [
        " Google Drive, iCloud",
        " Accessible anywhere",
        " Backed up automatically",
        " Needs internet",
        " Shared across devices"
    ]
    for i, point in enumerate(cloud_points):
        add_styled_textbox(
            slide, Inches(7.3), Inches(3 + i * 0.5), Inches(5), Inches(0.5),
            point, font_size=16, font_color=COLORS['dark_gray']
        )

    # Bottom note
    add_styled_textbox(
        slide, Inches(1), Inches(6.5), Inches(11.33), Inches(0.6),
        "Best practice: Use BOTH! Cloud for backup, local for speed.",
        font_size=18, font_color=COLORS['medium_blue'], bold=True,
        alignment=PP_ALIGN.CENTER
    )

def create_cloud_safety_slide(prs):
    """Slide 13: Cloud safety and privacy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Cloud Safety & Privacy")

    # Important box
    add_rounded_box(
        slide, Inches(0.8), Inches(1.7), Inches(11.73), Inches(1),
        COLORS['orange'],
        "Remember: When you put something in 'the cloud,' you're storing it on someone else's computers!",
        text_size=18, text_color=COLORS['white']
    )

    # Safety tips
    tips = [
        ("Use Strong Passwords", "Don't use 'password123'!"),
        ("Enable 2FA", "Two-factor authentication adds extra protection"),
        ("Know What You Share", "Be careful what you upload"),
        ("Check Privacy Settings", "Control who sees your stuff"),
    ]

    for i, (title, desc) in enumerate(tips):
        x = Inches(0.8 + (i % 2) * 6.3)
        y = Inches(3 + (i // 2) * 1.6)

        box = add_rounded_box(
            slide, x, y, Inches(5.8), Inches(1.4),
            COLORS['green']
        )

        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.2), Inches(5.4), Inches(0.5),
            title, font_size=18, font_color=COLORS['white'], bold=True
        )
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.7), Inches(5.4), Inches(0.5),
            desc, font_size=14, font_color=COLORS['white']
        )

    # Key reminder
    add_styled_textbox(
        slide, Inches(1), Inches(6.4), Inches(11.33), Inches(0.8),
        "The cloud is convenient, but always think before you upload!",
        font_size=18, font_color=COLORS['dark_blue'], bold=True,
        alignment=PP_ALIGN.CENTER
    )

def create_connection_slide(prs):
    """Slide 14: Connecting last week to this week."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Putting It All Together")

    # Visual journey
    # Client
    client = add_rounded_box(
        slide, Inches(0.5), Inches(2.5), Inches(2.3), Inches(1.6),
        COLORS['light_blue'],
        "Your Device\n(Client)",
        text_size=16, text_color=COLORS['dark_blue']
    )

    # Arrow 1
    add_arrow(slide, Inches(3), Inches(3.1), Inches(0.8), Inches(0.4), COLORS['orange'], 'right')
    add_styled_textbox(
        slide, Inches(2.9), Inches(2.4), Inches(1), Inches(0.4),
        "Request", font_size=12, font_color=COLORS['orange'],
        alignment=PP_ALIGN.CENTER
    )

    # Internet cloud
    internet = add_cloud_shape(
        slide, Inches(4), Inches(2.2), Inches(2.5), Inches(2),
        COLORS['sky_blue']
    )
    add_styled_textbox(
        slide, Inches(4), Inches(2.9), Inches(2.5), Inches(0.6),
        "Internet", font_size=16, font_color=COLORS['dark_blue'],
        alignment=PP_ALIGN.CENTER, bold=True
    )

    # Arrow 2
    add_arrow(slide, Inches(6.7), Inches(3.1), Inches(0.8), Inches(0.4), COLORS['orange'], 'right')

    # Data Center
    dc = add_rounded_box(
        slide, Inches(7.7), Inches(2.2), Inches(5), Inches(2.2),
        COLORS['data_center_gray']
    )
    add_styled_textbox(
        slide, Inches(7.7), Inches(2.4), Inches(5), Inches(0.5),
        "DATA CENTER", font_size=18, font_color=COLORS['white'],
        alignment=PP_ALIGN.CENTER, bold=True
    )

    # Servers inside
    for j in range(4):
        server = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(8 + j * 1.1), Inches(3),
            Inches(0.9), Inches(1.2)
        )
        set_shape_fill(server, COLORS['green'])
        server.line.fill.background()

    add_styled_textbox(
        slide, Inches(7.7), Inches(4.1), Inches(5), Inches(0.5),
        "Servers", font_size=12, font_color=COLORS['white'],
        alignment=PP_ALIGN.CENTER
    )

    # Summary points
    summary = [
        "Last week: You learned about CLIENTS and SERVERS",
        "This week: You learned WHERE those servers actually live!",
        "Next: How the internet connects everything together"
    ]

    for i, point in enumerate(summary):
        y = Inches(5 + i * 0.5)
        add_styled_textbox(
            slide, Inches(1), y, Inches(11.33), Inches(0.5),
            point, font_size=18, font_color=COLORS['dark_gray'],
            alignment=PP_ALIGN.CENTER
        )

def create_ai_needs_cloud_slide(prs):
    """Slide 15: Why AI needs the cloud."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Why AI Needs the Cloud", "The secret behind ChatGPT, DALL-E, and more!")

    # Main explanation
    add_rounded_box(
        slide, Inches(0.8), Inches(1.7), Inches(11.73), Inches(1.1),
        COLORS['purple'],
        "AI needs MASSIVE computing power - more than any single computer can provide!",
        text_size=24, text_color=COLORS['white']
    )

    # Three reasons
    reasons = [
        ("TRAINING", "Teaching AI requires\nprocessing BILLIONS\nof examples", "GPT-4 trained on\n45TB of text data!"),
        ("COMPUTING", "AI calculations need\nthousands of powerful\nGPUs working together", "One AI query = millions\nof calculations!"),
        ("STORAGE", "AI models are HUGE\nGPT-4 has 1.7 trillion\nparameters", "Too big for your\nlaptop or phone!"),
    ]

    for i, (title, desc, fact) in enumerate(reasons):
        x = Inches(0.7 + i * 4.1)

        # Reason box
        box = add_rounded_box(
            slide, x, Inches(3.1), Inches(3.9), Inches(2.8),
            COLORS['dark_blue']
        )

        add_styled_textbox(
            slide, x, Inches(3.25), Inches(3.9), Inches(0.5),
            title, font_size=20, font_color=COLORS['orange'],
            alignment=PP_ALIGN.CENTER, bold=True
        )

        add_styled_textbox(
            slide, x, Inches(3.8), Inches(3.9), Inches(1),
            desc, font_size=14, font_color=COLORS['white'],
            alignment=PP_ALIGN.CENTER
        )

        # Fun fact
        add_styled_textbox(
            slide, x, Inches(5), Inches(3.9), Inches(0.8),
            fact, font_size=12, font_color=COLORS['sky_blue'],
            alignment=PP_ALIGN.CENTER
        )

    # Bottom insight
    add_styled_textbox(
        slide, Inches(1), Inches(6.3), Inches(11.33), Inches(0.8),
        "Without cloud data centers, AI like ChatGPT simply couldn't exist!",
        font_size=20, font_color=COLORS['medium_blue'], bold=True,
        alignment=PP_ALIGN.CENTER
    )

def create_ai_chatbots_slide(prs):
    """Slide 16: Major AI chatbots and their clouds."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Meet the AI Chatbots!", "The most powerful AIs in the world")

    chatbots = [
        ("ChatGPT", "OpenAI", "Microsoft Azure", "Most popular AI chatbot\n175B+ parameters", COLORS['teal']),
        ("Grok", "xAI (Elon Musk)", "X.ai Cloud", "Built into X (Twitter)\nReal-time information", COLORS['dark_gray']),
        ("Claude", "Anthropic", "AWS + Google", "Known for being safe\nand helpful", COLORS['orange']),
        ("Gemini", "Google", "Google Cloud", "Multimodal AI\n(text, images, code)", COLORS['google_blue']),
    ]

    for i, (name, company, cloud, desc, color) in enumerate(chatbots):
        col = i % 2
        row = i // 2
        x = Inches(0.7 + col * 6.3)
        y = Inches(1.7 + row * 2.7)

        # Card
        card = add_rounded_box(
            slide, x, y, Inches(5.9), Inches(2.4),
            COLORS['light_gray'], border_color=color
        )

        # Header with name
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x, y,
            Inches(5.9), Inches(0.65)
        )
        set_shape_fill(header, color)
        header.line.fill.background()

        add_styled_textbox(
            slide, x, y + Inches(0.12), Inches(5.9), Inches(0.5),
            name, font_size=22, font_color=COLORS['white'],
            alignment=PP_ALIGN.CENTER, bold=True
        )

        # Company
        add_styled_textbox(
            slide, x, y + Inches(0.75), Inches(5.9), Inches(0.4),
            f"by {company}", font_size=12, font_color=COLORS['dark_gray'],
            alignment=PP_ALIGN.CENTER
        )

        # Cloud provider
        add_styled_textbox(
            slide, x, y + Inches(1.1), Inches(5.9), Inches(0.4),
            f"Runs on: {cloud}", font_size=13, font_color=color,
            alignment=PP_ALIGN.CENTER, bold=True
        )

        # Description
        add_styled_textbox(
            slide, x, y + Inches(1.5), Inches(5.9), Inches(0.8),
            desc, font_size=11, font_color=COLORS['dark_gray'],
            alignment=PP_ALIGN.CENTER
        )

    # Key message
    add_styled_textbox(
        slide, Inches(1), Inches(6.8), Inches(11.33), Inches(0.5),
        "All these AIs need THOUSANDS of cloud servers to answer a single question!",
        font_size=18, font_color=COLORS['purple'], bold=True,
        alignment=PP_ALIGN.CENTER
    )

def create_ai_services_slide(prs):
    """Slide 17: AI services powered by the cloud."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "AI in Apps You Use Daily!", "AI is everywhere...")

    services = [
        ("Google\nSearch", "AI understands\nyour questions", COLORS['google_blue']),
        ("YouTube", "AI recommends\nvideos for you", COLORS['red']),
        ("Instagram", "AI filters &\nrecommendations", COLORS['purple']),
        ("Snapchat", "AI face filters\nin real-time", COLORS['orange']),
        ("Spotify", "AI creates your\npersonalized playlists", COLORS['green']),
        ("TikTok", "AI picks videos\nfor your For You page", COLORS['dark_gray']),
        ("Grammarly", "AI fixes your\nwriting mistakes", COLORS['green']),
        ("Google\nTranslate", "AI translates\n100+ languages", COLORS['medium_blue']),
    ]

    for i, (name, desc, color) in enumerate(services):
        col = i % 4
        row = i // 4
        x = Inches(0.5 + col * 3.15)
        y = Inches(1.7 + row * 2.5)

        # Service card
        card = add_rounded_box(
            slide, x, y, Inches(2.95), Inches(2.2),
            COLORS['light_gray'], border_color=color
        )

        # Header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x, y,
            Inches(2.95), Inches(0.7)
        )
        set_shape_fill(header, color)
        header.line.fill.background()

        add_styled_textbox(
            slide, x, y + Inches(0.08), Inches(2.95), Inches(0.65),
            name, font_size=13, font_color=COLORS['white'],
            alignment=PP_ALIGN.CENTER, bold=True
        )

        add_styled_textbox(
            slide, x, y + Inches(0.9), Inches(2.95), Inches(1),
            desc, font_size=11, font_color=COLORS['dark_gray'],
            alignment=PP_ALIGN.CENTER
        )

    # Key insight
    add_styled_textbox(
        slide, Inches(1), Inches(6.7), Inches(11.33), Inches(0.5),
        "All these AI features run on cloud servers - not on your device!",
        font_size=18, font_color=COLORS['purple'], bold=True,
        alignment=PP_ALIGN.CENTER
    )

def create_ai_future_slide(prs):
    """Slide 17: The future of AI and cloud."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "The Future: AI + Cloud", "What's coming next?")

    # Current vs Future comparison
    # Current
    add_rounded_box(
        slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(2.5),
        COLORS['medium_blue']
    )
    add_styled_textbox(
        slide, Inches(0.8), Inches(1.85), Inches(5.8), Inches(0.5),
        "TODAY (2026)",
        font_size=22, font_color=COLORS['white'],
        alignment=PP_ALIGN.CENTER, bold=True
    )
    today_items = [
        "ChatGPT answers questions",
        "AI generates images & videos",
        "Voice assistants help you",
        "AI recommends content"
    ]
    for i, item in enumerate(today_items):
        add_styled_textbox(
            slide, Inches(1.1), Inches(2.4 + i * 0.45), Inches(5.2), Inches(0.45),
            f"  {item}", font_size=14, font_color=COLORS['white']
        )

    # Future
    add_rounded_box(
        slide, Inches(6.9), Inches(1.7), Inches(5.8), Inches(2.5),
        COLORS['purple']
    )
    add_styled_textbox(
        slide, Inches(6.9), Inches(1.85), Inches(5.8), Inches(0.5),
        "THE FUTURE",
        font_size=22, font_color=COLORS['white'],
        alignment=PP_ALIGN.CENTER, bold=True
    )
    future_items = [
        "AI creates movies & games",
        "AI helps doctors & scientists",
        "AI assistants for everything",
        "Self-driving cars everywhere"
    ]
    for i, item in enumerate(future_items):
        add_styled_textbox(
            slide, Inches(7.2), Inches(2.4 + i * 0.45), Inches(5.2), Inches(0.45),
            f"  {item}", font_size=14, font_color=COLORS['white']
        )

    # Key message
    add_rounded_box(
        slide, Inches(1.5), Inches(4.5), Inches(10.33), Inches(1.1),
        COLORS['orange'],
        "ALL of this is only possible because of CLOUD COMPUTING!\nMore powerful AI = More powerful data centers needed",
        text_size=18, text_color=COLORS['white']
    )

    # Why this matters to you
    add_styled_textbox(
        slide, Inches(0.8), Inches(5.85), Inches(11.73), Inches(0.5),
        "Why This Matters to YOU:",
        font_size=18, font_color=COLORS['dark_blue'], bold=True,
        alignment=PP_ALIGN.CENTER
    )

    careers = ["Cloud Engineer", "AI Developer", "Data Scientist", "ML Engineer"]
    for i, career in enumerate(careers):
        x = Inches(1 + i * 3)
        add_rounded_box(
            slide, x, Inches(6.35), Inches(2.7), Inches(0.65),
            COLORS['green'],
            career,
            text_size=13, text_color=COLORS['white']
        )

def create_recap_slide(prs):
    """Slide 19: Recap slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Today's Key Takeaways")

    takeaways = [
        ("THE CLOUD", "Servers in data centers (not in the sky!)", COLORS['cloud_blue']),
        ("DATA CENTERS", "Huge buildings with thousands of servers", COLORS['data_center_gray']),
        ("BIG 3 PROVIDERS", "AWS, Microsoft Azure, Google Cloud", COLORS['orange']),
        ("AI + CLOUD", "ChatGPT, Grok, Claude all need cloud power!", COLORS['purple']),
        ("CLOUD SAFETY", "Strong passwords & think before you share", COLORS['green']),
    ]

    for i, (term, definition, color) in enumerate(takeaways):
        y = Inches(1.65 + i * 1.0)

        # Term box
        term_box = add_rounded_box(
            slide, Inches(0.8), y, Inches(2.8), Inches(0.85),
            color,
            term, text_size=17, text_color=COLORS['white']
        )

        # Definition
        add_styled_textbox(
            slide, Inches(4), y + Inches(0.2), Inches(8.5), Inches(0.6),
            definition, font_size=18, font_color=COLORS['dark_gray']
        )

    # Kahoot time!
    add_styled_textbox(
        slide, Inches(1), Inches(6.8), Inches(11.33), Inches(0.5),
        "Great job today! Time for Kahoot!",
        font_size=24, font_color=COLORS['dark_blue'], bold=True,
        alignment=PP_ALIGN.CENTER
    )

def create_questions_slide(prs):
    """Slide 16: Questions slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    set_shape_fill(background, COLORS['dark_blue'])
    background.line.fill.background()

    # Decorative clouds
    add_cloud_shape(slide, Inches(0.5), Inches(0.5), Inches(3), Inches(1.8), COLORS['medium_blue'])
    add_cloud_shape(slide, Inches(9.5), Inches(5.5), Inches(3.5), Inches(2), COLORS['medium_blue'])

    # Question mark
    add_styled_textbox(
        slide, Inches(0), Inches(1.5), Inches(13.33), Inches(2),
        "?", font_size=120, font_color=COLORS['white'],
        alignment=PP_ALIGN.CENTER
    )

    # Title
    add_styled_textbox(
        slide, Inches(0), Inches(3.5), Inches(13.33), Inches(1),
        "Questions?",
        font_size=60, font_color=COLORS['white'], bold=True,
        alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold"
    )

    add_styled_textbox(
        slide, Inches(0), Inches(5), Inches(13.33), Inches(0.6),
        "Ask away before we start the Kahoot!",
        font_size=24, font_color=COLORS['sky_blue'],
        alignment=PP_ALIGN.CENTER
    )

# ============================================================================
# MAIN FUNCTION
# ============================================================================

def main():
    """Create the presentation."""
    # Create presentation with widescreen dimensions (16:9)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("Creating slides...")

    # Create all slides
    create_title_slide(prs)
    print("  [OK] Title slide")

    create_agenda_slide(prs)
    print("  [OK] Agenda slide")

    create_review_slide(prs)
    print("  [OK] Review slide")

    create_cloud_myth_slide(prs)
    print("  [OK] Cloud myth slide")

    create_data_center_intro_slide(prs)
    print("  [OK] Data center intro slide")

    create_data_center_scale_slide(prs)
    print("  [OK] Data center scale slide")

    create_cloud_services_intro_slide(prs)
    print("  [OK] Cloud services intro slide")

    create_cloud_you_use_slide(prs)
    print("  [OK] Cloud you use slide")

    create_no_cloud_slide(prs)
    print("  [OK] No cloud slide")

    create_big_three_slide(prs)
    print("  [OK] Big three slide")

    create_how_cloud_works_slide(prs)
    print("  [OK] How cloud works slide")

    create_local_vs_cloud_slide(prs)
    print("  [OK] Local vs cloud slide")

    create_cloud_safety_slide(prs)
    print("  [OK] Cloud safety slide")

    create_connection_slide(prs)
    print("  [OK] Connection slide")

    create_ai_needs_cloud_slide(prs)
    print("  [OK] AI needs cloud slide")

    create_ai_chatbots_slide(prs)
    print("  [OK] AI chatbots slide")

    create_ai_services_slide(prs)
    print("  [OK] AI services slide")

    create_ai_future_slide(prs)
    print("  [OK] AI future slide")

    create_recap_slide(prs)
    print("  [OK] Recap slide")

    create_questions_slide(prs)
    print("  [OK] Questions slide")

    # Save the presentation
    script_dir = os.path.dirname(os.path.abspath(__file__))
    parent_dir = os.path.dirname(script_dir)
    output_path = os.path.join(parent_dir, "2026-01-31_Cloud_Data_Centers_BEAUTIFUL.pptx")

    prs.save(output_path)
    print(f"\n[SUCCESS] Presentation saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
