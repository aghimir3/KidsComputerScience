"""
Add Real Data Center Images to Presentation
Downloads official Google data center photos and adds them to the presentation.

Run: python add_datacenter_images.py
"""

import os
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# URLs for official Google data center images
IMAGES = {
    'google_servers': {
        'url': 'https://www.gstatic.com/marketing-cms/assets/images/e0/7a/021a08cb491c8b43d190d663a2db/douglas-county-servers.jpg',
        'caption': 'Google Data Center - Server Rows (Douglas County, Georgia)',
        'credit': 'Photo: Google'
    },
    'google_pipes': {
        'url': 'https://www.gstatic.com/marketing-cms/assets/images/0b/76/ab8d87fe45078fe8a97144b02055/douglas-county-pipes.jpg',
        'caption': 'Colorful Cooling Pipes (Douglas County, Georgia)',
        'credit': 'Photo: Google'
    },
    'google_server_aisle': {
        'url': 'https://www.gstatic.com/marketing-cms/assets/images/19/43/b476c0984f2da3b2faa1a7f588ce/server-aisles-in-our-new-albany-data-center-building-in-central-ohio.jpg',
        'caption': 'Server Aisles (New Albany, Ohio)',
        'credit': 'Photo: Google'
    },
    'google_network': {
        'url': 'https://www.gstatic.com/marketing-cms/assets/images/05/67/4a986e4e49d6831b1e15fb67f1dc/council-bluffs-network-room.jpg',
        'caption': 'Network Room - Routers & Switches (Council Bluffs, Iowa)',
        'credit': 'Photo: Google'
    },
    'google_cooling': {
        'url': 'https://www.gstatic.com/marketing-cms/assets/images/6f/e1/50ab9c9647728ae4b81e36b1a730/the-dalles-cooling-towers.jpg',
        'caption': 'Cooling Towers (The Dalles, Oregon)',
        'credit': 'Photo: Google'
    },
}

# Colors
COLORS = {
    'dark_blue': RGBColor(30, 58, 95),
    'orange': RGBColor(255, 153, 51),
    'white': RGBColor(255, 255, 255),
    'light_gray': RGBColor(200, 200, 200),
    'dark_gray': RGBColor(66, 66, 66),
}


def download_image(url, save_path):
    """Download an image from URL."""
    try:
        response = requests.get(url, timeout=30)
        response.raise_for_status()
        with open(save_path, 'wb') as f:
            f.write(response.content)
        print(f"  Downloaded: {os.path.basename(save_path)}")
        return True
    except Exception as e:
        print(f"  Failed to download {url}: {e}")
        return False


def add_title_shape(slide, text, subtitle=None):
    """Add a styled title bar at the top of the slide."""
    # Title background bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(1.3)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLORS['dark_blue']
    title_bar.line.fill.background()

    # Title text
    textbox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(12.33), Inches(0.7)
    )
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.color.rgb = COLORS['white']
    p.font.bold = True
    p.font.name = "Segoe UI Semibold"

    # Accent line
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.3),
        Inches(13.33), Inches(0.08)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = COLORS['orange']
    accent_line.line.fill.background()

    if subtitle:
        sub_textbox = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.95), Inches(12.33), Inches(0.4)
        )
        tf = sub_textbox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(173, 216, 230)  # sky_blue
        p.font.name = "Segoe UI"


def create_photo_slide(prs, images_dir):
    """Create a slide with real data center photos."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Real Data Center Photos!", "Inside Google's Data Centers")

    # Photo positions (2x2 grid + 1 bottom)
    positions = [
        (0.5, 1.7, 4.1, 2.5),    # Top left
        (4.8, 1.7, 4.1, 2.5),    # Top middle
        (9.1, 1.7, 3.8, 2.5),    # Top right
        (2.5, 4.4, 4.1, 2.5),    # Bottom left
        (6.8, 4.4, 4.1, 2.5),    # Bottom right
    ]

    image_keys = ['google_server_aisle', 'google_servers', 'google_network', 'google_pipes', 'google_cooling']

    for i, key in enumerate(image_keys):
        if i >= len(positions):
            break

        img_path = os.path.join(images_dir, f"{key}.jpg")
        if not os.path.exists(img_path):
            continue

        left, top, width, height = positions[i]
        img_info = IMAGES[key]

        try:
            # Add image
            pic = slide.shapes.add_picture(
                img_path,
                Inches(left), Inches(top),
                Inches(width), Inches(height)
            )

            # Add caption below image
            caption_box = slide.shapes.add_textbox(
                Inches(left), Inches(top + height + 0.05),
                Inches(width), Inches(0.3)
            )
            tf = caption_box.text_frame
            p = tf.paragraphs[0]
            p.text = img_info['credit']
            p.font.size = Pt(9)
            p.font.color.rgb = COLORS['dark_gray']
            p.font.name = "Segoe UI"
            p.alignment = PP_ALIGN.CENTER

        except Exception as e:
            print(f"  Error adding image {key}: {e}")

    return slide


def main():
    """Main function to add data center images to presentation."""
    script_dir = os.path.dirname(os.path.abspath(__file__))
    parent_dir = os.path.dirname(script_dir)

    # Create images directory
    images_dir = os.path.join(parent_dir, "images")
    os.makedirs(images_dir, exist_ok=True)

    # Download images
    print("Downloading data center images...")
    for key, info in IMAGES.items():
        save_path = os.path.join(images_dir, f"{key}.jpg")
        if not os.path.exists(save_path):
            download_image(info['url'], save_path)
        else:
            print(f"  Already exists: {key}.jpg")

    # Open existing presentation
    pptx_path = os.path.join(parent_dir, "2026-01-31_Cloud_Data_Centers_BEAUTIFUL.pptx")
    if not os.path.exists(pptx_path):
        print(f"ERROR: Presentation not found at {pptx_path}")
        return

    print(f"\nOpening presentation: {pptx_path}")
    prs = Presentation(pptx_path)

    # Find the slide after "What is a Data Center?" (slide 5) to insert photos
    # We'll add a new slide after the data center intro
    print("Adding photo slide...")

    # Create the photo slide
    photo_slide = create_photo_slide(prs, images_dir)

    # Move the new slide to position 6 (after data center intro)
    # The slide is added at the end, we need to reorder
    slide_count = len(prs.slides)
    print(f"  Total slides: {slide_count}")

    # Save to a new file
    output_path = os.path.join(parent_dir, "2026-01-31_Cloud_Data_Centers_BEAUTIFUL.pptx")
    prs.save(output_path)
    print(f"\n[SUCCESS] Presentation updated: {output_path}")
    print("  Note: Photo slide added at the end - you can move it in PowerPoint")


if __name__ == "__main__":
    main()
