"""
PowerPoint Generator for: Clients, Servers & How Websites Work
Kids Computer Science Class - January 25, 2026

This script creates a beautifully designed presentation using python-pptx.
Run: pip install python-pptx
Then: python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# ============================================================================
# COLOR PALETTE - Modern, Professional, Kid-Friendly
# ============================================================================
COLORS = {
    # Primary colors
    'dark_blue': RGBColor(30, 58, 95),      # Headers, titles
    'medium_blue': RGBColor(52, 103, 166),  # Accents
    'light_blue': RGBColor(100, 149, 237),  # Highlights
    'sky_blue': RGBColor(173, 216, 230),    # Backgrounds

    # Accent colors
    'orange': RGBColor(255, 153, 51),       # Call-outs, emphasis
    'green': RGBColor(76, 175, 80),         # Success, positive
    'red': RGBColor(244, 67, 54),           # Warnings, important
    'purple': RGBColor(156, 39, 176),       # Special highlights

    # Neutrals
    'white': RGBColor(255, 255, 255),
    'light_gray': RGBColor(245, 245, 245),
    'dark_gray': RGBColor(66, 66, 66),
    'black': RGBColor(33, 33, 33),

    # Gradients (we'll use these for visual distinction)
    'client_color': RGBColor(100, 181, 246),   # Light blue for clients
    'server_color': RGBColor(129, 199, 132),   # Light green for servers
}

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def set_shape_fill(shape, color):
    """Set solid fill color for a shape."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_shape_line(shape, color, width_pt=1):
    """Set line/border color and width for a shape."""
    line = shape.line
    line.color.rgb = color
    line.width = Pt(width_pt)

def add_styled_textbox(slide, left, top, width, height, text,
                       font_size=18, font_color=None, bold=False,
                       alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Add a textbox with styled text."""
    if font_color is None:
        font_color = COLORS['dark_gray']

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment

    return textbox

def add_title_shape(slide, text, subtitle=None):
    """Add a styled title bar at the top of the slide."""
    # Title background bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(1.3)
    )
    set_shape_fill(title_bar, COLORS['dark_blue'])
    title_bar.line.fill.background()

    # Title text
    add_styled_textbox(
        slide, Inches(0.5), Inches(0.35), Inches(12.33), Inches(0.7),
        text, font_size=36, font_color=COLORS['white'], bold=True,
        alignment=PP_ALIGN.LEFT, font_name="Segoe UI Semibold"
    )

    # Accent line under title
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.3),
        Inches(13.33), Inches(0.08)
    )
    set_shape_fill(accent_line, COLORS['orange'])
    accent_line.line.fill.background()

    if subtitle:
        add_styled_textbox(
            slide, Inches(0.5), Inches(0.95), Inches(12.33), Inches(0.4),
            subtitle, font_size=16, font_color=COLORS['sky_blue'],
            alignment=PP_ALIGN.LEFT
        )

def add_icon_shape(slide, shape_type, left, top, size, fill_color, text=None, text_color=None):
    """Add a shape that acts as an icon with optional text inside."""
    shape = slide.shapes.add_shape(shape_type, left, top, size, size)
    set_shape_fill(shape, fill_color)
    shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(24)
        p.font.color.rgb = text_color or COLORS['white']
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.name = "Segoe UI"

    return shape

def add_bullet_points(slide, left, top, width, height, items, font_size=20):
    """Add a text frame with bullet points."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS['dark_gray']
        p.font.name = "Segoe UI"
        p.space_before = Pt(8)
        p.space_after = Pt(8)

    return textbox

def add_rounded_box(slide, left, top, width, height, fill_color, text=None,
                    text_size=16, text_color=None, border_color=None):
    """Add a rounded rectangle with optional text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    set_shape_fill(shape, fill_color)

    if border_color:
        set_shape_line(shape, border_color, 2)
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(text_size)
        p.font.color.rgb = text_color or COLORS['white']
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(10)

    # Center text vertically
    shape.text_frame.anchor = MSO_ANCHOR.MIDDLE

    return shape

def add_arrow(slide, left, top, width, height, fill_color, direction='right'):
    """Add an arrow shape."""
    if direction == 'right':
        shape_type = MSO_SHAPE.RIGHT_ARROW
    elif direction == 'left':
        shape_type = MSO_SHAPE.LEFT_ARROW
    elif direction == 'down':
        shape_type = MSO_SHAPE.DOWN_ARROW
    else:
        shape_type = MSO_SHAPE.UP_ARROW

    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    set_shape_fill(shape, fill_color)
    shape.line.fill.background()
    return shape

# ============================================================================
# SLIDE CREATION FUNCTIONS
# ============================================================================

def create_title_slide(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    set_shape_fill(background, COLORS['dark_blue'])
    background.line.fill.background()

    # Decorative circles
    for i, (x, y, size, color) in enumerate([
        (10, 0.5, 2, COLORS['medium_blue']),
        (11, 2, 1.5, COLORS['light_blue']),
        (0.5, 5, 1.8, COLORS['medium_blue']),
        (-0.5, 6, 1.2, COLORS['light_blue']),
    ]):
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y),
            Inches(size), Inches(size)
        )
        set_shape_fill(circle, color)
        circle.line.fill.background()
        # Make semi-transparent by using lighter shade

    # Main title
    add_styled_textbox(
        slide, Inches(0.5), Inches(2.2), Inches(12.33), Inches(1.2),
        "Clients, Servers &\nHow Websites Work",
        font_size=52, font_color=COLORS['white'], bold=True,
        alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold"
    )

    # Subtitle
    add_styled_textbox(
        slide, Inches(0.5), Inches(4.2), Inches(12.33), Inches(0.6),
        "Understanding the Internet",
        font_size=28, font_color=COLORS['sky_blue'],
        alignment=PP_ALIGN.CENTER
    )

    # Date
    add_styled_textbox(
        slide, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.5),
        "January 25, 2026",
        font_size=18, font_color=COLORS['light_blue'],
        alignment=PP_ALIGN.CENTER
    )

def create_agenda_slide(prs):
    """Slide 2: Today's agenda."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Today's Journey")

    topics = [
        ("1", "What are Clients & Servers?", COLORS['client_color']),
        ("2", "HTTP & HTTPS - The Web's Language", COLORS['orange']),
        ("3", "Ports - The Apartment Numbers", COLORS['green']),
        ("4", "How a Website Actually Loads", COLORS['purple']),
        ("5", "Hands-on: tracert Command", COLORS['medium_blue']),
    ]

    for i, (num, text, color) in enumerate(topics):
        y = Inches(1.8 + i * 1.0)

        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1), y,
            Inches(0.7), Inches(0.7)
        )
        set_shape_fill(circle, color)
        circle.line.fill.background()

        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['white']
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Topic text
        add_styled_textbox(
            slide, Inches(2), y + Inches(0.15), Inches(9), Inches(0.5),
            text, font_size=24, font_color=COLORS['dark_gray']
        )

def create_client_intro_slide(prs):
    """Slide 3: What is a Client?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "What is a CLIENT?")

    # Big definition box
    add_rounded_box(
        slide, Inches(1), Inches(1.8), Inches(11.33), Inches(1.2),
        COLORS['client_color'],
        "A device that ASKS for information",
        text_size=32, text_color=COLORS['dark_blue']
    )

    # Examples header
    add_styled_textbox(
        slide, Inches(1), Inches(3.3), Inches(4), Inches(0.5),
        "Examples of Clients:", font_size=22, font_color=COLORS['dark_blue'], bold=True
    )

    # Client examples with icons
    examples = [
        ("💻", "Your Laptop"),
        ("📱", "Your Phone"),
        ("🎮", "Gaming Console"),
        ("📺", "Smart TV"),
    ]

    for i, (icon, label) in enumerate(examples):
        x = Inches(1.5 + i * 2.8)
        y = Inches(4)

        box = add_rounded_box(
            slide, x, y, Inches(2.4), Inches(1.8),
            COLORS['light_gray'], border_color=COLORS['client_color']
        )

        add_styled_textbox(
            slide, x, y + Inches(0.3), Inches(2.4), Inches(0.8),
            icon, font_size=40, alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x, y + Inches(1.1), Inches(2.4), Inches(0.5),
            label, font_size=16, font_color=COLORS['dark_gray'],
            alignment=PP_ALIGN.CENTER, bold=True
        )

    # Key point
    add_styled_textbox(
        slide, Inches(1), Inches(6.3), Inches(11.33), Inches(0.6),
        "🔑 Key: When YOU browse the web, YOUR device is the client!",
        font_size=20, font_color=COLORS['medium_blue'], bold=True,
        alignment=PP_ALIGN.CENTER
    )

def create_server_intro_slide(prs):
    """Slide 4: What is a Server?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "What is a SERVER?")

    # Big definition box
    add_rounded_box(
        slide, Inches(1), Inches(1.8), Inches(11.33), Inches(1.2),
        COLORS['server_color'],
        "A computer that SERVES information to clients",
        text_size=32, text_color=COLORS['dark_blue']
    )

    # Server characteristics
    add_styled_textbox(
        slide, Inches(1), Inches(3.3), Inches(5), Inches(0.5),
        "What makes a server special?", font_size=22, font_color=COLORS['dark_blue'], bold=True
    )

    characteristics = [
        "Always ON (24/7)",
        "Connected to the internet",
        "Waiting for requests",
        "Stores websites, apps, data"
    ]

    add_bullet_points(
        slide, Inches(1), Inches(3.9), Inches(5.5), Inches(2.5),
        characteristics, font_size=20
    )

    # Visual: Server icon
    server_box = add_rounded_box(
        slide, Inches(7.5), Inches(3.5), Inches(4.5), Inches(2.8),
        COLORS['dark_blue']
    )

    add_styled_textbox(
        slide, Inches(7.5), Inches(4), Inches(4.5), Inches(1),
        "🖥️", font_size=60, alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(7.5), Inches(5.2), Inches(4.5), Inches(0.6),
        "SERVER", font_size=24, font_color=COLORS['white'],
        alignment=PP_ALIGN.CENTER, bold=True
    )

    # Key point
    add_styled_textbox(
        slide, Inches(1), Inches(6.5), Inches(11.33), Inches(0.6),
        "🔑 Google, Netflix, YouTube all run on powerful servers!",
        font_size=20, font_color=COLORS['green'], bold=True,
        alignment=PP_ALIGN.CENTER
    )

def create_client_server_diagram_slide(prs):
    """Slide 5: Client-Server interaction diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "How Clients & Servers Talk")

    # Client box (left)
    client_box = add_rounded_box(
        slide, Inches(0.8), Inches(3), Inches(3), Inches(2.2),
        COLORS['client_color']
    )
    add_styled_textbox(
        slide, Inches(0.8), Inches(3.3), Inches(3), Inches(0.8),
        "💻", font_size=50, alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(0.8), Inches(4.3), Inches(3), Inches(0.6),
        "CLIENT", font_size=22, font_color=COLORS['dark_blue'],
        alignment=PP_ALIGN.CENTER, bold=True
    )
    add_styled_textbox(
        slide, Inches(0.8), Inches(4.8), Inches(3), Inches(0.4),
        "(Your Computer)", font_size=14, font_color=COLORS['dark_blue'],
        alignment=PP_ALIGN.CENTER
    )

    # Server box (right)
    server_box = add_rounded_box(
        slide, Inches(9.5), Inches(3), Inches(3), Inches(2.2),
        COLORS['server_color']
    )
    add_styled_textbox(
        slide, Inches(9.5), Inches(3.3), Inches(3), Inches(0.8),
        "🖥️", font_size=50, alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(9.5), Inches(4.3), Inches(3), Inches(0.6),
        "SERVER", font_size=22, font_color=COLORS['dark_blue'],
        alignment=PP_ALIGN.CENTER, bold=True
    )
    add_styled_textbox(
        slide, Inches(9.5), Inches(4.8), Inches(3), Inches(0.4),
        "(Google's Computer)", font_size=14, font_color=COLORS['dark_blue'],
        alignment=PP_ALIGN.CENTER
    )

    # Request arrow (top)
    request_arrow = add_arrow(
        slide, Inches(4.2), Inches(3.2), Inches(5), Inches(0.6),
        COLORS['orange'], 'right'
    )
    add_styled_textbox(
        slide, Inches(4.5), Inches(2.5), Inches(4.4), Inches(0.6),
        "1. REQUEST", font_size=18, font_color=COLORS['orange'],
        alignment=PP_ALIGN.CENTER, bold=True
    )
    add_styled_textbox(
        slide, Inches(4.5), Inches(2.9), Inches(4.4), Inches(0.4),
        '"Hey, can I see google.com?"', font_size=14, font_color=COLORS['dark_gray'],
        alignment=PP_ALIGN.CENTER
    )

    # Response arrow (bottom)
    response_arrow = add_arrow(
        slide, Inches(4.2), Inches(4.7), Inches(5), Inches(0.6),
        COLORS['green'], 'left'
    )
    add_styled_textbox(
        slide, Inches(4.5), Inches(5.4), Inches(4.4), Inches(0.6),
        "2. RESPONSE", font_size=18, font_color=COLORS['green'],
        alignment=PP_ALIGN.CENTER, bold=True
    )
    add_styled_textbox(
        slide, Inches(4.5), Inches(5.8), Inches(4.4), Inches(0.4),
        '"Here\'s the webpage!"', font_size=14, font_color=COLORS['dark_gray'],
        alignment=PP_ALIGN.CENTER
    )

    # Key insight
    add_styled_textbox(
        slide, Inches(1), Inches(6.5), Inches(11.33), Inches(0.6),
        "This request-response pattern happens EVERY time you visit a website!",
        font_size=18, font_color=COLORS['medium_blue'], bold=True,
        alignment=PP_ALIGN.CENTER
    )

def create_http_intro_slide(prs):
    """Slide 6: What is HTTP?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "HTTP - The Web's Language")

    # Acronym breakdown
    letters = [
        ("H", "Hyper", COLORS['orange']),
        ("T", "Text", COLORS['green']),
        ("T", "Transfer", COLORS['purple']),
        ("P", "Protocol", COLORS['medium_blue']),
    ]

    for i, (letter, word, color) in enumerate(letters):
        x = Inches(1.5 + i * 2.8)

        # Letter circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x, Inches(2),
            Inches(1), Inches(1)
        )
        set_shape_fill(circle, color)
        circle.line.fill.background()

        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = letter
        p.font.size = Pt(40)
        p.font.color.rgb = COLORS['white']
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Word below
        add_styled_textbox(
            slide, x - Inches(0.3), Inches(3.1), Inches(1.6), Inches(0.5),
            word, font_size=18, font_color=color, bold=True,
            alignment=PP_ALIGN.CENTER
        )

    # Explanation box
    add_rounded_box(
        slide, Inches(1), Inches(4), Inches(11.33), Inches(1.4),
        COLORS['light_gray'], border_color=COLORS['medium_blue']
    )
    add_styled_textbox(
        slide, Inches(1.3), Inches(4.2), Inches(10.73), Inches(1.2),
        "HTTP is the set of rules that clients and servers use to communicate.\n"
        "It's like the language they both speak!",
        font_size=22, font_color=COLORS['dark_gray'],
        alignment=PP_ALIGN.CENTER
    )

    # Visual analogy
    add_styled_textbox(
        slide, Inches(1), Inches(5.8), Inches(11.33), Inches(1),
        "🗣️ Think of it like: When you order food, you speak the same language as the waiter.\n"
        "HTTP is the \"language\" computers use to order and deliver web pages!",
        font_size=18, font_color=COLORS['dark_blue'],
        alignment=PP_ALIGN.CENTER
    )

def create_https_slide(prs):
    """Slide 7: HTTP vs HTTPS."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "HTTP vs HTTPS", "What's the 'S' for?")

    # HTTP box (left) - Not secure
    http_box = add_rounded_box(
        slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5),
        COLORS['light_gray'], border_color=COLORS['red']
    )
    add_styled_textbox(
        slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.6),
        "HTTP", font_size=32, font_color=COLORS['red'],
        alignment=PP_ALIGN.CENTER, bold=True
    )
    add_styled_textbox(
        slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(0.8),
        "📬", font_size=50, alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(1), Inches(3.4), Inches(5.1), Inches(0.7),
        "Like a POSTCARD\nAnyone can read it!",
        font_size=16, font_color=COLORS['dark_gray'],
        alignment=PP_ALIGN.CENTER
    )

    # HTTPS box (right) - Secure
    https_box = add_rounded_box(
        slide, Inches(7), Inches(1.8), Inches(5.5), Inches(2.5),
        COLORS['light_gray'], border_color=COLORS['green']
    )
    add_styled_textbox(
        slide, Inches(7), Inches(1.9), Inches(5.5), Inches(0.6),
        "HTTPS", font_size=32, font_color=COLORS['green'],
        alignment=PP_ALIGN.CENTER, bold=True
    )
    add_styled_textbox(
        slide, Inches(7), Inches(2.5), Inches(5.5), Inches(0.8),
        "🔒", font_size=50, alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(7.2), Inches(3.4), Inches(5.1), Inches(0.7),
        "Like a LOCKED BOX\nOnly you & server can read!",
        font_size=16, font_color=COLORS['dark_gray'],
        alignment=PP_ALIGN.CENTER
    )

    # The S stands for...
    add_rounded_box(
        slide, Inches(3.5), Inches(4.6), Inches(6.33), Inches(1.2),
        COLORS['green'],
        "S = SECURE",
        text_size=36, text_color=COLORS['white']
    )

    # Why it matters
    add_styled_textbox(
        slide, Inches(1), Inches(6.2), Inches(11.33), Inches(1),
        "🔑 Always look for HTTPS and the 🔒 in your browser!\n"
        "Especially when entering passwords or credit cards.",
        font_size=18, font_color=COLORS['dark_blue'],
        alignment=PP_ALIGN.CENTER, bold=True
    )

def create_ports_intro_slide(prs):
    """Slide 8: What are Ports?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "What are PORTS?")

    # Analogy: Building
    add_styled_textbox(
        slide, Inches(0.8), Inches(1.7), Inches(11.73), Inches(0.6),
        "🏢 Think of a server like an APARTMENT BUILDING...",
        font_size=24, font_color=COLORS['dark_blue'], bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # Building visual
    building = add_rounded_box(
        slide, Inches(1.5), Inches(2.5), Inches(4), Inches(3.5),
        COLORS['dark_blue']
    )

    # Apartment numbers (ports)
    apartments = [
        ("80", "Web (HTTP)", Inches(1.8), Inches(2.8)),
        ("443", "Secure Web", Inches(1.8), Inches(3.8)),
        ("22", "SSH", Inches(3.4), Inches(2.8)),
        ("25", "Email", Inches(3.4), Inches(3.8)),
    ]

    for port, service, x, y in apartments:
        box = add_rounded_box(
            slide, x, y, Inches(1.4), Inches(0.8),
            COLORS['orange']
        )
        add_styled_textbox(
            slide, x, y + Inches(0.1), Inches(1.4), Inches(0.4),
            f"#{port}", font_size=16, font_color=COLORS['white'],
            alignment=PP_ALIGN.CENTER, bold=True
        )
        add_styled_textbox(
            slide, x, y + Inches(0.45), Inches(1.4), Inches(0.3),
            service, font_size=10, font_color=COLORS['white'],
            alignment=PP_ALIGN.CENTER
        )

    add_styled_textbox(
        slide, Inches(1.5), Inches(6.1), Inches(4), Inches(0.4),
        "SERVER", font_size=18, font_color=COLORS['dark_blue'],
        alignment=PP_ALIGN.CENTER, bold=True
    )

    # Explanation
    explanation = [
        "IP Address = Building Address",
        "Port = Apartment Number",
        "Different services live in different \"apartments\""
    ]

    add_styled_textbox(
        slide, Inches(6), Inches(2.7), Inches(6.5), Inches(0.5),
        "The Analogy:", font_size=22, font_color=COLORS['dark_blue'], bold=True
    )

    add_bullet_points(
        slide, Inches(6), Inches(3.3), Inches(6.5), Inches(2),
        explanation, font_size=20
    )

def create_ports_numbers_slide(prs):
    """Slide 9: Common Port Numbers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Common Port Numbers", "The ones you should know!")

    # Port cards
    ports = [
        ("80", "HTTP", "Regular websites", COLORS['orange']),
        ("443", "HTTPS", "Secure websites", COLORS['green']),
        ("22", "SSH", "Remote computer access", COLORS['purple']),
        ("53", "DNS", "Domain lookups", COLORS['medium_blue']),
    ]

    for i, (port, name, desc, color) in enumerate(ports):
        x = Inches(0.8 + (i % 2) * 6.3)
        y = Inches(1.8 + (i // 2) * 2.5)

        # Port card
        card = add_rounded_box(
            slide, x, y, Inches(5.8), Inches(2.2),
            COLORS['white'], border_color=color
        )

        # Port number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + Inches(0.3), y + Inches(0.5),
            Inches(1.2), Inches(1.2)
        )
        set_shape_fill(circle, color)
        circle.line.fill.background()

        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = port
        p.font.size = Pt(28)
        p.font.color.rgb = COLORS['white']
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Port name and description
        add_styled_textbox(
            slide, x + Inches(1.8), y + Inches(0.5), Inches(3.5), Inches(0.6),
            name, font_size=26, font_color=color, bold=True
        )
        add_styled_textbox(
            slide, x + Inches(1.8), y + Inches(1.1), Inches(3.5), Inches(0.8),
            desc, font_size=16, font_color=COLORS['dark_gray']
        )

    # Memory tip
    add_rounded_box(
        slide, Inches(2.5), Inches(6.5), Inches(8.33), Inches(0.8),
        COLORS['sky_blue'],
        "💡 Remember: 80 for regular, 443 for secure!",
        text_size=20, text_color=COLORS['dark_blue']
    )

def create_url_anatomy_slide(prs):
    """Slide 10: Anatomy of a URL."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Anatomy of a URL")

    # URL example
    add_rounded_box(
        slide, Inches(0.5), Inches(1.8), Inches(12.33), Inches(1),
        COLORS['light_gray'],
        "https://www.google.com/search?q=cats",
        text_size=28, text_color=COLORS['dark_gray']
    )

    # Breakdown with colored highlights
    parts = [
        ("https://", "Protocol", "The language", COLORS['green'], Inches(0.7)),
        ("www.", "Subdomain", "Optional prefix", COLORS['purple'], Inches(2.7)),
        ("google", "Domain", "The website name", COLORS['orange'], Inches(4.3)),
        (".com", "TLD", "Top-level domain", COLORS['medium_blue'], Inches(6.5)),
        ("/search", "Path", "Page location", COLORS['red'], Inches(8.3)),
        ("?q=cats", "Query", "Extra info", COLORS['dark_gray'], Inches(10.3)),
    ]

    for text, name, desc, color, x in parts:
        y = Inches(3.2)

        # Colored box for the part
        box = add_rounded_box(
            slide, x, y, Inches(1.8), Inches(1.8),
            color
        )

        add_styled_textbox(
            slide, x, y + Inches(0.2), Inches(1.8), Inches(0.5),
            text, font_size=14, font_color=COLORS['white'],
            alignment=PP_ALIGN.CENTER, bold=True
        )
        add_styled_textbox(
            slide, x, y + Inches(0.7), Inches(1.8), Inches(0.5),
            name, font_size=16, font_color=COLORS['white'],
            alignment=PP_ALIGN.CENTER, bold=True
        )
        add_styled_textbox(
            slide, x, y + Inches(1.2), Inches(1.8), Inches(0.5),
            desc, font_size=11, font_color=COLORS['white'],
            alignment=PP_ALIGN.CENTER
        )

    # Key takeaway
    add_styled_textbox(
        slide, Inches(1), Inches(5.5), Inches(11.33), Inches(1.5),
        "🔑 The DOMAIN (like 'google') is the main identifier!\n"
        "It's what you type to find a website.",
        font_size=20, font_color=COLORS['dark_blue'], bold=True,
        alignment=PP_ALIGN.CENTER
    )

def create_website_journey_slide(prs):
    """Slide 11: How a website loads - step by step."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "The Journey of Loading a Website")

    steps = [
        ("1", "You type\ngoogle.com", COLORS['client_color']),
        ("2", "Computer asks DNS\nfor IP address", COLORS['purple']),
        ("3", "DNS returns IP:\n142.250.80.46", COLORS['purple']),
        ("4", "Browser sends\nrequest to server", COLORS['orange']),
        ("5", "Server sends\nwebpage back", COLORS['green']),
        ("6", "Browser displays\nthe page!", COLORS['medium_blue']),
    ]

    for i, (num, text, color) in enumerate(steps):
        x = Inches(0.5 + (i % 3) * 4.2)
        y = Inches(1.8 + (i // 3) * 2.6)

        # Step box
        box = add_rounded_box(
            slide, x, y, Inches(3.8), Inches(2.2),
            color
        )

        # Step number
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + Inches(0.2), y + Inches(0.2),
            Inches(0.6), Inches(0.6)
        )
        set_shape_fill(circle, COLORS['white'])
        circle.line.fill.background()

        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(22)
        p.font.color.rgb = color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Step text
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.9), Inches(3.4), Inches(1.2),
            text, font_size=16, font_color=COLORS['white'],
            alignment=PP_ALIGN.CENTER, bold=True
        )

    # Key point
    add_styled_textbox(
        slide, Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.5),
        "⚡ All of this happens in less than a second!",
        font_size=18, font_color=COLORS['dark_blue'], bold=True,
        alignment=PP_ALIGN.CENTER
    )

def create_dns_slide(prs):
    """Slide 12: DNS explained."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "DNS - The Internet's Phone Book")

    # Analogy
    add_styled_textbox(
        slide, Inches(0.8), Inches(1.8), Inches(11.73), Inches(0.8),
        "You know your friend's NAME, but your phone needs their NUMBER to call them.\n"
        "DNS works the same way for websites!",
        font_size=20, font_color=COLORS['dark_gray'],
        alignment=PP_ALIGN.CENTER
    )

    # Visual: Name -> DNS -> IP
    # Domain name box
    add_rounded_box(
        slide, Inches(0.8), Inches(3), Inches(3.5), Inches(1.5),
        COLORS['client_color'],
        "google.com\n(Domain Name)",
        text_size=22, text_color=COLORS['dark_blue']
    )

    # Arrow
    add_arrow(slide, Inches(4.5), Inches(3.5), Inches(1.2), Inches(0.5), COLORS['orange'], 'right')

    # DNS box
    add_rounded_box(
        slide, Inches(5.9), Inches(3), Inches(2.5), Inches(1.5),
        COLORS['purple'],
        "📖\nDNS",
        text_size=22, text_color=COLORS['white']
    )

    # Arrow
    add_arrow(slide, Inches(8.6), Inches(3.5), Inches(1.2), Inches(0.5), COLORS['orange'], 'right')

    # IP address box
    add_rounded_box(
        slide, Inches(10), Inches(3), Inches(2.5), Inches(1.5),
        COLORS['server_color'],
        "142.250.80.46\n(IP Address)",
        text_size=18, text_color=COLORS['dark_blue']
    )

    # Explanation
    add_styled_textbox(
        slide, Inches(1), Inches(5), Inches(11.33), Inches(1.5),
        "DNS (Domain Name System) translates human-friendly names\n"
        "into computer-friendly IP addresses.\n\n"
        "Without DNS, you'd have to memorize numbers like 142.250.80.46!",
        font_size=18, font_color=COLORS['dark_gray'],
        alignment=PP_ALIGN.CENTER
    )

def create_tracert_slide(prs):
    """Slide 13: The tracert command."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "tracert - Trace the Route!", "Hands-on Activity")

    # What it does
    add_rounded_box(
        slide, Inches(0.8), Inches(1.8), Inches(11.73), Inches(1.2),
        COLORS['medium_blue'],
        "tracert shows every STOP (hop) your data makes on its journey to a website!",
        text_size=22, text_color=COLORS['white']
    )

    # Command box
    add_styled_textbox(
        slide, Inches(0.8), Inches(3.3), Inches(3), Inches(0.5),
        "Try this command:", font_size=18, font_color=COLORS['dark_blue'], bold=True
    )

    cmd_box = add_rounded_box(
        slide, Inches(0.8), Inches(3.9), Inches(11.73), Inches(0.9),
        COLORS['dark_gray'],
        "tracert google.com",
        text_size=24, text_color=COLORS['green']
    )

    # Visual: hops
    add_styled_textbox(
        slide, Inches(0.8), Inches(5.1), Inches(3), Inches(0.5),
        "What you'll see:", font_size=18, font_color=COLORS['dark_blue'], bold=True
    )

    hops = ["Your Router", "ISP", "Data Center", "Google Server"]

    for i, hop in enumerate(hops):
        x = Inches(0.8 + i * 3)

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x, Inches(5.7),
            Inches(1), Inches(1)
        )
        set_shape_fill(circle, COLORS['orange'] if i < 3 else COLORS['green'])
        circle.line.fill.background()

        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['white']
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        add_styled_textbox(
            slide, x - Inches(0.4), Inches(6.8), Inches(1.8), Inches(0.5),
            hop, font_size=12, font_color=COLORS['dark_gray'],
            alignment=PP_ALIGN.CENTER
        )

        if i < 3:
            # Arrow between hops
            add_arrow(slide, x + Inches(1.1), Inches(6), Inches(1.6), Inches(0.4),
                     COLORS['light_blue'], 'right')

def create_recap_slide(prs):
    """Slide 14: Recap slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_shape(slide, "Today's Key Takeaways")

    takeaways = [
        ("CLIENT", "A device that ASKS for information", COLORS['client_color']),
        ("SERVER", "A computer that SERVES information", COLORS['server_color']),
        ("HTTP/HTTPS", "The language of the web (S = Secure!)", COLORS['orange']),
        ("PORTS", "Like apartment numbers for services", COLORS['purple']),
        ("DNS", "Translates names to IP addresses", COLORS['medium_blue']),
    ]

    for i, (term, definition, color) in enumerate(takeaways):
        y = Inches(1.7 + i * 1.05)

        # Term box
        term_box = add_rounded_box(
            slide, Inches(0.8), y, Inches(2.8), Inches(0.85),
            color,
            term, text_size=20, text_color=COLORS['white']
        )

        # Definition
        add_styled_textbox(
            slide, Inches(4), y + Inches(0.2), Inches(8.5), Inches(0.6),
            definition, font_size=20, font_color=COLORS['dark_gray']
        )

    # Encouragement
    add_styled_textbox(
        slide, Inches(1), Inches(6.8), Inches(11.33), Inches(0.5),
        "🎉 Great job today! Time for Kahoot!",
        font_size=24, font_color=COLORS['dark_blue'], bold=True,
        alignment=PP_ALIGN.CENTER
    )

def create_questions_slide(prs):
    """Slide 15: Questions slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    set_shape_fill(background, COLORS['dark_blue'])
    background.line.fill.background()

    # Question mark decoration
    add_styled_textbox(
        slide, Inches(0), Inches(1.5), Inches(13.33), Inches(2),
        "❓", font_size=120, alignment=PP_ALIGN.CENTER
    )

    # Title
    add_styled_textbox(
        slide, Inches(0), Inches(3.5), Inches(13.33), Inches(1),
        "Questions?",
        font_size=60, font_color=COLORS['white'], bold=True,
        alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold"
    )

    add_styled_textbox(
        slide, Inches(0), Inches(5), Inches(13.33), Inches(0.6),
        "Ask away before we start the Kahoot!",
        font_size=24, font_color=COLORS['sky_blue'],
        alignment=PP_ALIGN.CENTER
    )

# ============================================================================
# MAIN FUNCTION
# ============================================================================

def main():
    """Create the presentation."""
    # Create presentation with widescreen dimensions (16:9)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("Creating slides...")

    # Create all slides
    create_title_slide(prs)
    print("  [OK] Title slide")

    create_agenda_slide(prs)
    print("  [OK] Agenda slide")

    create_client_intro_slide(prs)
    print("  [OK] Client intro slide")

    create_server_intro_slide(prs)
    print("  [OK] Server intro slide")

    create_client_server_diagram_slide(prs)
    print("  [OK] Client-Server diagram slide")

    create_http_intro_slide(prs)
    print("  [OK] HTTP intro slide")

    create_https_slide(prs)
    print("  [OK] HTTPS slide")

    create_ports_intro_slide(prs)
    print("  [OK] Ports intro slide")

    create_ports_numbers_slide(prs)
    print("  [OK] Port numbers slide")

    create_url_anatomy_slide(prs)
    print("  [OK] URL anatomy slide")

    create_website_journey_slide(prs)
    print("  [OK] Website journey slide")

    create_dns_slide(prs)
    print("  [OK] DNS slide")

    create_tracert_slide(prs)
    print("  [OK] tracert slide")

    create_recap_slide(prs)
    print("  [OK] Recap slide")

    create_questions_slide(prs)
    print("  [OK] Questions slide")

    # Save the presentation
    output_path = os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "2026-01-25_Clients_Servers_BEAUTIFUL.pptx"
    )
    prs.save(output_path)
    print(f"\n[SUCCESS] Presentation saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
