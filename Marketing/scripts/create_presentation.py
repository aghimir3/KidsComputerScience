"""
EverestITT Parent Marketing Presentation
Premium design — persuasive, modern, clean
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Premium Color Palette ───────────────────────────────────────────────────

COLORS = {
    # Brand
    "dark_blue": RGBColor(30, 58, 95),
    "orange": RGBColor(255, 153, 51),

    # Premium accents
    "medium_blue": RGBColor(52, 103, 166),
    "light_blue": RGBColor(100, 149, 237),
    "sky_blue": RGBColor(200, 220, 240),
    "ai_purple": RGBColor(103, 58, 183),
    "light_purple": RGBColor(237, 231, 246),
    "green": RGBColor(46, 125, 50),
    "light_green": RGBColor(232, 245, 233),
    "teal": RGBColor(0, 150, 136),
    "gold": RGBColor(255, 193, 7),
    "coral": RGBColor(255, 111, 97),

    # Neutrals
    "white": RGBColor(255, 255, 255),
    "off_white": RGBColor(248, 249, 250),
    "light_gray": RGBColor(240, 240, 240),
    "medium_gray": RGBColor(158, 158, 158),
    "dark_gray": RGBColor(55, 55, 55),
    "charcoal": RGBColor(33, 37, 41),
    "black": RGBColor(0, 0, 0),
}

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ─── Helper Functions ────────────────────────────────────────────────────────

def set_shape_fill(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_shape_line(shape, color=None, width_pt=0):
    line = shape.line
    if color is None:
        line.fill.background()
    else:
        line.color.rgb = color
        line.width = Pt(width_pt)


def add_rounded_box(slide, left, top, width, height, fill_color, border_color=None, border_width=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    set_shape_fill(shape, fill_color)
    if border_color:
        set_shape_line(shape, border_color, border_width)
    else:
        set_shape_line(shape)
    # Adjust corner rounding
    shape.adjustments[0] = 0.08
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                font_color=None, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Segoe UI", line_spacing=1.2, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color or COLORS["dark_gray"]
    p.font.bold = bold
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, font_size=16,
                          font_color=None, font_name="Segoe UI", line_spacing=1.5,
                          alignment=PP_ALIGN.LEFT, bullet=False):
    """Add a textbox with multiple paragraphs. lines is a list of (text, bold, color) tuples or strings."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, bold, color = line, False, font_color or COLORS["dark_gray"]
        else:
            text, bold, color = line[0], line[1] if len(line) > 1 else False, line[2] if len(line) > 2 else (font_color or COLORS["dark_gray"])

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if bullet:
            p.text = f"  •  {text}"
        else:
            p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(2)
        p.space_before = Pt(2)
        p.line_spacing = Pt(font_size * line_spacing)

    return txBox


def add_full_bg(slide, color):
    """Fill the entire slide background with a color."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    set_shape_fill(bg, color)
    set_shape_line(bg)
    return bg


def add_accent_line(slide, left, top, width, color, height=0.06):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    set_shape_fill(line, color)
    set_shape_line(line)
    return line


def add_icon_circle(slide, left, top, size, color, text, text_size=20, text_color=None):
    """Add a circle with text inside — used as an icon/number."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    set_shape_fill(shape, color)
    set_shape_line(shape)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(text_size)
    p.font.color.rgb = text_color or COLORS["white"]
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


# ─── SLIDE 1: Title ─────────────────────────────────────────────────────────

def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Full dark blue background
    add_full_bg(slide, COLORS["dark_blue"])

    # Subtle gradient overlay — lighter blue strip at top
    add_rounded_box(slide, 0.5, 0.4, 12.33, 0.08, COLORS["orange"])

    # Main title
    add_textbox(slide, 1.0, 1.2, 11.33, 1.2,
                "Kids Computer Science Class",
                font_size=48, font_color=COLORS["white"], bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold")

    # Subtitle
    add_textbox(slide, 1.0, 2.3, 11.33, 0.8,
                "Empowering Future Innovators Through Technology",
                font_size=22, font_color=COLORS["sky_blue"], bold=False,
                alignment=PP_ALIGN.CENTER, italic=True)

    # Orange accent line
    add_accent_line(slide, 4.5, 3.3, 4.33, COLORS["orange"], 0.05)

    # Three info boxes in a row
    box_y = 4.0
    box_h = 1.6
    box_w = 3.4
    gap = 0.55
    start_x = 1.5

    # Box 1: When
    add_rounded_box(slide, start_x, box_y, box_w, box_h, COLORS["medium_blue"], border_color=COLORS["light_blue"], border_width=1)
    add_textbox(slide, start_x, box_y + 0.2, box_w, 0.4, "WHEN", font_size=14,
                font_color=COLORS["sky_blue"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, start_x, box_y + 0.55, box_w, 0.4, "Every Saturday", font_size=22,
                font_color=COLORS["white"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, start_x, box_y + 1.0, box_w, 0.4, "9:00 AM - 1:00 PM Pacific", font_size=14,
                font_color=COLORS["sky_blue"], alignment=PP_ALIGN.CENTER)

    # Box 2: Who
    x2 = start_x + box_w + gap
    add_rounded_box(slide, x2, box_y, box_w, box_h, COLORS["medium_blue"], border_color=COLORS["light_blue"], border_width=1)
    add_textbox(slide, x2, box_y + 0.2, box_w, 0.4, "WHO", font_size=14,
                font_color=COLORS["sky_blue"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x2, box_y + 0.55, box_w, 0.4, "Ages 10 - 18", font_size=22,
                font_color=COLORS["white"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x2, box_y + 1.0, box_w, 0.4, "All experience levels welcome", font_size=14,
                font_color=COLORS["sky_blue"], alignment=PP_ALIGN.CENTER)

    # Box 3: How
    x3 = x2 + box_w + gap
    add_rounded_box(slide, x3, box_y, box_w, box_h, COLORS["medium_blue"], border_color=COLORS["light_blue"], border_width=1)
    add_textbox(slide, x3, box_y + 0.2, box_w, 0.4, "HOW", font_size=14,
                font_color=COLORS["sky_blue"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x3, box_y + 0.55, box_w, 0.4, "Live & Interactive", font_size=22,
                font_color=COLORS["white"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x3, box_y + 1.0, box_w, 0.4, "Real teachers, not recorded videos", font_size=14,
                font_color=COLORS["sky_blue"], alignment=PP_ALIGN.CENTER)

    # Bottom bar
    add_accent_line(slide, 0.5, 6.9, 12.33, COLORS["orange"], 0.05)
    add_textbox(slide, 1.0, 6.3, 11.33, 0.5,
                "EverestITT  |  everestitt.com  |  626-419-6649",
                font_size=16, font_color=COLORS["medium_gray"],
                alignment=PP_ALIGN.CENTER)


# ─── SLIDE 2: Why AI Matters ────────────────────────────────────────────────

def create_why_ai_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White background with dark blue header
    add_full_bg(slide, COLORS["white"])

    # Header bar
    add_rounded_box(slide, 0, 0, 13.333, 1.2, COLORS["dark_blue"])
    add_textbox(slide, 0.8, 0.25, 11.73, 0.7,
                "Why Your Child Needs AI Skills — Now",
                font_size=34, font_color=COLORS["white"], bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Segoe UI Semibold")
    add_accent_line(slide, 0, 1.2, 13.333, COLORS["orange"], 0.06)

    # Big stat — hero statement
    add_textbox(slide, 0.8, 1.7, 11.73, 0.8,
                "Every top university and employer is looking for students who understand AI.",
                font_size=24, font_color=COLORS["dark_blue"], bold=True,
                alignment=PP_ALIGN.CENTER)

    # Three impact cards
    card_w = 3.6
    card_h = 2.8
    card_y = 2.9
    gap = 0.55
    start_x = 1.2

    # Card 1: College
    add_rounded_box(slide, start_x, card_y, card_w, card_h, COLORS["off_white"],
                    border_color=COLORS["dark_blue"], border_width=2)
    add_icon_circle(slide, start_x + 1.4, card_y + 0.3, 0.8, COLORS["dark_blue"], "🎓", text_size=28)
    add_textbox(slide, start_x + 0.2, card_y + 1.3, card_w - 0.4, 0.4,
                "College Applications", font_size=18, font_color=COLORS["dark_blue"],
                bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, start_x + 0.2, card_y + 1.8, card_w - 0.4, 0.9,
                "Stanford, MIT, and Harvard value students who can build with AI — not just use it",
                font_size=14, font_color=COLORS["dark_gray"], alignment=PP_ALIGN.CENTER)

    # Card 2: Careers
    x2 = start_x + card_w + gap
    add_rounded_box(slide, x2, card_y, card_w, card_h, COLORS["off_white"],
                    border_color=COLORS["orange"], border_width=2)
    add_icon_circle(slide, x2 + 1.4, card_y + 0.3, 0.8, COLORS["orange"], "💼", text_size=28)
    add_textbox(slide, x2 + 0.2, card_y + 1.3, card_w - 0.4, 0.4,
                "Future Careers", font_size=18, font_color=COLORS["orange"],
                bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x2 + 0.2, card_y + 1.8, card_w - 0.4, 0.9,
                "AI skills are becoming required across every industry — tech, medicine, law, business",
                font_size=14, font_color=COLORS["dark_gray"], alignment=PP_ALIGN.CENTER)

    # Card 3: Head start
    x3 = x2 + card_w + gap
    add_rounded_box(slide, x3, card_y, card_w, card_h, COLORS["off_white"],
                    border_color=COLORS["green"], border_width=2)
    add_icon_circle(slide, x3 + 1.4, card_y + 0.3, 0.8, COLORS["green"], "🚀", text_size=28)
    add_textbox(slide, x3 + 0.2, card_y + 1.3, card_w - 0.4, 0.4,
                "Head Start", font_size=18, font_color=COLORS["green"],
                bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x3 + 0.2, card_y + 1.8, card_w - 0.4, 0.9,
                "The students who learn AI now will have a massive advantage over those who don't",
                font_size=14, font_color=COLORS["dark_gray"], alignment=PP_ALIGN.CENTER)

    # Bottom urgency banner
    add_rounded_box(slide, 0.8, 6.2, 11.73, 0.7, COLORS["dark_blue"])
    add_textbox(slide, 0.8, 6.25, 11.73, 0.6,
                "The question isn't IF your child needs AI skills — it's whether they'll be ready.",
                font_size=18, font_color=COLORS["orange"], bold=True, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 3: Meet the Team ─────────────────────────────────────────────────

def create_team_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_full_bg(slide, COLORS["white"])

    # Header
    add_rounded_box(slide, 0, 0, 13.333, 1.2, COLORS["dark_blue"])
    add_textbox(slide, 0.8, 0.25, 11.73, 0.7,
                "Learn From Industry Professionals",
                font_size=34, font_color=COLORS["white"], bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Segoe UI Semibold")
    add_accent_line(slide, 0, 1.2, 13.333, COLORS["orange"], 0.06)

    # Founder — featured prominently
    add_rounded_box(slide, 0.8, 1.6, 11.73, 2.0, COLORS["dark_blue"])
    add_accent_line(slide, 0.8, 1.6, 0.15, COLORS["orange"], 2.0)  # Left accent bar

    add_textbox(slide, 1.3, 1.7, 4.0, 0.5, "FOUNDER & LEAD INSTRUCTOR",
                font_size=12, font_color=COLORS["sky_blue"], bold=True)
    add_textbox(slide, 1.3, 2.05, 4.0, 0.5, "Samir Thapa (Bijay)",
                font_size=26, font_color=COLORS["orange"], bold=True)

    # Founder credentials in 3 columns
    cred_y = 2.65
    add_textbox(slide, 1.3, cred_y, 3.5, 0.7,
                "Principal Technology Specialist\nat Microsoft",
                font_size=14, font_color=COLORS["white"])
    add_textbox(slide, 5.0, cred_y, 3.5, 0.7,
                "30+ years of experience\nin the IT industry",
                font_size=14, font_color=COLORS["white"])
    add_textbox(slide, 8.7, cred_y, 3.5, 0.7,
                "Multiple Master's degrees\nElectronics & Communication Engineering",
                font_size=14, font_color=COLORS["white"])

    # Teachers — two boxes side by side
    teacher_y = 4.0
    teacher_w = 5.67
    teacher_h = 2.0

    # Abhigya
    add_rounded_box(slide, 0.8, teacher_y, teacher_w, teacher_h, COLORS["off_white"],
                    border_color=COLORS["medium_blue"], border_width=2)
    add_textbox(slide, 1.1, teacher_y + 0.15, 2.0, 0.3, "TEACHER",
                font_size=11, font_color=COLORS["medium_blue"], bold=True)
    add_textbox(slide, 1.1, teacher_y + 0.45, teacher_w - 0.6, 0.4, "Abhigya",
                font_size=22, font_color=COLORS["dark_blue"], bold=True)
    add_multiline_textbox(slide, 1.1, teacher_y + 0.9, teacher_w - 0.6, 1.0, [
        "Lead Software Engineer at Flagship Financial Group",
        "Master's in Software Engineering",
        "Open Source (FOSS) Contributor",
    ], font_size=13, font_color=COLORS["dark_gray"], line_spacing=1.4, bullet=True)

    # Sanskar
    x2 = 0.8 + teacher_w + 0.4
    add_rounded_box(slide, x2, teacher_y, teacher_w, teacher_h, COLORS["off_white"],
                    border_color=COLORS["teal"], border_width=2)
    add_textbox(slide, x2 + 0.3, teacher_y + 0.15, 2.0, 0.3, "TEACHER",
                font_size=11, font_color=COLORS["teal"], bold=True)
    add_textbox(slide, x2 + 0.3, teacher_y + 0.45, teacher_w - 0.6, 0.4, "Sanskar",
                font_size=22, font_color=COLORS["dark_blue"], bold=True)
    add_multiline_textbox(slide, x2 + 0.3, teacher_y + 0.9, teacher_w - 0.6, 1.0, [
        "CS Engineering 2025 — Cal State University",
        "World's youngest Microsoft Cloud Certified",
        "Software Engineer at Google, Northrop Grumman",
    ], font_size=13, font_color=COLORS["dark_gray"], line_spacing=1.4, bullet=True)

    # Support team — bottom row
    support_y = 6.3
    add_textbox(slide, 0.8, support_y, 11.73, 0.5,
                "Support Team:   Sachin sir (Assistant Teacher)   •   Ishwari maam (Assistant Teacher)   •   Pratikshya maam (Program Manager — Parent Contact)",
                font_size=14, font_color=COLORS["medium_gray"], alignment=PP_ALIGN.CENTER)


# ─── SLIDE 4: Year Roadmap ──────────────────────────────────────────────────

def create_roadmap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_full_bg(slide, COLORS["white"])

    # Header
    add_rounded_box(slide, 0, 0, 13.333, 1.2, COLORS["dark_blue"])
    add_textbox(slide, 0.8, 0.25, 11.73, 0.7,
                "A Complete Year of Computer Science",
                font_size=34, font_color=COLORS["white"], bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Segoe UI Semibold")
    add_accent_line(slide, 0, 1.2, 13.333, COLORS["orange"], 0.06)

    # Timeline — three phases
    phase_h = 1.45
    phase_w = 11.73
    start_y = 1.7

    # Phase 1: Jan-Mar
    p1_y = start_y
    add_rounded_box(slide, 0.8, p1_y, phase_w, phase_h, COLORS["off_white"],
                    border_color=COLORS["medium_blue"], border_width=2)
    # Month badge
    add_rounded_box(slide, 1.0, p1_y + 0.2, 2.2, 1.05, COLORS["medium_blue"])
    add_textbox(slide, 1.0, p1_y + 0.3, 2.2, 0.4, "JAN — MAR",
                font_size=20, font_color=COLORS["white"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, p1_y + 0.7, 2.2, 0.4, "Foundations",
                font_size=14, font_color=COLORS["sky_blue"], alignment=PP_ALIGN.CENTER)
    # Content
    add_textbox(slide, 3.5, p1_y + 0.2, 8.5, 0.4, "PC Hardware, Networking & Cloud",
                font_size=20, font_color=COLORS["dark_blue"], bold=True)
    add_textbox(slide, 3.5, p1_y + 0.7, 8.5, 0.6,
                "Computer components  •  How the internet works  •  Cloud computing  •  Linux & Command Line  •  Cloud collaboration",
                font_size=14, font_color=COLORS["dark_gray"])

    # Phase 2: Apr-Jun (HIGHLIGHTED)
    p2_y = p1_y + phase_h + 0.3
    add_rounded_box(slide, 0.8, p2_y, phase_w, phase_h + 0.3, COLORS["ai_purple"],
                    border_color=COLORS["gold"], border_width=3)
    # Month badge
    add_rounded_box(slide, 1.0, p2_y + 0.25, 2.2, 1.15, COLORS["white"])
    add_textbox(slide, 1.0, p2_y + 0.35, 2.2, 0.4, "APR — JUN",
                font_size=20, font_color=COLORS["ai_purple"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, p2_y + 0.8, 2.2, 0.4, "AI Unit",
                font_size=14, font_color=COLORS["ai_purple"], alignment=PP_ALIGN.CENTER)
    # Content
    add_textbox(slide, 3.5, p2_y + 0.15, 8.5, 0.4, "Artificial Intelligence",
                font_size=22, font_color=COLORS["white"], bold=True)
    add_textbox(slide, 3.5, p2_y + 0.6, 8.5, 0.5,
                "Prompt Engineering  •  ChatGPT  •  Claude Code  •  OpenClaw  •  AI Agents  •  Ethics & Safety",
                font_size=14, font_color=COLORS["sky_blue"])
    # NOW ENROLLING badge
    add_rounded_box(slide, 9.5, p2_y + 1.0, 2.8, 0.5, COLORS["gold"])
    add_textbox(slide, 9.5, p2_y + 1.02, 2.8, 0.45, "★  STARTS APRIL 4",
                font_size=14, font_color=COLORS["charcoal"], bold=True, alignment=PP_ALIGN.CENTER)

    # Phase 3: Jul-Dec
    p3_y = p2_y + phase_h + 0.6
    add_rounded_box(slide, 0.8, p3_y, phase_w, phase_h, COLORS["off_white"],
                    border_color=COLORS["green"], border_width=2)
    # Month badge
    add_rounded_box(slide, 1.0, p3_y + 0.2, 2.2, 1.05, COLORS["green"])
    add_textbox(slide, 1.0, p3_y + 0.3, 2.2, 0.4, "JUL — DEC",
                font_size=20, font_color=COLORS["white"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, p3_y + 0.7, 2.2, 0.4, "Programming",
                font_size=14, font_color=COLORS["light_green"], alignment=PP_ALIGN.CENTER)
    # Content
    add_textbox(slide, 3.5, p3_y + 0.2, 8.5, 0.4, "Python Programming",
                font_size=20, font_color=COLORS["dark_blue"], bold=True)
    add_textbox(slide, 3.5, p3_y + 0.7, 8.5, 0.6,
                "Python fundamentals  •  VS Code & GitHub  •  Logic & problem-solving  •  Data structures  •  Capstone project",
                font_size=14, font_color=COLORS["dark_gray"])

    # Bottom
    add_textbox(slide, 0.8, 6.7, 11.73, 0.5,
                "Each phase builds on the last — by December, your child has a complete tech foundation",
                font_size=15, font_color=COLORS["medium_gray"], alignment=PP_ALIGN.CENTER, italic=True)


# ─── SLIDE 5: AI Deep Dive ──────────────────────────────────────────────────

def create_ai_deep_dive_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full purple background for impact
    add_full_bg(slide, COLORS["ai_purple"])

    # Top accent
    add_accent_line(slide, 0.5, 0.4, 12.33, COLORS["gold"], 0.05)

    # Title
    add_textbox(slide, 0.8, 0.7, 11.73, 0.7,
                "STARTING APRIL: Artificial Intelligence",
                font_size=36, font_color=COLORS["white"], bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold")
    add_textbox(slide, 0.8, 1.4, 11.73, 0.5,
                "Your child won't just use AI — they'll understand how it works and build with it",
                font_size=17, font_color=COLORS["sky_blue"], alignment=PP_ALIGN.CENTER, italic=True)

    # Six topic cards in 2x3 grid
    card_w = 3.5
    card_h = 1.85
    row1_y = 2.3
    row2_y = 4.4
    gap_x = 0.55
    start_x = 1.2

    topics = [
        ("How AI Works", "Understand what's behind the\nbuzz — LLMs, training data,\nand neural networks", COLORS["dark_blue"]),
        ("Prompt Engineering", "The #1 skill for getting\nprofessional results from\nany AI tool", COLORS["orange"]),
        ("ChatGPT & Claude Code", "Hands-on with the same\ntools professionals use\nat work every day", COLORS["medium_blue"]),
        ("OpenClaw", "Build and customize AI\nagents that can work\nautonomously", COLORS["teal"]),
        ("AI Agents", "How autonomous AI systems\nwork — and how to create\nyour own", COLORS["green"]),
        ("Ethics & Safety", "Use AI responsibly —\nunderstand limitations,\nbias, and best practices", COLORS["coral"]),
    ]

    for i, (title, desc, color) in enumerate(topics):
        row = i // 3
        col = i % 3
        x = start_x + col * (card_w + gap_x)
        y = row1_y if row == 0 else row2_y

        add_rounded_box(slide, x, y, card_w, card_h, color)
        add_textbox(slide, x + 0.2, y + 0.2, card_w - 0.4, 0.4,
                    title, font_size=19, font_color=COLORS["white"],
                    bold=True, alignment=PP_ALIGN.LEFT)
        add_accent_line(slide, x + 0.2, y + 0.6, card_w - 0.4, COLORS["white"], 0.03)
        add_textbox(slide, x + 0.2, y + 0.75, card_w - 0.4, 1.0,
                    desc, font_size=13, font_color=COLORS["white"],
                    alignment=PP_ALIGN.LEFT)

    # Bottom accent
    add_accent_line(slide, 0.5, 6.9, 12.33, COLORS["gold"], 0.05)
    add_textbox(slide, 0.8, 6.4, 11.73, 0.5,
                "These are skills most adults don't have yet. Your child will.",
                font_size=16, font_color=COLORS["gold"], bold=True, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 6: What Your Child Will Be Able To Do ────────────────────────────

def create_outcomes_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_full_bg(slide, COLORS["white"])

    # Header
    add_rounded_box(slide, 0, 0, 13.333, 1.2, COLORS["dark_blue"])
    add_textbox(slide, 0.8, 0.25, 11.73, 0.7,
                "What Your Child Will Be Able To Do",
                font_size=34, font_color=COLORS["white"], bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Segoe UI Semibold")
    add_accent_line(slide, 0, 1.2, 13.333, COLORS["orange"], 0.06)

    # Four outcome rows with numbered circles
    outcomes = [
        ("Write effective prompts that get professional-quality results from any AI tool",
         COLORS["ai_purple"], "1"),
        ("Use ChatGPT, Claude Code, and OpenClaw the way real professionals do at work",
         COLORS["orange"], "2"),
        ("Understand how Large Language Models actually work — not just what buttons to press",
         COLORS["dark_blue"], "3"),
        ("Build their own AI-powered projects they can show off on college applications",
         COLORS["green"], "4"),
    ]

    start_y = 1.7
    row_h = 1.05
    row_gap = 0.15

    for i, (text, color, num) in enumerate(outcomes):
        y = start_y + i * (row_h + row_gap)

        # Background bar
        add_rounded_box(slide, 0.8, y, 11.73, row_h, COLORS["off_white"],
                        border_color=color, border_width=2)

        # Number circle
        add_icon_circle(slide, 1.2, y + 0.15, 0.75, color, num, text_size=24)

        # Text
        add_textbox(slide, 2.3, y + 0.2, 9.8, 0.65,
                    text, font_size=18, font_color=COLORS["dark_gray"])

    # Bottom highlight box
    bottom_y = start_y + 4 * (row_h + row_gap) + 0.3
    add_rounded_box(slide, 1.5, bottom_y, 10.33, 0.85, COLORS["gold"])
    add_textbox(slide, 1.5, bottom_y + 0.15, 10.33, 0.55,
                "Most adults don't have these skills yet — your child will be years ahead.",
                font_size=18, font_color=COLORS["charcoal"], bold=True, alignment=PP_ALIGN.CENTER)

    # Bottom note
    add_textbox(slide, 0.8, 6.5, 11.73, 0.4,
                "All skills are taught hands-on with real projects — not just theory",
                font_size=14, font_color=COLORS["medium_gray"], alignment=PP_ALIGN.CENTER, italic=True)


# ─── SLIDE 7: Why We're Different ────────────────────────────────────────────

def create_differentiators_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_full_bg(slide, COLORS["white"])

    # Header
    add_rounded_box(slide, 0, 0, 13.333, 1.2, COLORS["dark_blue"])
    add_textbox(slide, 0.8, 0.25, 11.73, 0.7,
                "This Isn't Your Average Computer Class",
                font_size=34, font_color=COLORS["white"], bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Segoe UI Semibold")
    add_accent_line(slide, 0, 1.2, 13.333, COLORS["orange"], 0.06)

    # Two columns — comparison
    col_w = 5.5
    col_h = 4.5
    col_y = 1.7

    # LEFT: Other Programs (gray, negative)
    left_x = 0.8
    add_rounded_box(slide, left_x, col_y, col_w, col_h, COLORS["off_white"],
                    border_color=COLORS["coral"], border_width=2)
    add_textbox(slide, left_x + 0.3, col_y + 0.2, col_w - 0.6, 0.4,
                "Other Programs", font_size=22, font_color=COLORS["coral"], bold=True)
    add_accent_line(slide, left_x + 0.3, col_y + 0.65, col_w - 0.6, COLORS["coral"], 0.03)

    other_items = [
        "Watch pre-recorded videos alone",
        "No real teacher interaction",
        "Generic, one-size-fits-all content",
        "No homework or accountability",
        "No feedback on progress",
    ]
    for j, item in enumerate(other_items):
        add_textbox(slide, left_x + 0.3, col_y + 0.9 + j * 0.6, col_w - 0.6, 0.5,
                    f"✗   {item}", font_size=15, font_color=COLORS["dark_gray"])

    # RIGHT: EverestITT (green, positive)
    right_x = 0.8 + col_w + 0.73
    add_rounded_box(slide, right_x, col_y, col_w, col_h, COLORS["off_white"],
                    border_color=COLORS["green"], border_width=2)
    add_textbox(slide, right_x + 0.3, col_y + 0.2, col_w - 0.6, 0.4,
                "EverestITT", font_size=22, font_color=COLORS["green"], bold=True)
    add_accent_line(slide, right_x + 0.3, col_y + 0.65, col_w - 0.6, COLORS["green"], 0.03)

    our_items = [
        "Live instruction every Saturday",
        "Industry professionals as teachers",
        "Structured, progressive curriculum",
        "Weekly homework with grading",
        "Real feedback and mentorship",
    ]
    for j, item in enumerate(our_items):
        add_textbox(slide, right_x + 0.3, col_y + 0.9 + j * 0.6, col_w - 0.6, 0.5,
                    f"✓   {item}", font_size=15, font_color=COLORS["dark_gray"])

    # VS circle in the middle
    vs_x = left_x + col_w - 0.15
    add_icon_circle(slide, vs_x, col_y + 1.8, 1.0, COLORS["dark_blue"], "VS", text_size=20)

    # Bottom
    add_textbox(slide, 0.8, 6.5, 11.73, 0.4,
                "Your child deserves more than a playlist of tutorials.",
                font_size=16, font_color=COLORS["dark_blue"], bold=True, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 8: Hangout Sessions ───────────────────────────────────────────────

def create_hangout_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_full_bg(slide, COLORS["white"])

    # Header
    add_rounded_box(slide, 0, 0, 13.333, 1.2, COLORS["dark_blue"])
    add_textbox(slide, 0.8, 0.25, 11.73, 0.7,
                "Support Beyond Saturday — Daily Hangout Sessions",
                font_size=32, font_color=COLORS["white"], bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Segoe UI Semibold")
    add_accent_line(slide, 0, 1.2, 13.333, COLORS["orange"], 0.06)

    # Big schedule banner
    add_rounded_box(slide, 1.5, 1.7, 10.33, 1.3, COLORS["dark_blue"])
    add_textbox(slide, 1.5, 1.8, 10.33, 0.5, "Monday — Friday",
                font_size=28, font_color=COLORS["orange"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.5, 2.3, 10.33, 0.5, "4:30 PM — 5:30 PM Pacific Time",
                font_size=22, font_color=COLORS["white"], bold=True, alignment=PP_ALIGN.CENTER)

    # Three benefit cards
    card_w = 3.5
    card_h = 2.6
    card_y = 3.4
    gap = 0.55
    start_x = 1.2

    # Card 1: Student Help
    add_rounded_box(slide, start_x, card_y, card_w, card_h, COLORS["off_white"],
                    border_color=COLORS["medium_blue"], border_width=2)
    add_icon_circle(slide, start_x + 1.35, card_y + 0.3, 0.8, COLORS["medium_blue"], "?", text_size=30)
    add_textbox(slide, start_x + 0.2, card_y + 1.3, card_w - 0.4, 0.4,
                "Get Help Anytime", font_size=18, font_color=COLORS["medium_blue"],
                bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, start_x + 0.2, card_y + 1.75, card_w - 0.4, 0.7,
                "Stuck on homework? Have a question? Assistant teachers are there to help every weekday.",
                font_size=14, font_color=COLORS["dark_gray"], alignment=PP_ALIGN.CENTER)

    # Card 2: Homework Support
    x2 = start_x + card_w + gap
    add_rounded_box(slide, x2, card_y, card_w, card_h, COLORS["off_white"],
                    border_color=COLORS["orange"], border_width=2)
    add_icon_circle(slide, x2 + 1.35, card_y + 0.3, 0.8, COLORS["orange"], "✓", text_size=28)
    add_textbox(slide, x2 + 0.2, card_y + 1.3, card_w - 0.4, 0.4,
                "Homework & Classwork", font_size=18, font_color=COLORS["orange"],
                bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x2 + 0.2, card_y + 1.75, card_w - 0.4, 0.7,
                "Students get guided support on all their assignments so they never fall behind.",
                font_size=14, font_color=COLORS["dark_gray"], alignment=PP_ALIGN.CENTER)

    # Card 3: Parent Communication
    x3 = x2 + card_w + gap
    add_rounded_box(slide, x3, card_y, card_w, card_h, COLORS["off_white"],
                    border_color=COLORS["green"], border_width=2)
    add_icon_circle(slide, x3 + 1.35, card_y + 0.3, 0.8, COLORS["green"], "💬", text_size=28)
    add_textbox(slide, x3 + 0.2, card_y + 1.3, card_w - 0.4, 0.4,
                "Parent Communication", font_size=18, font_color=COLORS["green"],
                bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x3 + 0.2, card_y + 1.75, card_w - 0.4, 0.7,
                "Parents can join to ask questions, discuss progress, and stay connected with the team.",
                font_size=14, font_color=COLORS["dark_gray"], alignment=PP_ALIGN.CENTER)

    # Bottom highlight
    add_rounded_box(slide, 1.5, 6.3, 10.33, 0.7, COLORS["gold"])
    add_textbox(slide, 1.5, 6.35, 10.33, 0.6,
                "Your child gets support 6 days a week — not just Saturdays.",
                font_size=17, font_color=COLORS["charcoal"], bold=True, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 9: Enroll Now CTA ────────────────────────────────────────────────

def create_enroll_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full dark blue background
    add_full_bg(slide, COLORS["dark_blue"])

    # Top accent
    add_accent_line(slide, 0.5, 0.4, 12.33, COLORS["orange"], 0.05)

    # Main message
    add_textbox(slide, 1.0, 0.8, 11.33, 1.0,
                "Give Your Child the Skills That Matter",
                font_size=42, font_color=COLORS["white"], bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold")

    add_textbox(slide, 1.0, 1.8, 11.33, 0.6,
                "The AI unit starts April 4th — don't let your child fall behind",
                font_size=20, font_color=COLORS["sky_blue"],
                alignment=PP_ALIGN.CENTER, italic=True)

    # Accent line
    add_accent_line(slide, 4.5, 2.6, 4.33, COLORS["orange"], 0.05)

    # Contact info — three boxes
    box_w = 3.4
    box_h = 2.2
    box_y = 3.1
    gap = 0.55
    start_x = 1.5

    # Box 1: Call
    add_rounded_box(slide, start_x, box_y, box_w, box_h, COLORS["medium_blue"],
                    border_color=COLORS["orange"], border_width=2)
    add_textbox(slide, start_x, box_y + 0.3, box_w, 0.4, "CALL US",
                font_size=14, font_color=COLORS["sky_blue"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, start_x, box_y + 0.8, box_w, 0.5, "626-419-6649",
                font_size=26, font_color=COLORS["white"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, start_x, box_y + 1.4, box_w, 0.5, "Talk to us about\nenrollment today",
                font_size=14, font_color=COLORS["sky_blue"], alignment=PP_ALIGN.CENTER)

    # Box 2: Visit
    x2 = start_x + box_w + gap
    add_rounded_box(slide, x2, box_y, box_w, box_h, COLORS["medium_blue"],
                    border_color=COLORS["orange"], border_width=2)
    add_textbox(slide, x2, box_y + 0.3, box_w, 0.4, "VISIT US ONLINE",
                font_size=14, font_color=COLORS["sky_blue"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x2, box_y + 0.8, box_w, 0.5, "everestitt.com",
                font_size=26, font_color=COLORS["white"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x2, box_y + 1.4, box_w, 0.5, "Learn more about\nour program",
                font_size=14, font_color=COLORS["sky_blue"], alignment=PP_ALIGN.CENTER)

    # Box 3: Contact
    x3 = x2 + box_w + gap
    add_rounded_box(slide, x3, box_y, box_w, box_h, COLORS["medium_blue"],
                    border_color=COLORS["orange"], border_width=2)
    add_textbox(slide, x3, box_y + 0.3, box_w, 0.4, "PARENT CONTACT",
                font_size=14, font_color=COLORS["sky_blue"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x3, box_y + 0.8, box_w, 0.5, "Pratikshya maam",
                font_size=26, font_color=COLORS["white"], bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x3, box_y + 1.4, box_w, 0.5, "Program Manager\nQuestions & enrollment",
                font_size=14, font_color=COLORS["sky_blue"], alignment=PP_ALIGN.CENTER)

    # Bottom section
    add_rounded_box(slide, 2.0, 5.8, 9.33, 0.8, COLORS["gold"])
    add_textbox(slide, 2.0, 5.85, 9.33, 0.7,
                "Every Saturday  •  9 AM - 1 PM Pacific  •  Ages 10-18  •  No Experience Needed",
                font_size=17, font_color=COLORS["charcoal"], bold=True, alignment=PP_ALIGN.CENTER)

    # Bottom accent
    add_accent_line(slide, 0.5, 6.9, 12.33, COLORS["orange"], 0.05)

    add_textbox(slide, 1.0, 6.95, 11.33, 0.4,
                "EverestITT — Empowering Future Innovators",
                font_size=14, font_color=COLORS["medium_gray"], alignment=PP_ALIGN.CENTER)


# ─── MAIN ────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    create_title_slide(prs)          # 1. Title
    create_why_ai_slide(prs)         # 2. Why AI matters
    create_team_slide(prs)           # 3. Meet the team
    create_roadmap_slide(prs)        # 4. Year roadmap
    create_ai_deep_dive_slide(prs)   # 5. AI deep dive
    create_outcomes_slide(prs)       # 6. What they'll be able to do
    create_differentiators_slide(prs) # 7. Why we're different
    create_hangout_slide(prs)        # 8. Daily hangout sessions
    create_enroll_slide(prs)         # 9. Enroll now

    script_dir = os.path.dirname(os.path.abspath(__file__))
    parent_dir = os.path.dirname(script_dir)
    output_path = os.path.join(parent_dir, "EverestITT_Parent_Presentation.pptx")
    prs.save(output_path)
    print(f"[SUCCESS] Presentation created: {output_path}")


if __name__ == "__main__":
    main()
