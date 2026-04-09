"""
PowerPoint Generator for: Prompt Engineering & Using AI Tools
Kids Computer Science Class - April 11, 2026

Uses the standard theme from tools/create_theme.py

Run:
  pip install python-pptx
  python create_presentation.py
"""

import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', 'tools'))
from create_theme import *

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


# Theme-specific colors
COLORS["ai_purple"] = RGBColor(103, 58, 183)
COLORS["prompt_green"] = RGBColor(56, 142, 60)
COLORS["warning_amber"] = RGBColor(255, 160, 0)
COLORS["deep_purple"] = RGBColor(74, 20, 140)


# ============================================================================
# SLIDE FUNCTIONS
# ============================================================================

def create_slide_1_title(prs):
    """Slide 1: Title Slide"""
    create_title_slide(
        prs,
        "Prompt Engineering &\nUsing AI Tools",
        "Learning to Talk to AI Effectively",
        "April 11, 2026",
        tagline="Week 2 of the AI Phase  |  Kids Computer Science Class"
    )


def create_slide_2_agenda(prs):
    """Slide 2: Agenda / Today's Journey"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Today's Journey")

    items = [
        ("Quick Recap: Last Week's AI Foundations", COLORS["light_blue"]),
        ("What is a Prompt?", COLORS["ai_purple"]),
        ("Why Prompts Matter", COLORS["orange"]),
        ("The Prompting Framework", COLORS["prompt_green"]),
        ("Good vs Bad Prompts", COLORS["red"]),
        ("Prompting Techniques", COLORS["teal"]),
        ("AI Limitations & Safety", COLORS["warning_amber"]),
        ("Hands-On Activity!", COLORS["green"]),
    ]

    for i, (text, color) in enumerate(items):
        add_agenda_item(slide, i + 1, text, 1.6 + i * 0.65, color)


def create_slide_3_recap(prs):
    """Slide 3: Quick Recap from Last Week"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Quick Recap: Last Week's AI Foundations")

    box_w = Inches(5.6)
    box_h = Inches(1.8)
    left_x = Inches(0.7)
    right_x = Inches(6.93)
    row1_y = Inches(1.7)
    row2_y = Inches(3.7)

    boxes = [
        (left_x, row1_y, "LLM = Large Language Model",
         "AI trained on massive text data to predict the next word",
         COLORS["medium_blue"]),
        (right_x, row1_y, "Transformer",
         "Architecture that processes all words at once using attention",
         COLORS["ai_purple"]),
        (left_x, row2_y, "Next Token Prediction",
         "The core process: predicting the most likely next word",
         COLORS["teal"]),
        (right_x, row2_y, "AI Agents",
         "AI that can plan, use tools, and take actions autonomously",
         COLORS["orange"]),
    ]

    for (x, y, title, desc, color) in boxes:
        add_rounded_box(slide, x, y, box_w, box_h, color)
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.2), box_w - Inches(0.4), Inches(0.5),
            title, font_size=20, font_color=COLORS["white"], bold=True,
            alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.8), box_w - Inches(0.4), Inches(0.8),
            desc, font_size=16, font_color=COLORS["light_gray"],
            alignment=PP_ALIGN.CENTER
        )

    add_rounded_box(
        slide, Inches(1.5), Inches(5.9), Inches(10.33), Inches(0.7),
        COLORS["prompt_green"],
        text="Now let's learn how to actually USE these AI tools!",
        text_size=20, text_color=COLORS["white"]
    )


def create_slide_4_what_is_prompt(prs):
    """Slide 4: What is a Prompt?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "What is a Prompt?")

    add_rounded_box(
        slide, Inches(1.0), Inches(1.6), Inches(11.33), Inches(0.9),
        COLORS["ai_purple"],
        text="A prompt is the text instruction you give to an AI to get a response",
        text_size=22, text_color=COLORS["white"]
    )

    examples = [
        ("Simple:", "\"What is the weather?\"",
         COLORS["light_gray"], COLORS["dark_gray"], None),
        ("Better:", "\"What's the weather forecast for Portland, Oregon this weekend?\"",
         COLORS["sky_blue"], COLORS["dark_blue"], None),
        ("Best:", "\"Give me a 3-day weather forecast for Portland, Oregon starting Friday. Include highs, lows, and rain chance.\"",
         COLORS["prompt_green"], COLORS["white"], COLORS["green"]),
    ]

    start_y = 2.9
    for i, (label, example, bg_color, txt_color, border) in enumerate(examples):
        y = start_y + i * 1.25
        add_rounded_box(
            slide, Inches(1.0), Inches(y), Inches(11.33), Inches(1.0),
            bg_color, border_color=border
        )
        add_styled_textbox(
            slide, Inches(1.3), Inches(y + 0.1), Inches(1.5), Inches(0.4),
            label, font_size=18, font_color=txt_color, bold=True
        )
        add_styled_textbox(
            slide, Inches(2.8), Inches(y + 0.1), Inches(9.2), Inches(0.8),
            example, font_size=16, font_color=txt_color
        )

    add_styled_textbox(
        slide, Inches(1.0), Inches(6.5), Inches(11.33), Inches(0.5),
        "The more specific your prompt, the better the AI's response!",
        font_size=20, font_color=COLORS["orange"], bold=True,
        alignment=PP_ALIGN.CENTER
    )


def create_slide_5_why_prompts_matter(prs):
    """Slide 5: Why Prompts Matter"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Why Prompts Matter: Same AI, Different Results")

    col_w = Inches(5.8)
    col_h = Inches(4.5)
    left_x = Inches(0.7)
    right_x = Inches(6.93)
    top_y = Inches(1.7)

    # Left - Vague
    add_rounded_box(slide, left_x, top_y, col_w, col_h,
                    COLORS["light_gray"], border_color=COLORS["red"])
    add_styled_textbox(
        slide, left_x + Inches(0.2), top_y + Inches(0.15),
        col_w - Inches(0.4), Inches(0.5),
        "Vague Prompt", font_size=24, font_color=COLORS["red"], bold=True,
        alignment=PP_ALIGN.CENTER
    )
    add_rounded_box(
        slide, left_x + Inches(0.3), top_y + Inches(0.8),
        col_w - Inches(0.6), Inches(0.7), COLORS["red"],
        text="\"Tell me about dogs\"", text_size=16, text_color=COLORS["white"]
    )
    add_styled_textbox(
        slide, left_x + Inches(0.3), top_y + Inches(1.8),
        col_w - Inches(0.6), Inches(2.2),
        "Result: A generic, unfocused response that covers random dog facts. "
        "Not useful for any specific purpose.",
        font_size=16, font_color=COLORS["dark_gray"],
        alignment=PP_ALIGN.CENTER
    )

    # Right - Specific
    add_rounded_box(slide, right_x, top_y, col_w, col_h,
                    COLORS["light_gray"], border_color=COLORS["green"])
    add_styled_textbox(
        slide, right_x + Inches(0.2), top_y + Inches(0.15),
        col_w - Inches(0.4), Inches(0.5),
        "Specific Prompt", font_size=24, font_color=COLORS["green"], bold=True,
        alignment=PP_ALIGN.CENTER
    )
    add_rounded_box(
        slide, right_x + Inches(0.3), top_y + Inches(0.8),
        col_w - Inches(0.6), Inches(0.7), COLORS["green"],
        text="\"Explain the top 3 most popular dog breeds for families with kids, including size and temperament\"",
        text_size=14, text_color=COLORS["white"]
    )
    add_styled_textbox(
        slide, right_x + Inches(0.3), top_y + Inches(1.8),
        col_w - Inches(0.6), Inches(2.2),
        "Result: A focused, helpful response with exactly what you need -- "
        "breed names, sizes, and temperaments.",
        font_size=16, font_color=COLORS["dark_gray"],
        alignment=PP_ALIGN.CENTER
    )

    add_takeaway_bar(slide,
        "The AI is the same -- YOUR prompt determines the quality of the answer!")


def create_slide_6_rtcf_framework(prs):
    """Slide 6: The Prompting Framework (RTCF)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "The Prompting Framework: R.T.C.F.")

    box_w = Inches(2.85)
    box_h = Inches(3.0)
    start_x = 0.55
    spacing = 3.1
    top_y = Inches(1.7)

    items = [
        ("R", "Role", "Who should the AI be?\n(teacher, scientist, editor)",
         COLORS["ai_purple"]),
        ("T", "Task", "What do you want it to do?\n(explain, write, compare)",
         COLORS["prompt_green"]),
        ("C", "Context", "What background info\ndoes it need?",
         COLORS["orange"]),
        ("F", "Format", "How should the answer look?\n(list, paragraph, table)",
         COLORS["teal"]),
    ]

    for i, (letter, label, desc, color) in enumerate(items):
        x = Inches(start_x + i * spacing)
        add_rounded_box(slide, x, top_y, box_w, box_h, color)
        add_styled_textbox(
            slide, x, top_y + Inches(0.2), box_w, Inches(0.7),
            letter, font_size=40, font_color=COLORS["white"], bold=True,
            alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold"
        )
        add_styled_textbox(
            slide, x, top_y + Inches(0.85), box_w, Inches(0.5),
            "= " + label, font_size=22, font_color=COLORS["white"], bold=True,
            alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x + Inches(0.15), top_y + Inches(1.5),
            box_w - Inches(0.3), Inches(1.2),
            desc, font_size=14, font_color=COLORS["light_gray"],
            alignment=PP_ALIGN.CENTER
        )

    # Example at bottom
    add_rounded_box(
        slide, Inches(0.55), Inches(5.0), Inches(12.23), Inches(1.2),
        COLORS["dark_blue"]
    )
    add_styled_textbox(
        slide, Inches(0.8), Inches(5.05), Inches(11.73), Inches(0.4),
        "Example using all four:", font_size=16, font_color=COLORS["orange"],
        bold=True
    )
    add_styled_textbox(
        slide, Inches(0.8), Inches(5.45), Inches(11.73), Inches(0.7),
        "\"You are a history teacher (Role). Explain the causes of WW2 (Task). "
        "My students are in 8th grade (Context). Use 5 bullet points (Format).\"",
        font_size=15, font_color=COLORS["white"]
    )

    add_styled_textbox(
        slide, Inches(1.0), Inches(6.5), Inches(11.33), Inches(0.4),
        "Remember: The more parts of RTCF you include, the better the response!",
        font_size=16, font_color=COLORS["orange"], bold=True,
        alignment=PP_ALIGN.CENTER
    )


def create_slide_7_rtcf_in_action(prs):
    """Slide 7: RTCF in Action"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Putting R.T.C.F. Together")

    parts = [
        ("Role:", "\"You are a friendly science teacher for 8th graders.\"",
         COLORS["ai_purple"]),
        ("Task:", "\"Explain how photosynthesis works.\"",
         COLORS["prompt_green"]),
        ("Context:", "\"The students have already learned about plant cells.\"",
         COLORS["orange"]),
        ("Format:", "\"Use 3 bullet points with simple language. Include an analogy.\"",
         COLORS["teal"]),
    ]

    start_y = 1.7
    for i, (label, text, color) in enumerate(parts):
        y = Inches(start_y + i * 0.6)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.7), y, Inches(0.15), Inches(0.45)
        )
        set_shape_fill(bar, color)
        bar.line.fill.background()

        add_styled_textbox(
            slide, Inches(1.0), y, Inches(1.5), Inches(0.45),
            label, font_size=18, font_color=color, bold=True
        )
        add_styled_textbox(
            slide, Inches(2.5), y, Inches(10.0), Inches(0.45),
            text, font_size=16, font_color=COLORS["dark_gray"]
        )

    add_rounded_box(
        slide, Inches(0.7), Inches(4.2), Inches(11.93), Inches(1.8),
        COLORS["dark_blue"]
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(4.3), Inches(11.33), Inches(0.4),
        "Combined Prompt:", font_size=18, font_color=COLORS["orange"], bold=True
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(4.75), Inches(11.33), Inches(1.1),
        "\"You are a friendly science teacher for 8th graders. Explain how "
        "photosynthesis works. The students have already learned about plant cells. "
        "Use 3 bullet points with simple language. Include an analogy.\"",
        font_size=15, font_color=COLORS["white"]
    )

    add_rounded_box(
        slide, Inches(1.5), Inches(6.3), Inches(10.33), Inches(0.7),
        COLORS["prompt_green"],
        text="You don't need ALL four every time -- but the more you include, the better!",
        text_size=18, text_color=COLORS["white"]
    )


def create_slide_8_good_vs_bad(prs):
    """Slide 8: Good vs Bad Prompts"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Good vs Bad Prompts: Side by Side")

    bad_w = Inches(5.0)
    good_w = Inches(7.0)
    bad_x = Inches(0.5)
    good_x = Inches(5.83)
    row_h = Inches(1.5)

    rows = [
        ("\"Write an essay\"",
         "\"Write a 5-paragraph essay about why recycling is important for a 7th grade science class\""),
        ("\"Fix this code\"",
         "\"Find the bug in this Python function that should add two numbers but returns None\""),
        ("\"Help with homework\"",
         "\"Explain how to solve quadratic equations step by step with an example using x\u00b2 + 5x + 6 = 0\""),
    ]

    add_styled_textbox(
        slide, bad_x, Inches(1.5), bad_w, Inches(0.4),
        "Bad Prompt", font_size=20, font_color=COLORS["red"], bold=True,
        alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, good_x, Inches(1.5), good_w, Inches(0.4),
        "Good Prompt", font_size=20, font_color=COLORS["green"], bold=True,
        alignment=PP_ALIGN.CENTER
    )

    start_y = 2.0
    for i, (bad, good) in enumerate(rows):
        y = Inches(start_y + i * 1.6)

        add_rounded_box(slide, bad_x, y, bad_w, row_h,
                        COLORS["light_gray"], border_color=COLORS["red"])
        add_styled_textbox(
            slide, bad_x + Inches(0.15), y + Inches(0.15),
            bad_w - Inches(0.3), row_h - Inches(0.3),
            bad, font_size=15, font_color=COLORS["red"], bold=True,
            alignment=PP_ALIGN.CENTER
        )

        add_styled_textbox(
            slide, Inches(5.4), y + Inches(0.3), Inches(0.5), Inches(0.5),
            "\u2192", font_size=28, font_color=COLORS["orange"], bold=True,
            alignment=PP_ALIGN.CENTER
        )

        add_rounded_box(slide, good_x, y, good_w, row_h,
                        COLORS["light_gray"], border_color=COLORS["green"])
        add_styled_textbox(
            slide, good_x + Inches(0.15), y + Inches(0.15),
            good_w - Inches(0.3), row_h - Inches(0.3),
            good, font_size=15, font_color=COLORS["prompt_green"], bold=True,
            alignment=PP_ALIGN.CENTER
        )

    add_styled_textbox(
        slide, Inches(1.0), Inches(6.9), Inches(11.33), Inches(0.4),
        "Specificity is the key to great AI responses!",
        font_size=18, font_color=COLORS["orange"], bold=True,
        alignment=PP_ALIGN.CENTER
    )


def create_slide_9_prompting_techniques(prs):
    """Slide 9: Prompting Techniques"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Prompting Techniques")

    techniques = [
        ("Zero-Shot", "Just ask directly. No examples needed.",
         "Example: \"Translate 'hello' to Spanish\"",
         COLORS["medium_blue"]),
        ("Few-Shot", "Give the AI a few examples first so it learns the pattern.",
         "Example: \"Happy = positive, Sad = negative, Excited = ?\"",
         COLORS["ai_purple"]),
        ("Chain-of-Thought", "Ask the AI to think step by step.",
         "Example: \"Solve this math problem. Show your work step by step.\"",
         COLORS["teal"]),
    ]

    start_y = 1.7
    box_h = Inches(1.5)
    spacing = 1.7

    for i, (name, desc, example, color) in enumerate(techniques):
        y = Inches(start_y + i * spacing)
        add_rounded_box(slide, Inches(0.7), y, Inches(11.93), box_h, color)
        add_styled_textbox(
            slide, Inches(1.0), y + Inches(0.1), Inches(11.33), Inches(0.45),
            name, font_size=24, font_color=COLORS["white"], bold=True
        )
        add_styled_textbox(
            slide, Inches(1.0), y + Inches(0.55), Inches(11.33), Inches(0.4),
            desc, font_size=16, font_color=COLORS["light_gray"]
        )
        add_styled_textbox(
            slide, Inches(1.0), y + Inches(0.95), Inches(11.33), Inches(0.4),
            example, font_size=14, font_color=COLORS["sky_blue"]
        )

    add_styled_textbox(
        slide, Inches(1.0), Inches(6.9), Inches(11.33), Inches(0.4),
        "Remember Few-Shot from last week? GPT-3 surprised everyone by learning from just a few examples!",
        font_size=14, font_color=COLORS["ai_purple"], bold=True,
        alignment=PP_ALIGN.CENTER
    )


def create_slide_10_chain_of_thought(prs):
    """Slide 10: Chain-of-Thought Deep Dive"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Chain-of-Thought: Making AI Think Step by Step")

    col_w = Inches(5.8)
    left_x = Inches(0.7)
    right_x = Inches(6.93)
    top_y = Inches(1.7)
    col_h = Inches(4.3)

    # Left - Without CoT
    add_rounded_box(slide, left_x, top_y, col_w, col_h,
                    COLORS["light_gray"], border_color=COLORS["red"])
    add_styled_textbox(
        slide, left_x + Inches(0.2), top_y + Inches(0.15),
        col_w - Inches(0.4), Inches(0.5),
        "Without Chain-of-Thought", font_size=22, font_color=COLORS["red"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    add_rounded_box(
        slide, left_x + Inches(0.3), top_y + Inches(0.8),
        col_w - Inches(0.6), Inches(0.6), COLORS["red"],
        text="\"What is 15% tip on $42.50?\"", text_size=15,
        text_color=COLORS["white"]
    )
    add_styled_textbox(
        slide, left_x + Inches(0.3), top_y + Inches(1.7),
        col_w - Inches(0.6), Inches(0.5),
        "AI response:", font_size=16, font_color=COLORS["dark_gray"], bold=True
    )
    add_styled_textbox(
        slide, left_x + Inches(0.3), top_y + Inches(2.2),
        col_w - Inches(0.6), Inches(1.5),
        "\"The tip would be $6.50\"\n\n(Might be wrong -- no work shown!)",
        font_size=15, font_color=COLORS["dark_gray"],
        alignment=PP_ALIGN.CENTER
    )

    # Right - With CoT
    add_rounded_box(slide, right_x, top_y, col_w, col_h,
                    COLORS["light_gray"], border_color=COLORS["green"])
    add_styled_textbox(
        slide, right_x + Inches(0.2), top_y + Inches(0.15),
        col_w - Inches(0.4), Inches(0.5),
        "With Chain-of-Thought", font_size=22, font_color=COLORS["green"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    add_rounded_box(
        slide, right_x + Inches(0.3), top_y + Inches(0.8),
        col_w - Inches(0.6), Inches(0.6), COLORS["green"],
        text="\"What is 15% tip on $42.50? Think step by step.\"",
        text_size=14, text_color=COLORS["white"]
    )
    add_styled_textbox(
        slide, right_x + Inches(0.3), top_y + Inches(1.7),
        col_w - Inches(0.6), Inches(0.5),
        "AI response:", font_size=16, font_color=COLORS["dark_gray"], bold=True
    )
    add_styled_textbox(
        slide, right_x + Inches(0.3), top_y + Inches(2.2),
        col_w - Inches(0.6), Inches(1.8),
        "Step 1: 10% of $42.50 = $4.25\n"
        "Step 2: 5% = $4.25 / 2 = $2.125\n"
        "Step 3: 15% = $4.25 + $2.125 = $6.38",
        font_size=15, font_color=COLORS["dark_gray"]
    )

    add_takeaway_bar(slide,
        "Adding 'think step by step' or 'explain your reasoning' makes AI much more accurate!",
        COLORS["teal"])


def create_slide_11_hallucinations(prs):
    """Slide 11: AI Limitations - Hallucinations"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "AI Limitations: Hallucinations")

    add_rounded_box(
        slide, Inches(0.7), Inches(1.6), Inches(11.93), Inches(1.0),
        COLORS["deep_purple"],
        text="Hallucination = When AI generates false information but presents it confidently",
        text_size=20, text_color=COLORS["white"]
    )

    examples = [
        ("AI might invent fake book titles that don't exist", COLORS["red"]),
        ("AI might cite made-up statistics or research papers", COLORS["orange"]),
        ("AI might give wrong answers to math or logic problems", COLORS["warning_amber"]),
    ]

    for i, (text, color) in enumerate(examples):
        y = Inches(3.0 + i * 1.1)
        add_rounded_box(
            slide, Inches(1.5), y, Inches(10.33), Inches(0.85), color,
            text=text, text_size=18, text_color=COLORS["white"]
        )

    add_takeaway_bar(slide,
        "ALWAYS fact-check important AI responses! AI is confident, not always correct.",
        COLORS["prompt_green"])


def create_slide_12_more_limitations(prs):
    """Slide 12: AI Limitations - Other Issues"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "More AI Limitations to Know")

    box_w = Inches(5.6)
    box_h = Inches(1.8)
    left_x = Inches(0.7)
    right_x = Inches(6.93)
    row1_y = Inches(1.7)
    row2_y = Inches(3.7)

    boxes = [
        (left_x, row1_y, "Knowledge Cutoff",
         "AI doesn't know about recent events after its training date",
         COLORS["medium_blue"]),
        (right_x, row1_y, "Bias",
         "AI can reflect biases from its training data",
         COLORS["orange"]),
        (left_x, row2_y, "No Real Understanding",
         "AI predicts words -- it doesn't truly 'understand' like humans",
         COLORS["ai_purple"]),
        (right_x, row2_y, "Context Window",
         "AI can only process a limited amount of text at once",
         COLORS["teal"]),
    ]

    for (x, y, title, desc, color) in boxes:
        add_rounded_box(slide, x, y, box_w, box_h, color)
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.2), box_w - Inches(0.4), Inches(0.5),
            title, font_size=22, font_color=COLORS["white"], bold=True,
            alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.85), box_w - Inches(0.4), Inches(0.7),
            desc, font_size=16, font_color=COLORS["light_gray"],
            alignment=PP_ALIGN.CENTER
        )

    add_rounded_box(
        slide, Inches(1.5), Inches(5.9), Inches(10.33), Inches(0.7),
        COLORS["dark_blue"],
        text="Understanding limitations makes you a SMARTER AI user!",
        text_size=20, text_color=COLORS["white"]
    )


def create_slide_13_ai_safety(prs):
    """Slide 13: AI Safety & Responsible Use"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Using AI Safely & Responsibly")

    col_w = Inches(5.8)
    col_h = Inches(4.0)
    left_x = Inches(0.7)
    right_x = Inches(6.93)
    top_y = Inches(1.7)

    # DO's
    add_rounded_box(slide, left_x, top_y, col_w, col_h,
                    COLORS["light_gray"], border_color=COLORS["green"])
    add_styled_textbox(
        slide, left_x + Inches(0.2), top_y + Inches(0.15),
        col_w - Inches(0.4), Inches(0.5),
        "DO's", font_size=28, font_color=COLORS["green"], bold=True,
        alignment=PP_ALIGN.CENTER
    )
    dos = [
        "Fact-check important information",
        "Use AI as a learning helper, not a replacement",
        "Tell teachers when you use AI for assignments",
        "Protect your personal information",
    ]
    for i, item in enumerate(dos):
        add_styled_textbox(
            slide, left_x + Inches(0.4), top_y + Inches(0.8 + i * 0.7),
            col_w - Inches(0.8), Inches(0.5),
            "\u2713  " + item, font_size=16, font_color=COLORS["prompt_green"],
            bold=True
        )

    # DON'T's
    add_rounded_box(slide, right_x, top_y, col_w, col_h,
                    COLORS["light_gray"], border_color=COLORS["red"])
    add_styled_textbox(
        slide, right_x + Inches(0.2), top_y + Inches(0.15),
        col_w - Inches(0.4), Inches(0.5),
        "DON'T's", font_size=28, font_color=COLORS["red"], bold=True,
        alignment=PP_ALIGN.CENTER
    )
    donts = [
        "Don't share personal info (address, passwords)",
        "Don't copy AI output and claim it as your own",
        "Don't trust AI blindly -- verify important facts",
        "Don't use AI to cheat on tests or assignments",
    ]
    for i, item in enumerate(donts):
        add_styled_textbox(
            slide, right_x + Inches(0.4), top_y + Inches(0.8 + i * 0.7),
            col_w - Inches(0.8), Inches(0.5),
            "\u2717  " + item, font_size=16, font_color=COLORS["red"], bold=True
        )

    add_takeaway_bar(slide,
        "AI is a powerful TOOL -- use it wisely and honestly!")


def create_slide_14_ai_tools(prs):
    """Slide 14: AI Tools You Can Try"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "AI Tools You Can Try")

    box_w = Inches(5.6)
    box_h = Inches(1.8)
    left_x = Inches(0.7)
    right_x = Inches(6.93)
    row1_y = Inches(1.7)
    row2_y = Inches(3.7)

    tools = [
        (left_x, row1_y, "ChatGPT", "by OpenAI",
         "Most popular chatbot, free tier available", COLORS["green"]),
        (right_x, row1_y, "Claude", "by Anthropic",
         "Focused on being helpful and safe", COLORS["ai_purple"]),
        (left_x, row2_y, "Gemini", "by Google",
         "Integrated with Google services", COLORS["medium_blue"]),
        (right_x, row2_y, "Copilot", "by Microsoft",
         "Built into Windows, Edge, and Office", COLORS["teal"]),
    ]

    for (x, y, name, company, desc, color) in tools:
        add_rounded_box(slide, x, y, box_w, box_h, color)
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.2), box_w - Inches(0.4), Inches(0.5),
            name, font_size=26, font_color=COLORS["white"], bold=True,
            alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.7), box_w - Inches(0.4), Inches(0.4),
            company, font_size=16, font_color=COLORS["light_gray"],
            alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(1.15), box_w - Inches(0.4), Inches(0.5),
            desc, font_size=15, font_color=COLORS["sky_blue"],
            alignment=PP_ALIGN.CENTER
        )

    add_takeaway_bar(slide,
        "All of these use LLMs -- the concepts you learned last week power these tools!")


def create_slide_15_key_takeaways(prs):
    """Slide 15: Key Takeaways"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Today's Key Takeaways")

    takeaways = [
        ("A prompt is how you talk to AI -- wording matters!", COLORS["ai_purple"]),
        ("Use R.T.C.F.: Role, Task, Context, Format", COLORS["prompt_green"]),
        ("Be specific: vague prompts = vague answers", COLORS["orange"]),
        ("Chain-of-thought: 'think step by step' improves accuracy", COLORS["teal"]),
        ("AI can hallucinate -- always fact-check", COLORS["red"]),
        ("Use AI responsibly -- it's a tool, not a shortcut", COLORS["medium_blue"]),
    ]

    for i, (text, color) in enumerate(takeaways):
        y = 1.7 + i * 0.75
        add_circle(slide, Inches(1.0), Inches(y), Inches(0.55), color,
                   text=str(i + 1), text_size=18)
        add_styled_textbox(
            slide, Inches(1.8), Inches(y + 0.05), Inches(10.5), Inches(0.5),
            text, font_size=22, font_color=COLORS["dark_gray"], bold=True
        )

    add_rounded_box(
        slide, Inches(2.5), Inches(6.4), Inches(8.33), Inches(0.7),
        COLORS["orange"],
        text="Time for the activity and Kahoot!",
        text_size=22, text_color=COLORS["white"]
    )


def create_slide_16_questions(prs):
    """Slide 16: Questions Slide"""
    create_questions_slide(prs, "Let's practice prompting!")


# ============================================================================
# MAIN
# ============================================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_slide_1_title(prs)
    create_slide_2_agenda(prs)
    create_slide_3_recap(prs)
    create_slide_4_what_is_prompt(prs)
    create_slide_5_why_prompts_matter(prs)
    create_slide_6_rtcf_framework(prs)
    create_slide_7_rtcf_in_action(prs)
    create_slide_8_good_vs_bad(prs)
    create_slide_9_prompting_techniques(prs)
    create_slide_10_chain_of_thought(prs)
    create_slide_11_hallucinations(prs)
    create_slide_12_more_limitations(prs)
    create_slide_13_ai_safety(prs)
    create_slide_14_ai_tools(prs)
    create_slide_15_key_takeaways(prs)
    create_slide_16_questions(prs)

    script_dir = os.path.dirname(os.path.abspath(__file__))
    parent_dir = os.path.dirname(script_dir)
    output_path = os.path.join(parent_dir, "2026-04-11_Prompt_Engineering_and_Using_AI_Tools_BEAUTIFUL.pptx")
    prs.save(output_path)
    print(f"[SUCCESS] Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
