"""
PowerPoint Generator for: Meta Prompting & AI Team Challenge
Kids Computer Science Class - April 18, 2026

Uses the standard theme from tools/create_theme.py

Run:
  pip install python-pptx
  python create_presentation.py
"""

import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', 'tools'))
from create_theme import *

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


# Theme-specific colors
COLORS["ai_purple"] = RGBColor(103, 58, 183)
COLORS["prompt_green"] = RGBColor(56, 142, 60)
COLORS["meta_gold"] = RGBColor(255, 179, 0)
COLORS["deep_purple"] = RGBColor(74, 20, 140)
COLORS["team_cyan"] = RGBColor(0, 172, 193)


# ============================================================================
# SLIDE FUNCTIONS
# ============================================================================

def create_slide_1_title(prs):
    """Slide 1: Title Slide"""
    create_title_slide(
        prs,
        "Meta Prompting &\nAI Team Challenge",
        "Using AI to Write Better Prompts + Hands-On Group Projects",
        "April 18, 2026",
        tagline="Week 3 of the AI Phase  |  Kids Computer Science Class"
    )


def create_slide_2_agenda(prs):
    """Slide 2: Agenda / Today's Journey"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Today's Journey")

    items = [
        ("Quick Recap: Weeks 1-2", COLORS["light_blue"]),
        ("What is Meta Prompting?", COLORS["meta_gold"]),
        ("Meta Prompting Techniques", COLORS["ai_purple"]),
        ("Live Demo: Meta Prompting in Action", COLORS["prompt_green"]),
        ("AI Tools You Can Use (Free!)", COLORS["teal"]),
        ("Group Challenge: AI Team Battle!", COLORS["team_cyan"]),
        ("Presentations & Voting", COLORS["orange"]),
        ("Kahoot!", COLORS["green"]),
    ]

    for i, (text, color) in enumerate(items):
        add_agenda_item(slide, i + 1, text, 1.6 + i * 0.65, color)


def create_slide_3_recap(prs):
    """Slide 3: Quick Recap from Weeks 1-2"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Quick Recap: What We've Learned So Far")

    box_w = Inches(3.8)
    box_h = Inches(1.5)
    start_x = 0.55
    spacing = 4.1
    row1_y = Inches(1.7)
    row2_y = Inches(3.4)

    boxes = [
        (Inches(start_x), row1_y, "LLMs & Transformers",
         "AI trained on text data using\nattention to predict next words",
         COLORS["medium_blue"]),
        (Inches(start_x + spacing), row1_y, "RTCF Framework",
         "Role, Task, Context, Format\nfor writing great prompts",
         COLORS["ai_purple"]),
        (Inches(start_x + spacing * 2), row1_y, "Prompting Techniques",
         "Zero-shot, Few-shot,\nChain-of-Thought",
         COLORS["teal"]),
        (Inches(start_x), row2_y, "AI Limitations",
         "Hallucinations, bias,\nknowledge cutoff",
         COLORS["orange"]),
        (Inches(start_x + spacing), row2_y, "AI Safety",
         "Fact-check, don't share\npersonal info, be honest",
         COLORS["prompt_green"]),
        (Inches(start_x + spacing * 2), row2_y, "AI Tools",
         "ChatGPT, Gemini, Grok\nand more",
         COLORS["team_cyan"]),
    ]

    for (x, y, title, desc, color) in boxes:
        add_rounded_box(slide, x, y, box_w, box_h, color)
        add_styled_textbox(
            slide, x + Inches(0.15), y + Inches(0.15), box_w - Inches(0.3), Inches(0.4),
            title, font_size=18, font_color=COLORS["white"], bold=True,
            alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x + Inches(0.15), y + Inches(0.6), box_w - Inches(0.3), Inches(0.7),
            desc, font_size=14, font_color=COLORS["light_gray"],
            alignment=PP_ALIGN.CENTER
        )

    add_rounded_box(
        slide, Inches(1.5), Inches(5.3), Inches(10.33), Inches(0.7),
        COLORS["meta_gold"],
        text="Today: Level up your prompting with META PROMPTING + team challenges!",
        text_size=20, text_color=COLORS["dark_blue"]
    )


def create_slide_4_what_is_meta(prs):
    """Slide 4: What is Meta Prompting?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "What is Meta Prompting?")

    # Big definition box
    add_rounded_box(
        slide, Inches(0.7), Inches(1.6), Inches(11.93), Inches(1.2),
        COLORS["deep_purple"]
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(1.7), Inches(11.33), Inches(0.5),
        "Meta Prompting = Using AI to help you write better prompts",
        font_size=24, font_color=COLORS["white"], bold=True,
        alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(2.2), Inches(11.33), Inches(0.4),
        "Instead of writing the perfect prompt yourself, ask AI to help you create it!",
        font_size=16, font_color=COLORS["light_gray"],
        alignment=PP_ALIGN.CENTER
    )

    # Analogy
    add_rounded_box(
        slide, Inches(0.7), Inches(3.2), Inches(5.6), Inches(3.0),
        COLORS["light_gray"], border_color=COLORS["ai_purple"]
    )
    add_styled_textbox(
        slide, Inches(0.9), Inches(3.35), Inches(5.2), Inches(0.4),
        "Think of it like this...", font_size=20, font_color=COLORS["ai_purple"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(0.9), Inches(3.9), Inches(5.2), Inches(2.0),
        "Regular prompting:\nYou write a prompt and hope it works.\n\n"
        "Meta prompting:\nYou ask AI to help you write the prompt,\n"
        "THEN you use that better prompt!",
        font_size=15, font_color=COLORS["dark_gray"],
        alignment=PP_ALIGN.CENTER
    )

    # Why it matters
    add_rounded_box(
        slide, Inches(6.93), Inches(3.2), Inches(5.7), Inches(3.0),
        COLORS["meta_gold"]
    )
    add_styled_textbox(
        slide, Inches(7.13), Inches(3.35), Inches(5.3), Inches(0.4),
        "Why does this matter?", font_size=20, font_color=COLORS["dark_blue"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(7.13), Inches(3.9), Inches(5.3), Inches(2.0),
        "- AI is great at writing prompts!\n"
        "- You get better results with less effort\n"
        "- It helps when you're stuck\n"
        "- It's how professionals use AI",
        font_size=15, font_color=COLORS["dark_blue"],
        alignment=PP_ALIGN.LEFT
    )

    add_takeaway_bar(slide,
        "Meta prompting is the #1 skill that separates beginners from experts!",
        COLORS["deep_purple"])


def create_slide_5_technique_1(prs):
    """Slide 5: Meta Prompting Technique 1 - Ask AI to Write Your Prompt"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Technique 1: Ask AI to Write Your Prompt")

    # Left side - what you say
    add_rounded_box(
        slide, Inches(0.7), Inches(1.7), Inches(5.6), Inches(2.2),
        COLORS["light_gray"], border_color=COLORS["red"]
    )
    add_styled_textbox(
        slide, Inches(0.9), Inches(1.85), Inches(5.2), Inches(0.4),
        "What you might normally type:", font_size=18, font_color=COLORS["red"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    add_rounded_box(
        slide, Inches(1.0), Inches(2.4), Inches(5.0), Inches(0.6),
        COLORS["red"],
        text="\"Tell me about black holes\"", text_size=16, text_color=COLORS["white"]
    )
    add_styled_textbox(
        slide, Inches(0.9), Inches(3.2), Inches(5.2), Inches(0.4),
        "Generic, unfocused result", font_size=14, font_color=COLORS["dark_gray"],
        alignment=PP_ALIGN.CENTER
    )

    # Right side - meta prompt
    add_rounded_box(
        slide, Inches(6.93), Inches(1.7), Inches(5.7), Inches(2.2),
        COLORS["light_gray"], border_color=COLORS["prompt_green"]
    )
    add_styled_textbox(
        slide, Inches(7.13), Inches(1.85), Inches(5.3), Inches(0.4),
        "The meta prompting way:", font_size=18, font_color=COLORS["prompt_green"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    add_rounded_box(
        slide, Inches(7.23), Inches(2.4), Inches(5.1), Inches(0.6),
        COLORS["prompt_green"],
        text="\"Write the best possible prompt for learning about black holes\"",
        text_size=14, text_color=COLORS["white"]
    )
    add_styled_textbox(
        slide, Inches(7.13), Inches(3.2), Inches(5.3), Inches(0.4),
        "AI creates a detailed, structured prompt for you!",
        font_size=14, font_color=COLORS["dark_gray"],
        alignment=PP_ALIGN.CENTER
    )

    # Arrow between
    add_styled_textbox(
        slide, Inches(6.1), Inches(2.5), Inches(1.0), Inches(0.5),
        "\u2192", font_size=36, font_color=COLORS["meta_gold"], bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # Example result
    add_rounded_box(
        slide, Inches(0.7), Inches(4.3), Inches(11.93), Inches(2.0),
        COLORS["dark_blue"]
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(4.4), Inches(11.33), Inches(0.4),
        "AI might generate a prompt like:", font_size=16, font_color=COLORS["meta_gold"],
        bold=True
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(4.85), Inches(11.33), Inches(1.3),
        "\"You are an astrophysics professor explaining to curious teenagers. Explain how "
        "black holes form, what happens at the event horizon, and why they matter for "
        "understanding the universe. Use analogies, include 5 key facts, and end with "
        "the most mind-blowing thing about black holes.\"",
        font_size=15, font_color=COLORS["white"]
    )


def create_slide_6_technique_2(prs):
    """Slide 6: Meta Prompting Technique 2 - Ask AI to Improve Your Prompt"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Technique 2: Ask AI to Improve Your Prompt")

    # Step-by-step flow
    step_colors = [COLORS["red"], COLORS["meta_gold"], COLORS["prompt_green"]]
    step_labels = ["Step 1: Your rough prompt", "Step 2: Ask AI to improve it", "Step 3: Use the improved prompt"]
    step_contents = [
        "\"Explain photosynthesis for my science class\"",
        "\"Make this prompt better using the RTCF framework:\nExplain photosynthesis for my science class\"",
        "AI rewrites it with Role, Task, Context, and Format!"
    ]

    box_w = Inches(11.93)
    box_h = Inches(1.3)

    for i in range(3):
        y = Inches(1.7 + i * 1.55)
        add_rounded_box(slide, Inches(0.7), y, box_w, box_h, step_colors[i])
        add_styled_textbox(
            slide, Inches(1.0), y + Inches(0.1), Inches(11.33), Inches(0.35),
            step_labels[i], font_size=18, font_color=COLORS["white"], bold=True
        )
        add_styled_textbox(
            slide, Inches(1.0), y + Inches(0.5), Inches(11.33), Inches(0.7),
            step_contents[i], font_size=15, font_color=COLORS["light_gray"]
        )

        # Arrow between steps
        if i < 2:
            arrow_y = Inches(1.7 + i * 1.55 + 1.3)
            add_styled_textbox(
                slide, Inches(6.0), arrow_y, Inches(1.33), Inches(0.3),
                "\u2193", font_size=24, font_color=COLORS["meta_gold"], bold=True,
                alignment=PP_ALIGN.CENTER
            )

    add_takeaway_bar(slide,
        "You don't need to write the perfect prompt -- let AI help you perfect it!",
        COLORS["deep_purple"])


def create_slide_7_technique_3(prs):
    """Slide 7: Meta Prompting Technique 3 - Ask AI What Questions to Ask"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Technique 3: Ask AI What Questions to Ask")

    # Problem
    add_rounded_box(
        slide, Inches(0.7), Inches(1.7), Inches(5.6), Inches(1.5),
        COLORS["light_gray"], border_color=COLORS["orange"]
    )
    add_styled_textbox(
        slide, Inches(0.9), Inches(1.85), Inches(5.2), Inches(0.4),
        "The Problem:", font_size=20, font_color=COLORS["orange"], bold=True,
        alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(0.9), Inches(2.3), Inches(5.2), Inches(0.7),
        "Sometimes you don't even know\nwhat to ask! You just know the topic.",
        font_size=16, font_color=COLORS["dark_gray"],
        alignment=PP_ALIGN.CENTER
    )

    # Solution
    add_rounded_box(
        slide, Inches(6.93), Inches(1.7), Inches(5.7), Inches(1.5),
        COLORS["light_gray"], border_color=COLORS["prompt_green"]
    )
    add_styled_textbox(
        slide, Inches(7.13), Inches(1.85), Inches(5.3), Inches(0.4),
        "The Meta Prompt:", font_size=20, font_color=COLORS["prompt_green"], bold=True,
        alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(7.13), Inches(2.3), Inches(5.3), Inches(0.7),
        "\"I'm researching climate change for\na school project. What are the 5 best\nquestions I should ask you?\"",
        font_size=15, font_color=COLORS["dark_gray"],
        alignment=PP_ALIGN.CENTER
    )

    # What AI returns
    add_rounded_box(
        slide, Inches(0.7), Inches(3.5), Inches(11.93), Inches(2.8),
        COLORS["dark_blue"]
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(3.6), Inches(11.33), Inches(0.4),
        "AI might suggest questions like:", font_size=18, font_color=COLORS["meta_gold"],
        bold=True
    )

    questions = [
        "1. What are the top 3 causes of climate change?",
        "2. How does climate change affect weather patterns?",
        "3. What are scientists doing to fight climate change?",
        "4. How does climate change impact animals and ecosystems?",
        "5. What can teenagers do to help reduce climate change?",
    ]
    for i, q in enumerate(questions):
        add_styled_textbox(
            slide, Inches(1.3), Inches(4.1 + i * 0.4), Inches(10.73), Inches(0.35),
            q, font_size=15, font_color=COLORS["white"]
        )

    add_takeaway_bar(slide,
        "When you're stuck, ask AI: \"What should I ask you about [topic]?\"",
        COLORS["meta_gold"])


def create_slide_8_live_demo(prs):
    """Slide 8: Live Demo - Before and After Meta Prompting"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Live Demo: Meta Prompting in Action")

    # Before
    add_rounded_box(
        slide, Inches(0.7), Inches(1.7), Inches(5.6), Inches(4.5),
        COLORS["light_gray"], border_color=COLORS["red"]
    )
    add_styled_textbox(
        slide, Inches(0.9), Inches(1.85), Inches(5.2), Inches(0.4),
        "BEFORE Meta Prompting", font_size=22, font_color=COLORS["red"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    add_rounded_box(
        slide, Inches(1.0), Inches(2.5), Inches(5.0), Inches(0.6),
        COLORS["red"],
        text="\"Help me study for my test\"", text_size=16, text_color=COLORS["white"]
    )
    add_styled_textbox(
        slide, Inches(0.9), Inches(3.4), Inches(5.2), Inches(2.5),
        "Result:\nGeneric study tips that could\napply to any subject.\n\n"
        "Not personalized.\nNot specific.\nNot very helpful.",
        font_size=15, font_color=COLORS["dark_gray"],
        alignment=PP_ALIGN.CENTER
    )

    # Arrow
    add_styled_textbox(
        slide, Inches(6.0), Inches(3.5), Inches(1.2), Inches(0.5),
        "\u2192", font_size=40, font_color=COLORS["meta_gold"], bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # After
    add_rounded_box(
        slide, Inches(6.93), Inches(1.7), Inches(5.7), Inches(4.5),
        COLORS["light_gray"], border_color=COLORS["prompt_green"]
    )
    add_styled_textbox(
        slide, Inches(7.13), Inches(1.85), Inches(5.3), Inches(0.4),
        "AFTER Meta Prompting", font_size=22, font_color=COLORS["prompt_green"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    add_rounded_box(
        slide, Inches(7.23), Inches(2.5), Inches(5.1), Inches(0.6),
        COLORS["prompt_green"],
        text="\"Write me the best prompt for studying for my 8th grade biology test on cells\"",
        text_size=13, text_color=COLORS["white"]
    )
    add_styled_textbox(
        slide, Inches(7.13), Inches(3.4), Inches(5.3), Inches(2.5),
        "Result:\nAI writes a detailed RTCF prompt.\nYou use THAT prompt.\n\n"
        "Personalized study guide.\nPractice questions included.\nExactly what you need!",
        font_size=15, font_color=COLORS["dark_gray"],
        alignment=PP_ALIGN.CENTER
    )

    add_takeaway_bar(slide,
        "Let's try this LIVE right now! Watch the difference in real time.",
        COLORS["team_cyan"])


def create_slide_9_ai_tools(prs):
    """Slide 9: AI Tools You Can Use (Free)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "AI Tools You Can Use (Free!)")

    box_w = Inches(3.8)
    box_h = Inches(3.5)
    start_x = 0.55
    spacing = 4.1
    top_y = Inches(1.7)

    tools = [
        ("ChatGPT", "by OpenAI", "chat.openai.com",
         "Most popular AI chatbot\nFree account with email or Google",
         COLORS["prompt_green"]),
        ("Gemini", "by Google", "gemini.google.com",
         "Works with your Google account\nGreat for research and writing",
         COLORS["medium_blue"]),
        ("Grok", "by xAI", "grok.com",
         "Free tier available\nSign up with email or X account",
         COLORS["dark_gray"]),
    ]

    for i, (name, company, url, desc, color) in enumerate(tools):
        x = Inches(start_x + i * spacing)
        add_rounded_box(slide, x, top_y, box_w, box_h, color)
        add_styled_textbox(
            slide, x + Inches(0.15), top_y + Inches(0.15), box_w - Inches(0.3), Inches(0.5),
            name, font_size=28, font_color=COLORS["white"], bold=True,
            alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x + Inches(0.15), top_y + Inches(0.6), box_w - Inches(0.3), Inches(0.35),
            company, font_size=14, font_color=COLORS["light_gray"],
            alignment=PP_ALIGN.CENTER
        )
        # URL box
        add_rounded_box(
            slide, x + Inches(0.2), top_y + Inches(1.1),
            box_w - Inches(0.4), Inches(0.5),
            COLORS["white"]
        )
        add_styled_textbox(
            slide, x + Inches(0.2), top_y + Inches(1.15),
            box_w - Inches(0.4), Inches(0.4),
            url, font_size=14, font_color=COLORS["dark_blue"], bold=True,
            alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x + Inches(0.15), top_y + Inches(1.8), box_w - Inches(0.3), Inches(1.2),
            desc, font_size=14, font_color=COLORS["light_gray"],
            alignment=PP_ALIGN.CENTER
        )

    # Safety reminder
    add_rounded_box(
        slide, Inches(0.55), Inches(5.5), Inches(12.23), Inches(0.7),
        COLORS["red"],
        text="REMINDER: Never share personal information (passwords, address, phone) with any AI tool!",
        text_size=18, text_color=COLORS["white"]
    )


def create_slide_10_team_roles(prs):
    """Slide 10: Working with AI as a Team"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Working with AI as a Team")

    # Intro text
    add_styled_textbox(
        slide, Inches(0.7), Inches(1.55), Inches(11.93), Inches(0.4),
        "Each team member has a role. Rotate roles between rounds!",
        font_size=18, font_color=COLORS["dark_gray"], bold=True,
        alignment=PP_ALIGN.CENTER
    )

    box_w = Inches(5.6)
    box_h = Inches(1.5)
    left_x = Inches(0.7)
    right_x = Inches(6.93)
    row1_y = Inches(2.2)
    row2_y = Inches(3.9)

    roles = [
        (left_x, row1_y, "Prompt Writer",
         "Types the prompts into the AI tool.\nTakes suggestions from the team.",
         COLORS["ai_purple"]),
        (right_x, row1_y, "Fact Checker",
         "Verifies AI responses are accurate.\nSearches for evidence online.",
         COLORS["red"]),
        (left_x, row2_y, "Note Taker",
         "Documents prompts, responses, and\nwhat the team changed each time.",
         COLORS["teal"]),
        (right_x, row2_y, "Idea Generator",
         "Suggests improvements to prompts.\nThinks of creative approaches.",
         COLORS["meta_gold"]),
    ]

    for (x, y, title, desc, color) in roles:
        add_rounded_box(slide, x, y, box_w, box_h, color)
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.1), box_w - Inches(0.4), Inches(0.4),
            title, font_size=22, font_color=COLORS["white"], bold=True,
            alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.6), box_w - Inches(0.4), Inches(0.7),
            desc, font_size=15, font_color=COLORS["light_gray"],
            alignment=PP_ALIGN.CENTER
        )

    add_takeaway_bar(slide,
        "You'll work in breakout rooms. One person shares screen -- everyone contributes!",
        COLORS["team_cyan"])


def create_slide_11_challenge_rules(prs):
    """Slide 11: Group Challenge Rules & Scoring"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Group Challenge Rules & Scoring")

    # Rules
    add_rounded_box(
        slide, Inches(0.7), Inches(1.7), Inches(5.6), Inches(4.5),
        COLORS["dark_blue"]
    )
    add_styled_textbox(
        slide, Inches(0.9), Inches(1.85), Inches(5.2), Inches(0.4),
        "Rules", font_size=24, font_color=COLORS["meta_gold"], bold=True,
        alignment=PP_ALIGN.CENTER
    )

    rules = [
        "Everyone participates -- no passengers!",
        "Use meta prompting to improve your prompts",
        "Document your process (prompts + changes)",
        "Fact-check at least 2 AI claims",
        "Be respectful of other teams' work",
        "Have fun and be creative!",
    ]

    for i, rule in enumerate(rules):
        add_styled_textbox(
            slide, Inches(1.1), Inches(2.5 + i * 0.55), Inches(4.8), Inches(0.4),
            f"{i + 1}. {rule}", font_size=14, font_color=COLORS["white"]
        )

    # Scoring
    add_rounded_box(
        slide, Inches(6.93), Inches(1.7), Inches(5.7), Inches(4.5),
        COLORS["meta_gold"]
    )
    add_styled_textbox(
        slide, Inches(7.13), Inches(1.85), Inches(5.3), Inches(0.4),
        "Scoring", font_size=24, font_color=COLORS["dark_blue"], bold=True,
        alignment=PP_ALIGN.CENTER
    )

    scores = [
        ("Prompt Quality", "Did you use RTCF + meta prompting?"),
        ("Iteration", "Did you improve your prompts?"),
        ("Fact-Checking", "Did you verify AI claims?"),
        ("Teamwork", "Did everyone contribute?"),
        ("Creativity", "Did you go above and beyond?"),
        ("Presentation", "Can you explain your process?"),
    ]

    for i, (category, desc) in enumerate(scores):
        y = Inches(2.5 + i * 0.55)
        add_styled_textbox(
            slide, Inches(7.3), y, Inches(2.5), Inches(0.4),
            category, font_size=14, font_color=COLORS["dark_blue"], bold=True
        )
        add_styled_textbox(
            slide, Inches(9.7), y, Inches(2.7), Inches(0.4),
            desc, font_size=12, font_color=COLORS["dark_gray"]
        )


def create_slide_12_round1(prs):
    """Slide 12: Round 1 - Study Guide Battle"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Round 1: AI Study Guide Battle", subtitle="40 minutes")

    # Instructions
    add_rounded_box(
        slide, Inches(0.7), Inches(1.7), Inches(11.93), Inches(1.3),
        COLORS["ai_purple"]
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(1.8), Inches(11.33), Inches(0.4),
        "Mission: Create the BEST study guide on a topic your team chooses!",
        font_size=20, font_color=COLORS["white"], bold=True,
        alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(2.3), Inches(11.33), Inches(0.5),
        "Pick any school subject -- math, science, history, English -- anything you're actually studying!",
        font_size=16, font_color=COLORS["light_gray"],
        alignment=PP_ALIGN.CENTER
    )

    # Steps
    steps = [
        ("1", "Pick a topic and write a basic prompt", COLORS["medium_blue"]),
        ("2", "Use meta prompting to improve your prompt", COLORS["meta_gold"]),
        ("3", "Run the improved prompt and review the result", COLORS["prompt_green"]),
        ("4", "Fact-check 2+ claims from the AI response", COLORS["red"]),
        ("5", "Iterate at least 3 times to get the best result", COLORS["teal"]),
    ]

    for i, (num, text, color) in enumerate(steps):
        y = Inches(3.3 + i * 0.6)
        add_circle(slide, Inches(1.0), y, Inches(0.45), color, text=num, text_size=18)
        add_styled_textbox(
            slide, Inches(1.7), y + Inches(0.05), Inches(10.5), Inches(0.4),
            text, font_size=18, font_color=COLORS["dark_gray"], bold=True
        )

    add_takeaway_bar(slide,
        "Document everything in your classwork sheet! You'll present your process later.",
        COLORS["team_cyan"])


def create_slide_13_ai_showdown(prs):
    """Slide 13: Activity 1 - AI Showdown"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Activity 1: AI Showdown", subtitle="20 minutes")

    # Instructions
    add_rounded_box(
        slide, Inches(0.7), Inches(1.7), Inches(11.93), Inches(1.2),
        COLORS["team_cyan"]
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(1.8), Inches(11.33), Inches(0.4),
        "Mission: Run the SAME prompt on different AI tools and compare the results!",
        font_size=20, font_color=COLORS["white"], bold=True,
        alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(2.3), Inches(11.33), Inches(0.4),
        "Each person in your group picks a different tool. Then compare side-by-side!",
        font_size=16, font_color=COLORS["light_gray"],
        alignment=PP_ALIGN.CENTER
    )

    # Steps
    steps = [
        ("1", "Everyone gets the SAME prompt from the teacher", COLORS["medium_blue"]),
        ("2", "Each person runs it on a DIFFERENT tool", COLORS["ai_purple"]),
        ("3", "Compare responses: accuracy, detail, and style", COLORS["meta_gold"]),
        ("4", "Rate each tool from 1-5 on your worksheet", COLORS["prompt_green"]),
        ("5", "Decide as a team: which AI did the best job?", COLORS["teal"]),
    ]

    for i, (num, text, color) in enumerate(steps):
        y = Inches(3.2 + i * 0.6)
        add_circle(slide, Inches(1.0), y, Inches(0.45), color, text=num, text_size=18)
        add_styled_textbox(
            slide, Inches(1.7), y + Inches(0.05), Inches(10.5), Inches(0.4),
            text, font_size=18, font_color=COLORS["dark_gray"], bold=True
        )

    add_takeaway_bar(slide,
        "Use the AI Showdown worksheet to document your results!",
        COLORS["team_cyan"])


def create_slide_14_creative_challenge(prs):
    """Slide 14: Activity 2 - Creative Challenge"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Activity 2: Creative Challenge", subtitle="35 minutes")

    # Instructions
    add_rounded_box(
        slide, Inches(0.7), Inches(1.7), Inches(11.93), Inches(1.2),
        COLORS["meta_gold"]
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(1.8), Inches(11.33), Inches(0.4),
        "Mission: Use meta prompting to create something amazing with AI!",
        font_size=20, font_color=COLORS["dark_blue"], bold=True,
        alignment=PP_ALIGN.CENTER
    )
    add_styled_textbox(
        slide, Inches(1.0), Inches(2.3), Inches(11.33), Inches(0.4),
        "Pick ONE project, iterate your prompts, and make it the best it can be!",
        font_size=16, font_color=COLORS["dark_gray"],
        alignment=PP_ALIGN.CENTER
    )

    # Project options
    box_w = Inches(5.6)
    box_h = Inches(1.0)
    left_x = Inches(0.7)
    right_x = Inches(6.93)

    options = [
        (left_x, Inches(3.2), "Write a Short Story or Poem",
         "Use AI to co-write something creative!", COLORS["ai_purple"]),
        (right_x, Inches(3.2), "Create a Quiz for Younger Kids",
         "Make a fun quiz on any school topic!", COLORS["prompt_green"]),
        (left_x, Inches(4.4), "Design a Lesson Plan",
         "Teach AI to teach! Pick any topic.", COLORS["teal"]),
        (right_x, Inches(4.4), "Build a How-To Guide",
         "Step-by-step guide for a hobby or skill!", COLORS["team_cyan"]),
    ]

    for (x, y, title, desc, color) in options:
        add_rounded_box(slide, x, y, box_w, box_h, color)
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.1), box_w - Inches(0.4), Inches(0.4),
            title, font_size=18, font_color=COLORS["white"], bold=True,
            alignment=PP_ALIGN.CENTER
        )
        add_styled_textbox(
            slide, x + Inches(0.2), y + Inches(0.5), box_w - Inches(0.4), Inches(0.4),
            desc, font_size=14, font_color=COLORS["light_gray"],
            alignment=PP_ALIGN.CENTER
        )

    add_takeaway_bar(slide,
        "You'll present your best project to the class! Make it awesome!",
        COLORS["deep_purple"])


def create_slide_15_presentations(prs):
    """Slide 15: Presentation & Voting Time"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Presentation & Voting Time!")

    add_styled_textbox(
        slide, Inches(0.7), Inches(1.6), Inches(11.93), Inches(0.5),
        "Each team gets 2-3 minutes to share their best work!",
        font_size=22, font_color=COLORS["dark_gray"], bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # What to present
    add_rounded_box(
        slide, Inches(0.7), Inches(2.3), Inches(5.6), Inches(2.8),
        COLORS["dark_blue"]
    )
    add_styled_textbox(
        slide, Inches(0.9), Inches(2.45), Inches(5.2), Inches(0.4),
        "What to Share:", font_size=22, font_color=COLORS["meta_gold"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    share_items = [
        "Your best study guide or creative project",
        "Your original prompt vs. final prompt",
        "What meta prompting techniques you used",
        "What you fact-checked and what you found",
        "What surprised you about AI's responses",
    ]
    for i, item in enumerate(share_items):
        add_styled_textbox(
            slide, Inches(1.1), Inches(3.0 + i * 0.4), Inches(4.8), Inches(0.35),
            f"- {item}", font_size=14, font_color=COLORS["white"]
        )

    # Awards
    add_rounded_box(
        slide, Inches(6.93), Inches(2.3), Inches(5.7), Inches(2.8),
        COLORS["meta_gold"]
    )
    add_styled_textbox(
        slide, Inches(7.13), Inches(2.45), Inches(5.3), Inches(0.4),
        "Awards:", font_size=22, font_color=COLORS["dark_blue"],
        bold=True, alignment=PP_ALIGN.CENTER
    )
    awards = [
        ("Best Study Guide", COLORS["ai_purple"]),
        ("Best Creative Project", COLORS["prompt_green"]),
        ("Best Prompt Strategy", COLORS["meta_gold"]),
        ("Best Teamwork", COLORS["team_cyan"]),
    ]
    for i, (award, color) in enumerate(awards):
        y = Inches(3.1 + i * 0.5)
        add_circle(slide, Inches(7.5), y, Inches(0.35), color,
                   text=str(i + 1), text_size=14)
        add_styled_textbox(
            slide, Inches(8.1), y + Inches(0.03), Inches(4.3), Inches(0.35),
            award, font_size=18, font_color=COLORS["dark_blue"], bold=True
        )

    add_takeaway_bar(slide,
        "Class votes on the winners! May the best team win!",
        COLORS["team_cyan"])


def create_slide_16_key_takeaways(prs):
    """Slide 16: Key Takeaways"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Today's Key Takeaways")

    takeaways = [
        ("Meta prompting = using AI to help write better prompts", COLORS["meta_gold"]),
        ("Ask AI to write, improve, or suggest prompts for you", COLORS["ai_purple"]),
        ("Iteration is key -- your 3rd prompt is always better than your 1st", COLORS["prompt_green"]),
        ("Different AI tools give different results -- compare them!", COLORS["team_cyan"]),
        ("Always fact-check AI responses, even good-looking ones", COLORS["red"]),
        ("AI is a collaboration tool -- teams + AI = amazing results", COLORS["teal"]),
    ]

    for i, (text, color) in enumerate(takeaways):
        y = 1.7 + i * 0.75
        add_circle(slide, Inches(1.0), Inches(y), Inches(0.55), color,
                   text=str(i + 1), text_size=18)
        add_styled_textbox(
            slide, Inches(1.8), Inches(y + 0.05), Inches(10.5), Inches(0.5),
            text, font_size=22, font_color=COLORS["dark_gray"], bold=True
        )

    add_rounded_box(
        slide, Inches(2.5), Inches(6.4), Inches(8.33), Inches(0.7),
        COLORS["orange"],
        text="Time for Kahoot!",
        text_size=22, text_color=COLORS["white"]
    )


def create_slide_17_questions(prs):
    """Slide 17: Questions Slide"""
    create_questions_slide(prs, "Great work today, teams!")


# ============================================================================
# MAIN
# ============================================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_slide_1_title(prs)
    create_slide_2_agenda(prs)
    create_slide_3_recap(prs)
    create_slide_4_what_is_meta(prs)
    create_slide_5_technique_1(prs)
    create_slide_6_technique_2(prs)
    create_slide_7_technique_3(prs)
    create_slide_8_live_demo(prs)
    create_slide_9_ai_tools(prs)
    create_slide_10_team_roles(prs)
    create_slide_11_challenge_rules(prs)
    create_slide_12_round1(prs)
    create_slide_13_ai_showdown(prs)
    create_slide_14_creative_challenge(prs)
    create_slide_15_presentations(prs)
    create_slide_16_key_takeaways(prs)
    create_slide_17_questions(prs)

    script_dir = os.path.dirname(os.path.abspath(__file__))
    parent_dir = os.path.dirname(script_dir)
    output_path = os.path.join(parent_dir, "2026-04-18_Meta_Prompting_and_AI_Team_Challenge_BEAUTIFUL.pptx")
    prs.save(output_path)
    print(f"[SUCCESS] Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
