"""
PowerPoint Generator for: How AI Agents Actually Work
Kids Computer Science Class - May 2, 2026

Uses the standard theme from tools/create_theme.py.
This is week 4 of the AI phase, building on 4/25 Intro to AI Agents.

Run:
    py create_presentation.py

Dependencies:
    pip install python-pptx
"""

import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', 'tools'))
from create_theme import *

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


# ---- Theme-specific colors (extend the standard palette) --------------------
COLORS["agent_purple"] = RGBColor(103, 58, 183)
COLORS["think_blue"] = RGBColor(25, 118, 210)
COLORS["act_orange"] = RGBColor(245, 124, 0)
COLORS["observe_green"] = RGBColor(56, 142, 60)
COLORS["loop_teal"] = RGBColor(0, 137, 123)
COLORS["meta_gold"] = RGBColor(255, 179, 0)


# =============================================================================
# SLIDE FUNCTIONS
# =============================================================================

def slide_1_title(prs):
    create_title_slide(
        prs,
        "How AI Agents\nActually Work",
        "Inside the Loop — Tools, Memory, and Why It All Works",
        "May 2, 2026",
        tagline="Week 4 of the AI Phase  |  Kids Computer Science Class"
    )


def slide_2_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Today's Journey")

    items = [
        ("Quick Recap: What Is an AI Agent?", COLORS["light_blue"]),
        ("The LLM Brain — How It Thinks", COLORS["agent_purple"]),
        ("System Prompts & Context Windows", COLORS["meta_gold"]),
        ("Tools — How an Agent Takes Action", COLORS["act_orange"]),
        ("Live Demo: Same Task, Three Ways", COLORS["observe_green"]),
        ("Why Agents Sometimes Fail", COLORS["red"]),
        ("Hands-On: Agent Trace Lab", COLORS["loop_teal"]),
        ("Kahoot!", COLORS["green"]),
    ]
    for i, (text, color) in enumerate(items):
        add_agenda_item(slide, i + 1, text, 1.6 + i * 0.65, color)


def slide_3_recap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Quick Recap: Last Week's Agent Loop")

    # Four loop steps as colored circles in a row
    steps = [
        ("THINK", COLORS["think_blue"], "Plan what to do"),
        ("ACT", COLORS["act_orange"], "Use a tool"),
        ("OBSERVE", COLORS["observe_green"], "Look at result"),
        ("REPEAT", COLORS["loop_teal"], "Until done"),
    ]
    for i, (label, color, desc) in enumerate(steps):
        x = 0.7 + i * 3.1
        add_circle(slide, Inches(x + 0.45), Inches(2.0), Inches(1.6), color,
                   text=label, text_size=18)
        add_styled_textbox(
            slide, Inches(x), Inches(3.8), Inches(2.5), Inches(0.5),
            desc, font_size=16, font_color=COLORS["dark_gray"],
            alignment=PP_ALIGN.CENTER
        )

    # Arrows between circles
    for i in range(3):
        x = 2.3 + i * 3.1
        add_arrow(slide, Inches(x), Inches(2.7), Inches(0.7), Inches(0.4),
                  COLORS["dark_gray"])

    add_takeaway_bar(slide,
        "Today: HOW does the agent decide what to do at each step?",
        COLORS["meta_gold"])


def slide_4_question(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLORS["dark_blue"])

    add_styled_textbox(
        slide, Inches(0.5), Inches(2.0), Inches(12.33), Inches(1.0),
        "But HOW does it actually work?",
        font_size=46, font_color=COLORS["white"], bold=True,
        alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold"
    )

    add_styled_textbox(
        slide, Inches(0.5), Inches(3.3), Inches(12.33), Inches(0.7),
        "Last week we saw what agents can DO.",
        font_size=24, font_color=COLORS["sky_blue"],
        alignment=PP_ALIGN.CENTER
    )

    add_styled_textbox(
        slide, Inches(0.5), Inches(4.0), Inches(12.33), Inches(0.7),
        "Today we open the hood.",
        font_size=24, font_color=COLORS["orange"], bold=True,
        alignment=PP_ALIGN.CENTER
    )


def slide_5_llm_brain(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "The LLM Brain", "It's just predicting the next word")

    # Prompt -> brain -> next word visualization
    add_rounded_box(slide, Inches(0.6), Inches(2.0), Inches(3.3), Inches(1.5),
                    COLORS["light_blue"])
    add_styled_textbox(slide, Inches(0.6), Inches(2.2), Inches(3.3), Inches(0.5),
                       "YOU TYPE", font_size=16, font_color=COLORS["dark_blue"],
                       bold=True, alignment=PP_ALIGN.CENTER)
    add_styled_textbox(slide, Inches(0.6), Inches(2.6), Inches(3.3), Inches(0.8),
                       '"The capital of France is..."',
                       font_size=14, font_color=COLORS["dark_gray"],
                       alignment=PP_ALIGN.CENTER)

    add_arrow(slide, Inches(4.0), Inches(2.55), Inches(0.8), Inches(0.4),
              COLORS["dark_gray"])

    add_rounded_box(slide, Inches(4.9), Inches(2.0), Inches(3.5), Inches(1.5),
                    COLORS["agent_purple"])
    add_styled_textbox(slide, Inches(4.9), Inches(2.2), Inches(3.5), Inches(0.5),
                       "LLM PREDICTS", font_size=16, font_color=COLORS["white"],
                       bold=True, alignment=PP_ALIGN.CENTER)
    add_styled_textbox(slide, Inches(4.9), Inches(2.6), Inches(3.5), Inches(0.8),
                       'Most likely next word?',
                       font_size=14, font_color=COLORS["white"],
                       alignment=PP_ALIGN.CENTER)

    add_arrow(slide, Inches(8.5), Inches(2.55), Inches(0.8), Inches(0.4),
              COLORS["dark_gray"])

    add_rounded_box(slide, Inches(9.4), Inches(2.0), Inches(3.3), Inches(1.5),
                    COLORS["green"])
    add_styled_textbox(slide, Inches(9.4), Inches(2.2), Inches(3.3), Inches(0.5),
                       "OUTPUT", font_size=16, font_color=COLORS["white"],
                       bold=True, alignment=PP_ALIGN.CENTER)
    add_styled_textbox(slide, Inches(9.4), Inches(2.6), Inches(3.3), Inches(0.8),
                       '"Paris"',
                       font_size=18, font_color=COLORS["white"], bold=True,
                       alignment=PP_ALIGN.CENTER)

    # Bottom explanation
    add_rounded_box(slide, Inches(1.0), Inches(4.3), Inches(11.33), Inches(1.5),
                    COLORS["light_gray"], border_color=COLORS["medium_blue"])
    add_styled_textbox(slide, Inches(1.3), Inches(4.45), Inches(10.7), Inches(0.5),
                       "It does this OVER and OVER, one word at a time.",
                       font_size=18, font_color=COLORS["dark_blue"], bold=True,
                       alignment=PP_ALIGN.CENTER)
    add_styled_textbox(slide, Inches(1.3), Inches(4.85), Inches(10.7), Inches(0.85),
                       'A long answer is just hundreds of "predict the next word" steps in a row.',
                       font_size=16, font_color=COLORS["dark_gray"],
                       alignment=PP_ALIGN.CENTER)

    add_takeaway_bar(slide,
        "An LLM doesn't 'know' answers — it predicts text that sounds right.",
        COLORS["dark_blue"])


def slide_6_system_prompt(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "The System Prompt", "Hidden instructions you never see")

    # Visual: hidden message + user message stack
    add_rounded_box(slide, Inches(1.0), Inches(1.7), Inches(11.33), Inches(1.5),
                    COLORS["red"], border_color=COLORS["dark_blue"])
    add_styled_textbox(slide, Inches(1.3), Inches(1.85), Inches(10.7), Inches(0.5),
                       "SYSTEM PROMPT (hidden)", font_size=14,
                       font_color=COLORS["white"], bold=True)
    add_styled_textbox(slide, Inches(1.3), Inches(2.2), Inches(10.7), Inches(0.95),
                       '"You are a helpful tutor for kids ages 10-15. Be patient. '
                       'Never share personal info. If unsure, say so."',
                       font_size=14, font_color=COLORS["white"])

    add_rounded_box(slide, Inches(1.0), Inches(3.4), Inches(11.33), Inches(1.0),
                    COLORS["light_blue"])
    add_styled_textbox(slide, Inches(1.3), Inches(3.55), Inches(10.7), Inches(0.4),
                       "USER MESSAGE (what you type)", font_size=14,
                       font_color=COLORS["dark_blue"], bold=True)
    add_styled_textbox(slide, Inches(1.3), Inches(3.85), Inches(10.7), Inches(0.5),
                       '"Help me with my math homework."',
                       font_size=14, font_color=COLORS["dark_gray"])

    # Two facts
    facts = [
        ("Set by the AI company or developer", COLORS["medium_blue"]),
        ("Tells the AI HOW to behave & rules to follow", COLORS["agent_purple"]),
    ]
    for i, (text, color) in enumerate(facts):
        x = 0.7 + i * 6.2
        add_rounded_box(slide, Inches(x), Inches(4.7), Inches(5.9), Inches(0.9),
                        color, text=text, text_size=15)

    add_takeaway_bar(slide,
        "Every AI you use has a hidden system prompt shaping its behavior.",
        COLORS["agent_purple"])


def slide_7_context_window(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "The Context Window", "The agent's working memory")

    # Whiteboard analogy
    add_rounded_box(slide, Inches(0.7), Inches(1.7), Inches(5.8), Inches(3.6),
                    COLORS["light_gray"], border_color=COLORS["dark_gray"])
    add_styled_textbox(slide, Inches(0.7), Inches(1.85), Inches(5.8), Inches(0.5),
                       "Think of a Whiteboard", font_size=18,
                       font_color=COLORS["dark_blue"], bold=True,
                       alignment=PP_ALIGN.CENTER)
    items_a = [
        "It only fits so much text at once",
        "Once full, old stuff gets erased",
        "The AI can ONLY see what's on it",
        "When you start a new chat -> new whiteboard",
    ]
    for i, item in enumerate(items_a):
        add_styled_textbox(slide, Inches(0.95), Inches(2.5 + i * 0.55),
                           Inches(5.5), Inches(0.5),
                           f"• {item}", font_size=14,
                           font_color=COLORS["dark_gray"])

    # Real numbers
    add_rounded_box(slide, Inches(6.85), Inches(1.7), Inches(5.8), Inches(3.6),
                    COLORS["agent_purple"])
    add_styled_textbox(slide, Inches(6.85), Inches(1.85), Inches(5.8), Inches(0.5),
                       "How Big Is It?", font_size=18,
                       font_color=COLORS["white"], bold=True,
                       alignment=PP_ALIGN.CENTER)
    sizes = [
        ("Older chatbots", "~4,000 words"),
        ("Modern AI (GPT-4o)", "~100,000 words"),
        ("Claude / Gemini Pro", "~750,000 - 1.5M words"),
        ("Still — there is ALWAYS a limit", ""),
    ]
    for i, (label, val) in enumerate(sizes):
        add_styled_textbox(slide, Inches(7.1), Inches(2.5 + i * 0.55),
                           Inches(3.0), Inches(0.5),
                           label, font_size=13, font_color=COLORS["white"],
                           bold=True)
        add_styled_textbox(slide, Inches(10.0), Inches(2.5 + i * 0.55),
                           Inches(2.6), Inches(0.5),
                           val, font_size=13, font_color=COLORS["sky_blue"])

    add_takeaway_bar(slide,
        "If the chat gets too long, the AI 'forgets' the early parts.",
        COLORS["dark_blue"])


def slide_8_tools_intro(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Tools",
                  "What turns an LLM into an AGENT")

    add_styled_textbox(
        slide, Inches(0.5), Inches(1.7), Inches(12.33), Inches(0.6),
        "An LLM only outputs text. Tools let the agent DO things.",
        font_size=20, font_color=COLORS["dark_blue"], bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # 6 tool cards in 2 rows
    tools = [
        ("Web Search", "Look up live info", COLORS["medium_blue"]),
        ("Code Execution", "Run real code", COLORS["agent_purple"]),
        ("Image Generation", "Create pictures", COLORS["pink"]),
        ("File Reader", "Read your docs", COLORS["teal"]),
        ("Email / Calendar", "Send & schedule", COLORS["act_orange"]),
        ("Browser Control", "Click & type for you", COLORS["observe_green"]),
    ]
    for i, (name, desc, color) in enumerate(tools):
        col = i % 3
        row = i // 3
        x = 0.55 + col * 4.3
        y = 2.6 + row * 1.7
        add_rounded_box(slide, Inches(x), Inches(y), Inches(4.0), Inches(1.4),
                        color)
        add_styled_textbox(slide, Inches(x + 0.1), Inches(y + 0.2),
                           Inches(3.8), Inches(0.5),
                           name, font_size=18, font_color=COLORS["white"],
                           bold=True, alignment=PP_ALIGN.CENTER)
        add_styled_textbox(slide, Inches(x + 0.1), Inches(y + 0.75),
                           Inches(3.8), Inches(0.5),
                           desc, font_size=14, font_color=COLORS["light_gray"],
                           alignment=PP_ALIGN.CENTER)

    add_takeaway_bar(slide,
        "Different tools = different superpowers. More tools = more capable agent.",
        COLORS["loop_teal"])


def slide_9_tool_call(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "How a Tool Call Actually Works",
                  "Like ordering at a restaurant")

    # Three columns: AI -> request -> tool runs -> result back
    cols = [
        ("1. AI says\n'I need to search'", COLORS["agent_purple"], "AI THINKS"),
        ("2. Request goes\nto a server",   COLORS["medium_blue"],  "NETWORK"),
        ("3. Server runs\nthe search",     COLORS["act_orange"],   "TOOL RUNS"),
        ("4. Result comes\nback to AI",    COLORS["observe_green"], "AI READS"),
    ]
    for i, (text, color, label) in enumerate(cols):
        x = 0.5 + i * 3.2
        add_rounded_box(slide, Inches(x), Inches(2.0), Inches(2.9), Inches(2.4),
                        color)
        add_styled_textbox(slide, Inches(x + 0.1), Inches(2.15),
                           Inches(2.7), Inches(0.5),
                           label, font_size=14, font_color=COLORS["white"],
                           bold=True, alignment=PP_ALIGN.CENTER)
        add_styled_textbox(slide, Inches(x + 0.1), Inches(2.7),
                           Inches(2.7), Inches(1.5),
                           text, font_size=15, font_color=COLORS["white"],
                           alignment=PP_ALIGN.CENTER)
        if i < 3:
            add_arrow(slide, Inches(x + 2.95), Inches(3.0),
                      Inches(0.25), Inches(0.4), COLORS["dark_gray"])

    add_rounded_box(slide, Inches(1.0), Inches(4.7), Inches(11.33), Inches(1.0),
                    COLORS["meta_gold"])
    add_styled_textbox(slide, Inches(1.2), Inches(4.85), Inches(11.0),
                       Inches(0.7),
                       "AI CONNECTION: Every tool call is an HTTP request — "
                       "just like loading a website!",
                       font_size=18, font_color=COLORS["dark_blue"], bold=True,
                       alignment=PP_ALIGN.CENTER)

    add_takeaway_bar(slide,
        "The AI doesn't 'do' the search. It asks another program to do it.",
        COLORS["dark_blue"])


def slide_10_loop_walkthrough(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Walkthrough: Booking a Trip",
                  "The loop in action")

    user_q = '"Find me a cheap flight from SFO to NYC next month and email it to my mom."'
    add_rounded_box(slide, Inches(0.7), Inches(1.7), Inches(11.93), Inches(0.7),
                    COLORS["light_blue"])
    add_styled_textbox(slide, Inches(0.9), Inches(1.85), Inches(11.5),
                       Inches(0.5), user_q,
                       font_size=15, font_color=COLORS["dark_blue"], bold=True,
                       alignment=PP_ALIGN.CENTER)

    steps = [
        ("THINK", "I need flight prices first.", COLORS["think_blue"]),
        ("ACT", "Search: 'SFO to NYC June flights'", COLORS["act_orange"]),
        ("OBSERVE", "Cheapest is $189 on JetBlue.", COLORS["observe_green"]),
        ("THINK", "Now I need mom's email address.", COLORS["think_blue"]),
        ("ACT", "Send email with flight details.", COLORS["act_orange"]),
        ("OBSERVE", "Email sent successfully.", COLORS["observe_green"]),
    ]
    for i, (label, text, color) in enumerate(steps):
        col = i % 2
        row = i // 2
        x = 0.7 + col * 6.15
        y = 2.65 + row * 1.25
        add_rounded_box(slide, Inches(x), Inches(y), Inches(1.4),
                        Inches(1.05), color)
        add_styled_textbox(slide, Inches(x + 0.05), Inches(y + 0.3),
                           Inches(1.3), Inches(0.5),
                           label, font_size=14, font_color=COLORS["white"],
                           bold=True, alignment=PP_ALIGN.CENTER)
        add_rounded_box(slide, Inches(x + 1.5), Inches(y), Inches(4.4),
                        Inches(1.05), COLORS["light_gray"],
                        border_color=color)
        add_styled_textbox(slide, Inches(x + 1.6), Inches(y + 0.3),
                           Inches(4.2), Inches(0.5),
                           text, font_size=13, font_color=COLORS["dark_gray"],
                           alignment=PP_ALIGN.LEFT)

    add_takeaway_bar(slide,
        "One user request = many small loop steps behind the scenes.",
        COLORS["loop_teal"])


def slide_11_reasoning(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Reasoning",
                  "When the AI 'thinks out loud'")

    add_styled_textbox(
        slide, Inches(0.5), Inches(1.7), Inches(12.33), Inches(0.6),
        "Modern agents show their thinking BEFORE they answer.",
        font_size=18, font_color=COLORS["dark_blue"],
        alignment=PP_ALIGN.CENTER
    )

    # Show fake "thinking" output
    add_rounded_box(slide, Inches(1.5), Inches(2.5), Inches(10.33),
                    Inches(2.6), COLORS["light_gray"],
                    border_color=COLORS["agent_purple"])
    add_styled_textbox(slide, Inches(1.7), Inches(2.65), Inches(10.0),
                       Inches(0.4), "[Agent thinking...]",
                       font_size=14, font_color=COLORS["agent_purple"],
                       bold=True)
    thoughts = [
        "Hmm, the user wants the cheapest flight.",
        "First I should figure out their dates.",
        "They said 'next month' — that's June 2026.",
        "I should search a flight site for SFO -> NYC in June.",
        "Then compare prices and pick the lowest.",
        "Then I'll need to find mom's email...",
    ]
    for i, t in enumerate(thoughts):
        add_styled_textbox(slide, Inches(1.7), Inches(3.05 + i * 0.32),
                           Inches(10.0), Inches(0.32),
                           f"• {t}", font_size=13,
                           font_color=COLORS["dark_gray"])

    add_rounded_box(slide, Inches(1.0), Inches(5.4), Inches(11.33),
                    Inches(0.85), COLORS["meta_gold"])
    add_styled_textbox(slide, Inches(1.2), Inches(5.55), Inches(11.0),
                       Inches(0.6),
                       "WHY IT MATTERS: You can SEE if the agent's plan is "
                       "right — and stop it before it does something wrong.",
                       font_size=15, font_color=COLORS["dark_blue"], bold=True,
                       alignment=PP_ALIGN.CENTER)

    add_takeaway_bar(slide,
        "Reasoning makes agents more trustworthy — you see the plan first.",
        COLORS["agent_purple"])


def slide_12_chatbot_vs_agent(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Chatbot vs. Assistant vs. Agent")

    rows = [
        ("Capability",      "Chatbot",       "Assistant",     "Agent"),
        ("Answers questions",   "YES",       "YES",           "YES"),
        ("Remembers chat",      "NO",        "YES",           "YES"),
        ("Uses tools",          "NO",        "NO",            "YES"),
        ("Multi-step tasks",    "NO",        "Sort of",       "YES"),
        ("Can search live info", "NO",       "NO",            "YES"),
    ]
    col_widths = [3.6, 2.9, 2.9, 2.9]
    col_starts = [0.55, 4.25, 7.25, 10.25]

    # Header row
    add_table_row(slide, 1.7, rows[0], col_widths, col_starts,
                  COLORS["dark_blue"], font_size=15, bold=True,
                  row_height=0.6)

    # Body rows
    body_colors = [COLORS["light_blue"], COLORS["sky_blue"], COLORS["light_blue"],
                   COLORS["sky_blue"], COLORS["light_blue"]]
    for i, row in enumerate(rows[1:]):
        add_table_row(slide, 2.4 + i * 0.7, row, col_widths, col_starts,
                      body_colors[i],
                      text_color=COLORS["dark_blue"],
                      font_size=14, row_height=0.6)

    add_takeaway_bar(slide,
        "Agents = LLM brain + memory + tools + planning, all working together.",
        COLORS["loop_teal"])


def slide_13_demo_intro(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, COLORS["observe_green"])

    add_styled_textbox(
        slide, Inches(0.5), Inches(1.5), Inches(12.33), Inches(1.0),
        "LIVE DEMO",
        font_size=56, font_color=COLORS["white"], bold=True,
        alignment=PP_ALIGN.CENTER, font_name="Segoe UI Semibold"
    )

    add_styled_textbox(
        slide, Inches(0.5), Inches(2.7), Inches(12.33), Inches(0.7),
        "Same Task, Three Ways",
        font_size=32, font_color=COLORS["white"],
        alignment=PP_ALIGN.CENTER
    )

    # The prompt
    add_rounded_box(slide, Inches(1.5), Inches(3.9), Inches(10.33),
                    Inches(1.6), COLORS["white"])
    add_styled_textbox(slide, Inches(1.7), Inches(4.05), Inches(10.0),
                       Inches(0.4),
                       "WE'LL ASK ALL THREE:",
                       font_size=14, font_color=COLORS["dark_blue"],
                       bold=True, alignment=PP_ALIGN.CENTER)
    add_styled_textbox(slide, Inches(1.7), Inches(4.45), Inches(10.0),
                       Inches(1.0),
                       '"What is the most popular video game right now, '
                       'how many people are playing it, and write a '
                       'short poem about it?"',
                       font_size=18, font_color=COLORS["dark_gray"],
                       alignment=PP_ALIGN.CENTER)

    add_styled_textbox(
        slide, Inches(0.5), Inches(5.9), Inches(12.33), Inches(0.6),
        "Watch what happens differently in each one.",
        font_size=20, font_color=COLORS["white"],
        alignment=PP_ALIGN.CENTER
    )


def slide_14_why_fail(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Why Agents Sometimes Fail",
                  "They are powerful — not perfect")

    fails = [
        ("Wrong Tool",
         "Picks web search when it should run code",
         COLORS["red"]),
        ("Bad Plan",
         "Plans 5 steps but skips a key one",
         COLORS["act_orange"]),
        ("Stuck Loop",
         "Keeps retrying the same broken approach",
         COLORS["agent_purple"]),
        ("Hallucination",
         "Confidently makes up info when uncertain",
         COLORS["pink"]),
    ]
    for i, (name, desc, color) in enumerate(fails):
        col = i % 2
        row = i // 2
        x = 0.7 + col * 6.15
        y = 1.8 + row * 1.95
        add_rounded_box(slide, Inches(x), Inches(y), Inches(5.9),
                        Inches(1.7), color)
        add_styled_textbox(slide, Inches(x + 0.2), Inches(y + 0.2),
                           Inches(5.5), Inches(0.5),
                           name, font_size=20, font_color=COLORS["white"],
                           bold=True, alignment=PP_ALIGN.CENTER)
        add_styled_textbox(slide, Inches(x + 0.2), Inches(y + 0.85),
                           Inches(5.5), Inches(0.7),
                           desc, font_size=15, font_color=COLORS["light_gray"],
                           alignment=PP_ALIGN.CENTER)

    add_takeaway_bar(slide,
        "Always check the agent's work — especially for important tasks.",
        COLORS["red"])


def slide_15_agi(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Wait — Is This AGI?",
                  "A common question — let's clear it up")

    add_rounded_box(slide, Inches(0.6), Inches(1.7), Inches(6.0),
                    Inches(4.2), COLORS["medium_blue"])
    add_styled_textbox(slide, Inches(0.6), Inches(1.9), Inches(6.0),
                       Inches(0.5),
                       "AI AGENT (today)",
                       font_size=22, font_color=COLORS["white"], bold=True,
                       alignment=PP_ALIGN.CENTER)
    agent_facts = [
        "An LLM with tools, memory, and planning",
        "Great at narrow, specific tasks",
        "Still makes mistakes & hallucinates",
        "Needs human supervision",
        "EXISTS RIGHT NOW",
    ]
    for i, f in enumerate(agent_facts):
        add_styled_textbox(slide, Inches(0.9), Inches(2.7 + i * 0.55),
                           Inches(5.6), Inches(0.5),
                           f"• {f}", font_size=15,
                           font_color=COLORS["white"])

    add_rounded_box(slide, Inches(6.85), Inches(1.7), Inches(6.0),
                    Inches(4.2), COLORS["dark_gray"])
    add_styled_textbox(slide, Inches(6.85), Inches(1.9), Inches(6.0),
                       Inches(0.5),
                       "AGI (someday?)",
                       font_size=22, font_color=COLORS["white"], bold=True,
                       alignment=PP_ALIGN.CENTER)
    agi_facts = [
        "Artificial General Intelligence",
        "Human-level reasoning across ALL tasks",
        "Could learn anything a person can",
        "Researchers debate when (if!) it arrives",
        "DOES NOT EXIST YET",
    ]
    for i, f in enumerate(agi_facts):
        add_styled_textbox(slide, Inches(7.15), Inches(2.7 + i * 0.55),
                           Inches(5.6), Inches(0.5),
                           f"• {f}", font_size=15,
                           font_color=COLORS["light_gray"])

    add_takeaway_bar(slide,
        "Agents are powerful tools today. AGI is a question for tomorrow.",
        COLORS["dark_blue"])


def slide_16_classwork(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Today's Classwork: Agent Trace Lab")

    add_rounded_box(slide, Inches(0.7), Inches(1.7), Inches(11.93),
                    Inches(0.9), COLORS["loop_teal"])
    add_styled_textbox(slide, Inches(0.9), Inches(1.85), Inches(11.5),
                       Inches(0.6),
                       "Run a real multi-step prompt and label what you see.",
                       font_size=20, font_color=COLORS["white"], bold=True,
                       alignment=PP_ALIGN.CENTER)

    sections = [
        ("Section 1", "Recap the agent loop",       "15 pts",
         COLORS["agent_purple"]),
        ("Section 2", "Spot the agent skills",       "25 pts",
         COLORS["think_blue"]),
        ("Section 3", "Annotate the trace",          "35 pts",
         COLORS["act_orange"]),
        ("Section 4", "Same task, two ways",         "25 pts",
         COLORS["observe_green"]),
    ]
    for i, (sec, desc, pts, color) in enumerate(sections):
        y = 2.85 + i * 0.85
        add_rounded_box(slide, Inches(0.7), Inches(y), Inches(2.0),
                        Inches(0.7), color, text=sec, text_size=16)
        add_rounded_box(slide, Inches(2.85), Inches(y), Inches(7.5),
                        Inches(0.7), COLORS["light_gray"],
                        border_color=color)
        add_styled_textbox(slide, Inches(3.0), Inches(y + 0.18),
                           Inches(7.3), Inches(0.4),
                           desc, font_size=15, font_color=COLORS["dark_gray"])
        add_rounded_box(slide, Inches(10.5), Inches(y), Inches(2.13),
                        Inches(0.7), color, text=pts, text_size=14)

    add_takeaway_bar(slide,
        "Total: 100 points + 10 bonus  |  Submit on Microsoft Teams",
        COLORS["dark_blue"])


def slide_17_homework(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Homework: Build Your Agent Workflow")

    add_rounded_box(slide, Inches(0.7), Inches(1.7), Inches(11.93),
                    Inches(1.0), COLORS["meta_gold"])
    add_styled_textbox(slide, Inches(0.9), Inches(1.85), Inches(11.5),
                       Inches(0.7),
                       "Design your own multi-step AI agent for a real task in your life.",
                       font_size=18, font_color=COLORS["dark_blue"], bold=True,
                       alignment=PP_ALIGN.CENTER)

    cards = [
        ("Pick a Task", "Something with multiple steps — not a single question",
         COLORS["medium_blue"]),
        ("Choose Tools", "Web search? Email? Code? Pick at least 2",
         COLORS["agent_purple"]),
        ("Write System Prompt", "How should YOUR agent behave?",
         COLORS["act_orange"]),
        ("Map 5 Steps", "What does the agent do at each loop?",
         COLORS["observe_green"]),
    ]
    for i, (title, desc, color) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = 0.7 + col * 6.15
        y = 3.0 + row * 1.55
        add_rounded_box(slide, Inches(x), Inches(y), Inches(5.9),
                        Inches(1.4), color)
        add_styled_textbox(slide, Inches(x + 0.2), Inches(y + 0.15),
                           Inches(5.5), Inches(0.5),
                           title, font_size=18, font_color=COLORS["white"],
                           bold=True, alignment=PP_ALIGN.CENTER)
        add_styled_textbox(slide, Inches(x + 0.2), Inches(y + 0.7),
                           Inches(5.5), Inches(0.6),
                           desc, font_size=14, font_color=COLORS["light_gray"],
                           alignment=PP_ALIGN.CENTER)

    add_takeaway_bar(slide,
        "Total: 100 points + 5 bonus  |  Due next class",
        COLORS["dark_blue"])


def slide_18_takeaways(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Today's Big Takeaways")

    takeaways = [
        ("LLMs predict text, one word at a time", COLORS["agent_purple"]),
        ("Tools turn an LLM into an AGENT", COLORS["act_orange"]),
        ("Every tool call is a network request", COLORS["medium_blue"]),
        ("Memory & context windows have limits", COLORS["loop_teal"]),
        ("Agents can fail — always verify their work", COLORS["red"]),
        ("Agents are NOT AGI — that doesn't exist yet", COLORS["observe_green"]),
    ]
    for i, (text, color) in enumerate(takeaways):
        y = 1.7 + i * 0.75
        add_circle(slide, Inches(1.0), Inches(y + 0.05), Inches(0.55), color,
                   text=str(i + 1), text_size=18)
        add_styled_textbox(slide, Inches(2.0), Inches(y + 0.1),
                           Inches(10.5), Inches(0.6),
                           text, font_size=20, font_color=COLORS["dark_gray"])

    add_takeaway_bar(slide,
        "Next week: We'll start BUILDING with AI tools.",
        COLORS["dark_blue"])


def slide_19_questions(prs):
    create_questions_slide(prs, "Ask away before we start the Kahoot!")


# =============================================================================
# MAIN
# =============================================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    builders = [
        slide_1_title,
        slide_2_agenda,
        slide_3_recap,
        slide_4_question,
        slide_5_llm_brain,
        slide_6_system_prompt,
        slide_7_context_window,
        slide_8_tools_intro,
        slide_9_tool_call,
        slide_10_loop_walkthrough,
        slide_11_reasoning,
        slide_12_chatbot_vs_agent,
        slide_13_demo_intro,
        slide_14_why_fail,
        slide_15_agi,
        slide_16_classwork,
        slide_17_homework,
        slide_18_takeaways,
        slide_19_questions,
    ]
    for fn in builders:
        fn(prs)

    script_dir = os.path.dirname(os.path.abspath(__file__))
    parent_dir = os.path.dirname(script_dir)
    output_path = os.path.join(parent_dir, "2026-05-02_How_AI_Agents_Work.pptx")
    prs.save(output_path)
    print(f"[SUCCESS] Presentation saved: {output_path}")
    print(f"  Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
